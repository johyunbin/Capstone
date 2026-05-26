#!/usr/bin/env python3
"""
v4.pptx → v5.pptx
slide 3·14에 영상 hyperlink overlay 추가:
- 전체 슬라이드 영역에 transparent hyperlink shape (어디 클릭이든 활성)
- 우하단에 작은 cyan ▶ 버튼 (시각 hint)
- v4 디자인 그대로 유지
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CYAN = RGBColor(0x4F, 0xB3, 0xD9)
NAVY = RGBColor(0x1B, 0x2A, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ANIM1_URL = "https://claude.ai/design/p/455bb91a-3c1e-48a4-b8ee-8932f4d8881d?file=index.html&present=1"
ANIM2_URL = "https://claude.ai/design/p/585e3fa7-daae-4cfc-8d8f-a0cf966299ea?file=video.html&present=1"


def add_full_slide_hyperlink(slide, url, prs):
    """슬라이드 전체 영역에 transparent hyperlink shape overlay"""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), prs.slide_width, prs.slide_height
    )
    sh.shadow.inherit = False
    sh.fill.background()  # transparent fill
    sh.line.fill.background()  # no border
    sh.click_action.hyperlink.address = url
    return sh


def add_play_button(slide, prs, url, color=CYAN):
    """우하단에 작은 ▶ cyan circle 버튼"""
    size = Inches(0.9)
    left = prs.slide_width - size - Inches(0.4)
    top = prs.slide_height - size - Inches(0.4)
    btn = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    btn.shadow.inherit = False
    btn.fill.solid()
    btn.fill.fore_color.rgb = color
    btn.line.color.rgb = NAVY
    btn.line.width = Pt(2)

    tf = btn.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_obj = tf.paragraphs[0]
    p_obj.alignment = PP_ALIGN.CENTER
    r = p_obj.add_run()
    r.text = "▶"
    r.font.name = "Apple SD Gothic Neo"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.hyperlink.address = url

    btn.click_action.hyperlink.address = url
    return btn


def add_label(slide, prs, label, url):
    """▶ 버튼 옆 작은 라벨 박스"""
    width = Inches(3.4)
    height = Inches(0.5)
    left = prs.slide_width - width - Inches(1.5)
    top = prs.slide_height - height - Inches(0.6)
    lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    lbl.shadow.inherit = False
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = NAVY
    lbl.line.fill.background()

    tf = lbl.text_frame
    tf.margin_left = Emu(60000); tf.margin_right = Emu(60000)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_obj = tf.paragraphs[0]
    p_obj.alignment = PP_ALIGN.CENTER
    r = p_obj.add_run()
    r.text = label
    r.font.name = "Apple SD Gothic Neo"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.hyperlink.address = url

    lbl.click_action.hyperlink.address = url
    return lbl


def main():
    src = "/Users/hyunbin/Capstone/submission/_drafts/속도는벡터_기말발표_v4.pptx"
    prs = Presentation(src)

    slots = [
        (3, ANIM1_URL, "▶ 영상 ① 재생 · 32초"),
        (14, ANIM2_URL, "▶ 영상 ② 재생 · 45초"),
    ]

    for idx, url, label in slots:
        slide = prs.slides[idx - 1]
        # 전체 슬라이드 클릭 hyperlink (transparent overlay)
        add_full_slide_hyperlink(slide, url, prs)
        # 우하단 라벨 박스
        add_label(slide, prs, label, url)
        # 우하단 ▶ 작은 cyan 버튼
        add_play_button(slide, prs, url)
        print(f"  slide {idx:>2} → hyperlink + ▶ overlay 추가 · URL: {url[:60]}...")

    out = "/Users/hyunbin/Capstone/submission/_drafts/속도는벡터_기말발표_v5_영상hyperlink_20260526_160000.pptx"
    prs.save(out)
    print(f"\n✓ saved: {out}")
    print(f"  슬라이드 수: {len(prs.slides)} (v4 그대로 + 영상 hyperlink 2슬라이드 overlay)")


if __name__ == "__main__":
    main()
