#!/usr/bin/env python3
"""
build_midterm_pptx.py — 중간발표 v2 PPT 생성

사용법:
    python3 scripts/build_midterm_pptx.py

출력:
    submission/속도는벡터_중간발표_20260417_0000.pptx

방침:
- 16:9 와이드 (13.333 x 7.5 inch)
- 한글 폰트: Apple SD Gothic Neo (본문), 제목은 굵게
- 팔레트: 검정 헤더 배경, 파랑 강조(#1f77b4), 주황 Level 1(#e07a00)
- Figure 5장 실제 삽입 (기존 1/2/6 + 신규 7/8)
- 표는 PowerPoint 네이티브 테이블
- 발표자 스크립트는 speaker notes 로 첨부
- 15 슬라이드 (표지 + S2~S14 + 부록)
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# --- 색상 ---
C_BLACK = RGBColor(0x15, 0x15, 0x15)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
C_GRAY_MID = RGBColor(0x9A, 0xA0, 0xA6)
C_BLUE = RGBColor(0x1F, 0x77, 0xB4)
C_ORANGE = RGBColor(0xE0, 0x7A, 0x00)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_BG_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)

# --- 폰트 ---
FONT_MAIN = "Apple SD Gothic Neo"
FONT_CODE = "D2Coding"

# --- 경로 ---
ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "submission" / "속도는벡터_중간발표_20260417_0000.pptx"
FIG_RQ1 = ROOT / "experiments" / "figures" / "rq1_motivation"
FIG_RQ2 = ROOT / "experiments" / "figures" / "rq2_aware"

# --- 슬라이드 크기 16:9 ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run_font(run, size=16, bold=False, color=C_BLACK, font=FONT_MAIN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한국어/CJK 폰트를 동일하게 강제
    rPr = run._r.get_or_add_rPr()
    # eastAsia typeface
    for east in rPr.findall(qn("a:ea")):
        rPr.remove(east)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)


def add_slide_bg(slide, color=C_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title_text, subtitle=None):
    """검정 배경 헤더 (제목 + 선택적 부제)."""
    # 헤더 사각형 — 전체 너비, 높이 1 inch
    left = Inches(0)
    top = Inches(0)
    width = SLIDE_W
    height = Inches(1.0)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BLACK
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_run_font(run, size=26, bold=True, color=C_WHITE)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        set_run_font(r2, size=14, bold=False, color=C_GRAY_LIGHT)

    return shape


def add_textbox(slide, left, top, width, height, text_lines, body_size=14,
                align=PP_ALIGN.LEFT):
    """text_lines: list of (text, {options}) or plain str.
    options: bold, color, size, bullet."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, item in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align

        if isinstance(item, str):
            runs_spec = [(item, {})]
        elif isinstance(item, tuple) and len(item) == 2 and isinstance(item[0], str):
            runs_spec = [item]
        else:
            # list of (text, options)
            runs_spec = item

        for text, opts in runs_spec:
            if opts.get("bullet"):
                # 불릿은 앞에 "• " prefix 로 간단히
                pass
            run = p.add_run()
            run.text = text
            set_run_font(
                run,
                size=opts.get("size", body_size),
                bold=opts.get("bold", False),
                color=opts.get("color", C_BLACK),
                font=opts.get("font", FONT_MAIN),
            )

    return tb


def add_bullet_list(slide, left, top, width, height, items, body_size=14):
    """items: list of str (각 항목은 맨 앞에 • 포함하여 작성)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        # item 은 list of (text, opts) 형식 또는 str
        if isinstance(item, str):
            run = p.add_run()
            run.text = "•  " + item
            set_run_font(run, size=body_size, color=C_BLACK)
        else:
            # rich — 첫 run 앞에 불릿
            first = True
            for text, opts in item:
                run = p.add_run()
                run.text = ("•  " + text) if first else text
                set_run_font(
                    run,
                    size=opts.get("size", body_size),
                    bold=opts.get("bold", False),
                    color=opts.get("color", C_BLACK),
                )
                first = False
    return tb


def add_table(slide, left, top, width, height, headers, rows, font_size=11,
              header_bg=C_BLACK, header_fg=C_WHITE):
    n_cols = len(headers)
    n_rows = len(rows) + 1  # header + data

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 헤더 셀
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        set_run_font(run, size=font_size, bold=True, color=header_fg)

    # 데이터 셀
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            # 짝수행 배경
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BG_LIGHT
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # cell_text 가 tuple 이면 (text, bold, color) 형식
            if isinstance(cell_text, tuple):
                text, bold, color = cell_text
            else:
                text = cell_text
                bold = False
                color = C_BLACK
            run = p.add_run()
            run.text = text
            set_run_font(run, size=font_size, bold=bold, color=color)

    return table


def add_picture(slide, path, left, top, width=None, height=None):
    if not path.exists():
        print(f"경고: figure 없음 {path}")
        return None
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return pic


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, size=11, color=C_BLACK)


# =========================================================================
# 슬라이드 빌더
# =========================================================================

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank

    # -------- S1 표지 --------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_bg(slide, C_BLACK)

    # 상단 장식 줄
    decor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.4),
                                    SLIDE_W, Inches(0.08))
    decor.fill.solid()
    decor.fill.fore_color.rgb = C_BLUE
    decor.line.fill.background()

    # 제목
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(12), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Exqutor 벡터 카디널리티 추정의"
    set_run_font(r, size=36, bold=True, color=C_WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "단일 테이블 사각지대와 Skew-Aware Sampling"
    set_run_font(r2, size=36, bold=True, color=C_WHITE)

    # 부제
    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(12), Inches(0.6))
    tf2 = tb2.text_frame
    p3 = tf2.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Pivot A BERNOULLI 교체 + Pivot C KM20 Stratified + 3 데이터셋 5-seed CI 검증"
    set_run_font(r3, size=18, bold=False, color=C_GRAY_LIGHT)

    # 팀 정보
    tb3 = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(12), Inches(1))
    tf3 = tb3.text_frame
    p4 = tf3.paragraphs[0]
    r4 = p4.add_run()
    r4.text = "속도는벡터 · 연세대학교 컴퓨터과학과 2026-1 캡스톤 디자인"
    set_run_font(r4, size=15, color=C_GRAY_LIGHT)
    p5 = tf3.add_paragraph()
    r5 = p5.add_run()
    r5.text = "박세은(팀장) · 강재현 · 조현빈 · 이동욱   |   2026-04-28"
    set_run_font(r5, size=13, color=C_GRAY_MID)

    set_notes(slide,
        "안녕하세요. 속도는벡터 팀입니다. Exqutor 논문의 Adaptive Sampling 이 단일 테이블 "
        "벡터 범위 쿼리에서 사각지대를 가진다는 발견과, 이를 두 단계로 해소한 결과, "
        "그리고 3 데이터셋 교차 검증까지 발표하겠습니다.")

    # -------- S2 실험 타임라인 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 2 · 실험 타임라인",
               "4/3 합의 → 4/14 착수 → 4/14~15 본선 → 4/16~17 보강")

    # 4단계 타임라인 박스
    stages = [
        ("4/3 교수님 합의", "가설 H1\nskewness ↑\n= Q-error ↑", C_GRAY_MID),
        ("4/14 착수", "[Phase 1~5]\nH1 기각\n8 지표 negative\n→ data-side pivot", C_ORANGE),
        ("4/14~15 본선", "[Phase 4 / 6 / 7]\nPivot A 4 구간 p<0.001\nPivot C s=0.500 ★★★\nPhase 7 artifact", C_BLUE),
        ("4/16~17 보강", "[5-seed CI / Two-Level]\n3 데이터셋 외적 재현\nSIFT 2배 + HHI/CV\n→ 외적 타당성 CONSISTENT", C_GREEN),
    ]

    box_w = Inches(3.0)
    box_h = Inches(4.3)
    y = Inches(1.5)
    x = Inches(0.4)
    for i, (title, body, color) in enumerate(stages):
        # 헤더 박스
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x, y, box_w, Inches(0.6))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        htf = head.text_frame
        htf.margin_left = Inches(0.1)
        htf.margin_top = Inches(0.05)
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run()
        hr.text = title
        set_run_font(hr, size=14, bold=True, color=C_WHITE)

        # 본문
        body_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           x, Inches(2.1), box_w, Inches(3.7))
        body_box.fill.solid()
        body_box.fill.fore_color.rgb = C_GRAY_LIGHT
        body_box.line.color.rgb = color
        body_box.line.width = Pt(1)
        btf = body_box.text_frame
        btf.margin_left = Inches(0.15)
        btf.margin_top = Inches(0.15)
        btf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            if j == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.alignment = PP_ALIGN.LEFT
            bp.space_after = Pt(4)
            br = bp.add_run()
            br.text = line
            set_run_font(br, size=12, color=C_BLACK)

        # 화살표 (마지막 제외)
        if i < 3:
            arr_x = x + box_w + Inches(0.02)
            arr_y = Inches(3.2)
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          arr_x, arr_y, Inches(0.3), Inches(0.5))
            arr.fill.solid()
            arr.fill.fore_color.rgb = C_GRAY_MID
            arr.line.fill.background()

        x += box_w + Inches(0.35)

    # 하단 캡션
    add_textbox(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.8), [
        [("v2 의 4/16~17 보강 branch 가 외적 타당성 결손(Phase 7 artifact)을 해소하였고, "
          "Two-Level 분해로 ", {"size": 13, "color": C_BLACK}),
         ("공간 인식 기여를 수리적으로 분리", {"size": 13, "bold": True, "color": C_GREEN})],
    ])

    set_notes(slide,
        "4월 3일 교수님 합의에서 출발해 4월 17일까지의 실험 흐름입니다. 원래 가설은 기각되었지만 "
        "Exqutor의 구조적 사각지대를 발견했고, 두 가지 개선을 네이티브로 구현했습니다. "
        "4월 16~17일의 보강 측정에서 5-seed 신뢰구간과 3 데이터셋 교차 검증으로 외적 타당성을 확보했으며, "
        "Two-Level 분해로 공간 인식 기여를 정량화했습니다.")

    # -------- S3 연구 배경 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 3 · 연구 배경",
               "Exqutor 의 벡터 카디널리티 추정 문제 해결과 잠재적 약점")

    items = [
        [("pgvector / VBASE / DuckDB 의 벡터 조건 선택도 ", {}),
         ("고정 비율 추정", {"bold": True, "color": C_ORANGE}),
         (" (33.3% / 50% / 100%)", {})],
        [("실제 선택도: ", {}),
         ("0.001% ~ 90% 넓은 범위", {"bold": True}),
         (" → 옵티마이저 잘못된 실행 계획 선택", {})],
        [("Exqutor (arXiv:2512.09695v2, 2024) 의 두 해법", {"bold": True, "color": C_BLUE})],
        [("  — ECQO: HNSW 인덱스 있을 때 range query 직접 실행 (1~2 ms)", {"size": 13})],
        [("  — Adaptive Sampling: 인덱스 없을 때 모멘텀 기반 동적 샘플링", {"size": 13})],
        [("pgvector 최대 1000배, VBASE 10000배, DuckDB 1.5~37배 성능 향상", {})],
        [("", {})],
        [("본 연구의 공략점", {"size": 16, "bold": True, "color": C_BLUE})],
        [("  Adaptive Sampling 의 정확도 개선 여지 — 특히 skewed 분포에서", {"size": 14})],
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
                    items, body_size=15)

    set_notes(slide,
        "Exqutor 는 벡터 데이터베이스 카디널리티 추정 정확도 문제의 최신 해법입니다. "
        "기존 시스템은 고정 비율 추정이라 옵티마이저의 실행 계획이 잘못되는 경우가 많았고, "
        "Exqutor 는 ECQO 와 Adaptive Sampling 두 가지 해법으로 이를 해소했습니다. "
        "본 연구는 Adaptive Sampling 의 정확도 개선 여지를 공략합니다.")

    # -------- S4 연구 질문 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 4 · 연구 질문",
               "RQ1 Motivation / RQ2 Distribution-Aware / RQ3 Distribution-Agnostic")

    rqs = [
        ("RQ1", "Motivation 검증", "Exqutor Adaptive Sampling 은 실제로 정확도 한계를 가지는가? 한계의 구조는?", C_ORANGE),
        ("RQ2", "Distribution-Aware", "분포를 알 때, stratified sampling 이 카디널리티 추정 정확도를 얼마나 개선하는가?", C_BLUE),
        ("RQ3", "Distribution-Agnostic", "분포를 모를 때, 사전 학습 없이 KM20 의 공간 인식 효과를 얼마나 회수할 수 있는가?\nRecovery Rate = (방법X − RANDOM20) / (KM20 − RANDOM20)\n3 패러다임 (Offline Partition / Online Query-Adaptive / Weight-based) × 7 가지 방법", C_GREEN),
    ]
    y = Inches(1.4)
    for tag, title, desc, color in rqs:
        # 태그 박스
        tag_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(0.7), y, Inches(1.2), Inches(0.9))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = color
        tag_box.line.fill.background()
        ttf = tag_box.text_frame
        ttf.margin_left = Inches(0.1)
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = tag
        set_run_font(tr, size=22, bold=True, color=C_WHITE)

        # 본문 박스
        body_box = slide.shapes.add_textbox(Inches(2.1), y, Inches(10.8), Inches(0.9))
        btf = body_box.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = title
        set_run_font(br, size=15, bold=True, color=color)
        bp2 = btf.add_paragraph()
        bp2.alignment = PP_ALIGN.LEFT
        br2 = bp2.add_run()
        br2.text = desc
        set_run_font(br2, size=12, color=C_BLACK)

        y += Inches(1.5)

    # 하단 메모
    mem = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.7))
    mtf = mem.text_frame
    mp = mtf.paragraphs[0]
    mr = mp.add_run()
    mr.text = "본 발표 범위: RQ1 Motivation + RQ2 본선 + 5-seed CI + RQ3 설계안 (실험은 W7)"
    set_run_font(mr, size=13, bold=True, color=C_BLUE)

    set_notes(slide,
        "연구 질문은 3 단계입니다. RQ1 에서 정확도 한계를 정량 측정하고, RQ2 에서 분포를 알 때의 개선을 "
        "측정하며, RQ3 에서 분포를 모를 때 공간 인식 효과를 얼마나 회수할 수 있는지를 Recovery Rate "
        "프레임워크로 비교합니다. 오늘은 RQ1 과 RQ2 의 결과와 RQ3 의 설계안까지 다루고, RQ3 실험은 "
        "최종발표에서 보고합니다.")

    # -------- S5 H1 기각 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 5 · H1 기각과 8 지표 negative",
               "query-conditional layer 의 사전 식별 불가능성 정량 확인")

    # 왼쪽 bullet
    left_items = [
        [("실험", {"bold": True}),
         (": 1M subset × 96d DEEP, 100 query × 6 sel × 5 seed = 3000 측정", {})],
        [("가설 H1 (4/3)", {"bold": True}),
         (": Fisher γ 큰 그룹의 Q-error 가 작은 그룹 대비 2 배 악화", {})],
        [("결과", {"bold": True, "color": C_ORANGE}),
         (": 12 조합 모두 ratio ≈ 1.0, 24 조합 |Spearman ρ| < 0.2", {})],
        [("100 query 전부 Fisher γ < 0 (left-skewed)", {"size": 13})],
        [("tail ratio P99/P50 ∈ [1.07, 1.17] 좁은 범위", {"size": 13})],
        [("→ ", {}),
         ("고차원 거리 집중 (distance concentration)", {"bold": True, "color": C_BLUE})],
        [("", {})],
        [("로컬 4 지표 추가 검증 (Phase 5)", {"bold": True})],
        [("  k-NN entropy, PCA EVR1, KDE modality, NN clustering", {"size": 12})],
        [("  24 조합 전수 |ρ| < 0.2 — 글로벌 4 + 로컬 4 ", {}),
         ("= 8 지표 전수 negative", {"bold": True, "color": C_ORANGE})],
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(5.5),
                    left_items, body_size=13)

    # 오른쪽 figure
    add_picture(slide, FIG_RQ1 / "figure_6_phase5_heatmap.png",
                Inches(7.2), Inches(1.4), width=Inches(5.8))

    set_notes(slide,
        "4월 14일 1 차 실험에서 원래 가설을 기각했습니다. Fisher γ 등 글로벌 4 지표 모두 Q-error 와의 "
        "Spearman 상관이 0.2 미만이었고, 로컬 4 지표까지 확인했지만 24 조합 모두 |ρ| < 0.2 로 무효화되었습니다. "
        "학술적 의미는 고차원 벡터에서 query-side feature 로 layer 를 사전 정의하는 것이 불가능하다는 "
        "정량 negative result 이며, 이는 본 연구 motivation 의 한 축이 됩니다.")

    # -------- S6 Design Constraint --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 6 · Exqutor Design Constraint 발견",
               "단일 테이블 vector range query 는 Exqutor 의 명시되지 않은 사각지대")

    items = [
        [("(i) Hook Trigger 사각지대", {"size": 16, "bold": True, "color": C_ORANGE})],
        [("  vector.c line 243 if (table_count > 2) — 단일 테이블은 hook 경로 배제", {"size": 13})],
        [("  모든 query 가 PG default selectivity 1/3 (= 333,333 행) 로 fall-through", {"size": 13})],
        [("  → Exqutor Adaptive Sampling 함수 자체에 도달 불가", {"size": 13, "color": C_BLUE})],
        [("", {})],
        [("(ii) Plan Replacement 부작용", {"size": 16, "bold": True, "color": C_ORANGE})],
        [("  hook 을 >= 1 로 우회하면 outer query plan-tree 가 Sample Scan 으로 교체", {"size": 13})],
        [("  true = 100,000 의 SELECT count(*) 가 sample 안 부분 카운트 32 로 격하", {"size": 13})],
        [("  → 단일 테이블에서는 outer query 의 의미 자체가 깨짐", {"size": 13, "color": C_BLUE})],
        [("", {})],
        [("두 finding 은 Exqutor 원논문 본문에 명시되지 않음 — 소스 검증으로 첫 정량 확인", {"size": 13, "bold": True, "color": C_BLACK})],
        [("실무 함의: 이미지 검색, 추천, RAG retrieval = 단일 테이블 vector range query 의 전형", {"size": 13})],
        [("→ ", {"size": 15}),
         ("본 연구 motivation 의 새 첫 줄", {"size": 15, "bold": True, "color": C_GREEN})],
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.4), Inches(7.0), Inches(5.8),
                    items, body_size=13)

    # vector.c 스니펫 이미지 (오른쪽)
    add_picture(slide, FIG_RQ1 / "slide6_vector_c_snippet.png",
                Inches(7.7), Inches(1.5), width=Inches(5.5))

    set_notes(slide,
        "저희는 Python 재구현과 Exqutor 네이티브의 동치성을 검증하려는 과정에서 두 가지 design constraint 를 "
        "발견했습니다. 첫째, Exqutor 의 Adaptive Sampling hook 이 table_count > 2 조건으로 단일 테이블 "
        "vector range query 를 경로에서 배제합니다. 둘째, hook 을 한 줄 우회해도 outer query 의 plan-tree 가 "
        "Sample Scan 으로 교체되어 outer query 결과 자체가 부분 카운트로 격하됩니다. 두 finding 은 Exqutor 가 "
        "multi-table join 시나리오를 가정한 해법임을 보여주며, 단일 테이블 vector range query — 이미지 검색, "
        "추천, RAG retrieval 같은 실무 시나리오 — 는 명시적으로 배제된 사각지대입니다. 본 두 발견이 저희 연구의 "
        "새 motivation 첫 줄입니다.")

    # -------- S7 Pivot A + Pivot C --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 7 · 수정 노선: Pivot A + Pivot C 병합",
               "4/14 17:35 확정 — 두 sanitize 의 직교 ablation")

    # 왼쪽 Pivot A
    left_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.6), Inches(1.4),
                                       Inches(6.0), Inches(5.3))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = C_GRAY_LIGHT
    left_box.line.color.rgb = C_BLUE
    left_box.line.width = Pt(1.5)

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), [
        [("Pivot A — Block sampling bias 제거", {"size": 17, "bold": True, "color": C_BLUE})],
        [("", {})],
        [("vector.c L889 한 줄 교체:", {"size": 13, "bold": True})],
        [("  TABLESAMPLE SYSTEM(%f)", {"size": 12, "font": FONT_CODE, "color": C_ORANGE})],
        [("  → TABLESAMPLE BERNOULLI(%f)", {"size": 12, "font": FONT_CODE, "color": C_BLUE})],
        [("", {})],
        [("역할: Fair baseline 도구", {"size": 13})],
        [("  단독 기여가 아닌 Pivot C 의 측정 단위", {"size": 12, "color": C_GRAY_MID})],
        [("  block 단위 샘플링의 row-correlation bias 제거", {"size": 12, "color": C_GRAY_MID})],
        [("", {})],
        [("효과 (Phase 4): 4 selectivity 구간", {"size": 13})],
        [("  paired Wilcoxon p < 0.001", {"size": 13, "bold": True, "color": C_GREEN})],
    ])

    # 오른쪽 Pivot C
    right_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(6.8), Inches(1.4),
                                         Inches(6.0), Inches(5.3))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = C_GRAY_LIGHT
    right_box.line.color.rgb = C_GREEN
    right_box.line.width = Pt(1.5)

    add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.0), [
        [("Pivot C — data-side KM20 Stratified", {"size": 17, "bold": True, "color": C_GREEN})],
        [("", {})],
        [("vector.c +228 줄 구현:", {"size": 13, "bold": True})],
        [("  GUC 3 개 + stratum_sizes 캐시", {"size": 12, "font": FONT_CODE})],
        [("  estimate_cardinality_with_stratified_sampling", {"size": 11, "font": FONT_CODE})],
        [("", {})],
        [("데이터 측면:", {"size": 13, "bold": True})],
        [("  numpy mini-batch k-means K=20 (seed=42)", {"size": 12})],
        [("  stratum_id smallint + index 사전 부여", {"size": 12})],
        [("  → query 시점 overhead 0", {"size": 12, "color": C_BLUE})],
        [("", {})],
        [("Horvitz-Thompson 가중 추정:", {"size": 13, "bold": True})],
        [("  est = Σ (n_i × cnt_i / s_i)", {"size": 12, "font": FONT_CODE, "color": C_GREEN})],
    ])

    # 하단 메시지
    add_textbox(slide, Inches(0.6), Inches(6.9), Inches(12.2), Inches(0.4), [
        [("두 Pivot 은 직교적 → SYSTEM / BERNOULLI / STRATIFIED × dataset × selectivity 의 ablation matrix 로 각 기여 분리",
          {"size": 13, "bold": True, "color": C_BLACK})],
    ], align=PP_ALIGN.CENTER)

    set_notes(slide,
        "본 두 finding 위에서 두 가지 sanitize 를 병합한 수정 노선을 확정했습니다. Pivot A 는 Exqutor 의 "
        "TABLESAMPLE SYSTEM 을 BERNOULLI 로 한 줄 교체해 block sampling bias 를 제거하는 fair baseline "
        "도구입니다. Pivot C 는 data-side k-means K=20 으로 데이터를 partition 한 후 stratified sampling 을 "
        "native 로 구현합니다. 두 Pivot 은 직교적이므로 ablation matrix 로 각 기여를 분리 측정합니다.")

    # -------- S8 Phase 4 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 8 · Phase 4 결과 — Pivot A (BERNOULLI) Native 검증 ★",
               "4 selectivity 구간 모두 paired Wilcoxon p < 0.001")

    # 표
    headers_p4 = ["selectivity", "SYS med", "BERN med", "diff%", "p", "유의"]
    rows_p4 = [
        ["0.001", "2.597", "2.597", "+0.0", "0.596", "—"],
        ["0.010", "1.558", "1.299", "+20.0", "0.240", "—"],
        [("0.050", True, C_BLACK), "1.203", "1.195", "+0.7", ("0.0009", True, C_GREEN), ("★★", True, C_GREEN)],
        [("0.100", True, C_BLACK), "1.212", "1.132", ("+7.0", True, C_BLACK), ("< 0.001", True, C_GREEN), ("★★★", True, C_GREEN)],
        [("0.300", True, C_BLACK), "1.210", "1.079", ("+12.0", True, C_BLACK), ("< 0.001", True, C_GREEN), ("★★★", True, C_GREEN)],
        [("0.500", True, C_BLACK), "1.153", "1.052", ("+9.6", True, C_BLACK), ("< 0.001", True, C_GREEN), ("★★★", True, C_GREEN)],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(6.8), Inches(3.5),
              headers_p4, rows_p4, font_size=11)

    # 오른쪽 figure
    add_picture(slide, FIG_RQ1 / "figure_1_phase4_scatter.png",
                Inches(7.5), Inches(1.4), width=Inches(5.5))

    # 하단 결론
    add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12.5), Inches(1.8), [
        [("s=0.500 에서 median 기준 9.6% 개선, 78/100 query 가 SYSTEM 쪽 불리",
          {"size": 13, "color": C_BLACK})],
        [("vector.c 한 줄 sed 교체로 Exqutor 카디널리티 추정 정확도 개선 경로가 native 로 직접 확인됨",
          {"size": 14, "bold": True, "color": C_GREEN})],
        [("Block sampling bias 는 통계적으로 유의한 구조적 약점 — 이후 Pivot C 의 fair baseline 확보",
          {"size": 13, "color": C_BLUE})],
    ])

    set_notes(slide,
        "Pivot A 의 정량 검증입니다. 네이티브에서 4 selectivity 구간 — s=0.05, 0.1, 0.3, 0.5 — 모두 "
        "paired Wilcoxon p < 0.001 로 SYSTEM 의 q-error 가 BERNOULLI 보다 통계적으로 유의하게 큽니다. "
        "s=0.5 에서 median 기준 약 9.6% 개선으로 가장 큰 효과를 보였습니다. TABLESAMPLE SYSTEM 에 block 단위 "
        "샘플링의 구조적 bias 가 존재하며, 한 줄을 BERNOULLI 로 교체하는 것만으로 제거할 수 있다는 것이 "
        "native 로 직접 확인되었습니다.")

    # -------- S9 Phase 6 Step 4 + 5-seed CI --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 9 · Phase 6 Step 4 + 5-seed 신뢰구간 ★★★",
               "s=0.500 p=4.01e-05 Bonferroni 유일 생존 + 5-seed CI 하한 양수 확정")

    # 표
    headers_p6 = ["selectivity", "5-seed 평균", "95% CI", "유의한 seed"]
    rows_p6 = [
        [("50%", True, C_BLACK), ("+1.64%", True, C_GREEN), ("[+1.11, +2.18]", True, C_GREEN), "5/5"],
        [("30%", True, C_BLACK), ("+2.62%", True, C_GREEN), ("[+1.45, +3.80]", True, C_GREEN), "5/5 (p<0.001)"],
        [("10%", True, C_BLACK), ("+4.19%", True, C_GREEN), ("[+0.72, +7.66]", True, C_GREEN), "5/5 (p<0.004)"],
        ["5%", "+1.85%", "(단일)", "3/5"],
        ["1%", "+8.93%", "(단일)", "4/5"],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(7.3), Inches(3.3),
              headers_p6, rows_p6, font_size=12)

    # 오른쪽 figure
    add_picture(slide, FIG_RQ1 / "figure_2_phase6_box.png",
                Inches(8.0), Inches(1.4), width=Inches(5.0))

    # 하단 결론
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(12.5), Inches(2.1), [
        [("s=0.500 단일 seed: p = 4.01e-05 (전체 110 검정 중 Bonferroni 보정 후 유일 생존)",
          {"size": 13})],
        [("5-seed 재현: ", {"size": 14}),
         ("mean +1.64%, CI [+1.11, +2.18]", {"size": 14, "bold": True, "color": C_GREEN}),
         (" — 하한 양수로 real effect 확정", {"size": 14})],
        [("50% / 30% / 10% 세 구간 모두 95% CI 가 0 을 포함하지 않음 → 확률적 우연이 아닌 통계적 확정",
          {"size": 13, "color": C_BLUE})],
    ])

    set_notes(slide,
        "Pivot C 의 핵심 결과입니다. vector.c 에 +228 줄의 stratified sampling 을 직접 구현하고 5 seed 반복으로 "
        "측정했습니다. s=0.500 에서 mean +1.64%, 95% 신뢰구간 [+1.11, +2.18] 로 하한이 양수로 확정되었습니다. "
        "30% 와 10% 구간도 각각 [+1.45, +3.80], [+0.72, +7.66] 로 CI 하한 양수입니다. 3 구간 모두 5/5 seed 가 "
        "유의하며, p=4.01e-05 의 단일 seed 수치가 확률적 우연이 아닌 real effect 임이 통계적으로 확정되었습니다. "
        "+1.64% 가 작아 보이지만 이미 좋은 시스템 (Exqutor BERNOULLI) 위에서의 추가 개선입니다.")

    # -------- S10 External Validity I (SIFT) --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 10 · External Validity I — SIFT 1.5M 에서 효과 2배  ★ 신규",
               "D_target 재계산 후 5-seed 재측정 — 쏠림 정량(CV 68%)과 인과 연결")

    # 좌측 표
    headers_s = ["selectivity", "KM20 5-seed 평균", "95% CI", "유의 seed"]
    rows_s = [
        [("50%", True, C_BLACK), ("+3.07%", True, C_GREEN), ("[+2.66, +3.48]", True, C_GREEN), "5/5 (p<1e-8)"],
        [("5%", True, C_BLACK), ("+4.39%", True, C_GREEN), ("[+2.63, +6.15]", True, C_GREEN), "5/5"],
        ["1%", "−0.53%", "[−3.18, +2.11]", "1/5"],
    ]
    add_table(slide, Inches(0.4), Inches(1.4), Inches(6.5), Inches(2.1),
              headers_s, rows_s, font_size=11)

    # 좌측 수치 해설
    add_textbox(slide, Inches(0.4), Inches(3.7), Inches(6.5), Inches(3.2), [
        [("SIFT 1.5M (128d) — 더 쏠린 분포", {"size": 14, "bold": True, "color": C_ORANGE})],
        [("", {})],
        [("• DEEP 1M CI [+1.11, +2.18] 과 SIFT CI [+2.66, +3.48]", {"size": 12})],
        [("  ", {}),
         ("겹치지 않음 → 통계적으로 구별", {"size": 12, "bold": True, "color": C_GREEN})],
        [("", {})],
        [("• 쏠림 정량 (HHI / CV)", {"size": 13, "bold": True})],
        [("  DEEP CV 0.234 vs SIFT CV ", {"size": 12}),
         ("0.394 (68% 높음)", {"size": 12, "bold": True, "color": C_ORANGE})],
        [("", {})],
        [("• SIFT 최대 클러스터 14.8만 (기대값의 2배)", {"size": 12})],
        [("• 상위 2 클러스터 전체의 19% (기대 10%)", {"size": 12})],
    ])

    # 우측 figure
    add_picture(slide, FIG_RQ2 / "figure_8_cross_dataset_bar.png",
                Inches(7.2), Inches(1.4), width=Inches(5.9))

    set_notes(slide,
        "두 번째 신규 기여는 외적 타당성입니다. D_target 을 각 데이터셋별로 재계산한 뒤 SIFT 1.5M 과 DEEP 8M 을 "
        "5 seed 로 재측정했습니다. SIFT 에서 s=0.500 의 KM20 개선이 mean +3.07%, CI [+2.66, +3.48] 로 "
        "DEEP 의 [+1.11, +2.18] 과 겹치지 않아 효과가 약 2 배로 확대됨이 통계적으로 구별됩니다. 쏠림도 CV 가 "
        "DEEP 0.234 대비 SIFT 0.394 로 68% 높으며, SIFT 최대 클러스터가 기대값의 2 배에 달합니다. "
        "데이터가 더 쏠려 있을수록 공간 인식 샘플링의 가치가 크다는 본 연구 가설이 HHI/CV 정량 지표와 "
        "Q-error 개선폭 사이의 일관된 단조 관계로 뒷받침됩니다.")

    # -------- S11 External Validity II (Two-Level) --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 11 · External Validity II — KM20 vs RANDOM20 + Two-Level  ★ 신규",
               "개선 원인이 '20 등분' 이 아니라 '공간 구조를 아는 분할' 임을 증명")

    # 좌측 figure (gradient)
    add_picture(slide, FIG_RQ2 / "figure_7_selectivity_gradient.png",
                Inches(0.3), Inches(1.4), width=Inches(6.8))

    # 우측 표 — KM20 vs RANDOM20 격차
    headers_g = ["selectivity", "KM20 (CI)", "RANDOM20 (CI)", "격차 (Level 2)"]
    rows_g = [
        ["50%", "+1.64%", "+2.20%", "~0"],
        ["30%", "+2.62%", "+0.26%", ("+2.4%p", True, C_GREEN)],
        ["10%", "+4.19%", "+1.74%", ("+2.5%p", True, C_GREEN)],
        ["5%", "+1.85%", "+0.79%", "+1.1%p"],
        [("1%", True, C_BLACK), ("+8.93%", True, C_GREEN), ("−10.67%", True, C_ORANGE),
         ("+19.6%p", True, C_GREEN)],
    ]
    add_table(slide, Inches(7.3), Inches(1.4), Inches(5.8), Inches(3.2),
              headers_g, rows_g, font_size=11)

    # 우측 하단 Two-Level 요약
    add_textbox(slide, Inches(7.3), Inches(4.8), Inches(5.8), Inches(2.4), [
        [("Two-Level Decomposition", {"size": 14, "bold": True, "color": C_BLUE})],
        [("KM20 개선 = Level 1 (비례 배분) + Level 2 (공간 인식)", {"size": 12})],
        [("", {})],
        [("s=0.010 에서 ", {"size": 13}),
         ("Level 2 단독 +19.60%p", {"size": 13, "bold": True, "color": C_GREEN})],
        [("  (Level 1 −10.67%, Total +8.93%)", {"size": 12, "color": C_GRAY_MID})],
        [("", {})],
        [("→ selectivity 가 좁아질수록 공간 인식이 지배적", {"size": 12, "bold": True})],
    ])

    set_notes(slide,
        "세 번째 신규 기여는 Two-Level Decomposition 입니다. 같은 층화 샘플링 코드를 KM20 (공간 인식) 과 "
        "RANDOM20 (무작위 분할) 두 파티션으로 비교하여, 개선 원인이 20 등분이 아니라 공간 구조를 아는 분할 "
        "임을 증명했습니다. s=0.010 에서 KM20 은 +8.93% 개선, RANDOM20 은 −10.67% 악화로 격차가 19.6%p 에 "
        "달합니다. 이를 수학적으로 분해하면 Level 1 (비례 배분 효과, RANDOM20 개선분) 과 Level 2 (공간 인식 "
        "효과, KM20-RANDOM20 격차) 로 나뉩니다. s=0.010 에서 Level 2 단독으로 +19.60%p 를 기여하여 KM20 의 "
        "공간 인식이 개선의 원동력임이 수학적으로 입증되었습니다.")

    # -------- S12 Phase 7 후속 해소 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 12 · Phase 7 — 외적 확장 시도와 후속 해소",
               "artifact 발견·철회 + D_target 재계산 후 boundary 조건 확정")

    # 3단 박스
    boxes = [
        ("1차 시도 (4/15)", C_ORANGE,
         "8M / sift 128d 에 1M derive D_target\n그대로 적용\n→ actual selectivity 0.0001 로 이동\n\nplan_rows 기반 초기 분석:\nSTRAT 9.60 vs BERN 20810\n(8M, 2168× 극단)"),
        ("artifact 발견·철회", C_BLACK,
         "검증: STRAT plan_rows 100/100 query\n에서 정확히 20\n= stratum LIMIT 상수 artifact\n\nhook_est 재분석:\nSTRAT −26.5% (8M), −325% (sift)\n→ 극극소 selectivity boundary"),
        ("후속 해소 (4/16)", C_GREEN,
         "D_target 재계산 → §4.8 재측정:\nDEEP 8M +1.76% [+0.65, +2.86]\n= 1M 과 CONSISTENT\nSIFT +3.07% [+2.66, +3.48]\n= DEEP 대비 2배 효과\n\n→ 안정 하한 s ≥ 0.050 확정"),
    ]
    x = Inches(0.5)
    for title, color, body in boxes:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, Inches(1.4), Inches(4.0), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        btf = box.text_frame
        btf.margin_left = Inches(0.1)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = title
        set_run_font(br, size=16, bold=True, color=C_WHITE)

        body_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, Inches(2.2), Inches(4.0), Inches(4.2))
        body_b.fill.solid()
        body_b.fill.fore_color.rgb = C_GRAY_LIGHT
        body_b.line.color.rgb = color
        body_b.line.width = Pt(1)
        btf2 = body_b.text_frame
        btf2.margin_left = Inches(0.2)
        btf2.margin_top = Inches(0.15)
        btf2.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            if j == 0:
                bp2 = btf2.paragraphs[0]
            else:
                bp2 = btf2.add_paragraph()
            bp2.alignment = PP_ALIGN.LEFT
            bp2.space_after = Pt(3)
            br2 = bp2.add_run()
            br2.text = line
            set_run_font(br2, size=12, color=C_BLACK)

        # 화살표
        if x < Inches(8.5):
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          x + Inches(4.05), Inches(3.9),
                                          Inches(0.3), Inches(0.5))
            arr.fill.solid()
            arr.fill.fore_color.rgb = C_GRAY_MID
            arr.line.fill.background()

        x += Inches(4.35)

    # 하단 요약
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.5), Inches(0.6), [
        [("Phase 7 의 boundary negative finding 은 Pivot C 효과 범위 (s ≥ 0.050) 의 정량 확정 — 학술적 가치 보존",
          {"size": 13, "bold": True, "color": C_BLACK})],
    ], align=PP_ALIGN.CENTER)

    set_notes(slide,
        "Phase 7 의 초기 외적 확장 시도에서는 D_target 을 재계산하지 않아 actual selectivity 가 극극소 영역으로 "
        "이동했고, 초기 plan_rows 기반 분석에 artifact 가 있음을 발견하여 철회했습니다. 이후 D_target 을 각 "
        "데이터셋별로 재계산한 §4.8 의 후속 측정에서 DEEP 8M 은 1M 과 CONSISTENT, SIFT 는 2 배 효과로 정상 "
        "selectivity 영역의 외적 타당성이 회복되었습니다. 남는 경계 조건은 s ≥ 0.050 에서만 층화 샘플링 효과가 "
        "안정적으로 발현된다는 것이며, 이는 Exqutor 기본 sample_size=385 의 구조적 하한의 결과로 정량적으로 "
        "확정되었습니다.")

    # -------- S13 한계 + 로드맵 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 13 · 한계와 남은 작업 / 최종발표 로드맵",
               "honest report 4건 + W5~W8 로드맵")

    # 상단 — 한계
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.5), Inches(2.4), [
        [("중간발표까지의 한계 (honest report)", {"size": 15, "bold": True, "color": C_ORANGE})],
        [("(a) 단일 dataset family — 96~128d 좁은 영역 (wiki 768d 는 W5 검증)", {"size": 13})],
        [("(b) SIFT mid-selectivity (s=0.10/0.30) CI 결손 — 서버 ", {"size": 13}),
         ("unrecognized node type", {"size": 12, "font": FONT_CODE, "color": C_ORANGE}),
         (" 오류 디버깅 필요 (W5)", {"size": 13})],
        [("(c) boundary condition: s ≥ 0.050 에서만 KM20 효과 안정 (W5 의 sample_size sensitivity 로 완화 탐색)",
          {"size": 13})],
        [("(d) ", {"size": 13}),
         ("vector.c", {"size": 12, "font": FONT_CODE}),
         (" hook 우회 수정의 reviewer 관점 위험 — 자문 회신 후 narrative 조정", {"size": 13})],
    ])

    # 하단 — 로드맵 박스 4개
    weeks = [
        ("W5\n4/28~5/4", "최적화 + wiki 검증",
         "per-stratum BERNOULLI\nLayer sweep (K=30/40/50)\nwiki 768d 검증\nSIFT mid-sel 보충", C_BLUE),
        ("W6\n5/5~5/11", "전면 ablation",
         "3 dataset × 3 mode\n× 6 sel × 10 seed\n× 500 query\n= 270k 측정", C_GREEN),
        ("W7\n5/12~5/18", "RQ3 Recovery Rate",
         "7-way 비교\n3 패러다임\nHilbert Curve\n최초 적용", C_ORANGE),
        ("W8\n5/19~5/25", "최종 산출물",
         "최종보고서\n포스터\n전시회 데모\n리허설 2회", C_GRAY_MID),
    ]
    x = Inches(0.5)
    for tag, title, body, color in weeks:
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x, Inches(4.4), Inches(3.0), Inches(0.9))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        htf = head.text_frame
        htf.margin_left = Inches(0.1)
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        htf.word_wrap = True
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run()
        hr.text = tag + " · " + title
        set_run_font(hr, size=12, bold=True, color=C_WHITE)

        body_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, Inches(5.4), Inches(3.0), Inches(2.0))
        body_b.fill.solid()
        body_b.fill.fore_color.rgb = C_GRAY_LIGHT
        body_b.line.color.rgb = color
        body_b.line.width = Pt(1)
        btf = body_b.text_frame
        btf.margin_left = Inches(0.15)
        btf.margin_top = Inches(0.1)
        btf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            if j == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.space_after = Pt(2)
            br = bp.add_run()
            br.text = line
            set_run_font(br, size=11, color=C_BLACK)

        x += Inches(3.15)

    set_notes(slide,
        "현재 네 가지 한계를 인지하고 있습니다: 단일 dataset family, SIFT mid-selectivity CI 결손, "
        "s ≥ 0.050 boundary, vector.c 우회 수정의 reviewer 위험. 이를 4 주 로드맵으로 해소합니다. "
        "W5 에 최적화와 wiki 검증, W6 에 3 데이터셋 × 500 query × 10 seed 전면 ablation, "
        "W7 에 RQ3 Recovery Rate 7-way 비교, W8 에 최종보고서와 포스터를 완성합니다. 특히 RQ3 는 세 "
        "패러다임의 7 가지 방법을 체계적으로 비교하며, 공간 DB 의 Hilbert Curve 기법을 벡터 카디널리티 "
        "추정에 최초로 적용합니다.")

    # -------- S14 결론 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 14 · 결론",
               "4 가지 기여 — 구조적 발견 / 방법론 구현 / 정량 효과 / 경계 조건")

    contribs = [
        ("기여 1 — 구조적 발견", C_ORANGE,
         "Exqutor 의 단일 테이블 사각지대 (Hook trigger + plan replacement + 5 design constraint) +\nquery-side feature 사전 식별 불가능성 (8 지표 × 24 조합 전수 |ρ| < 0.2)"),
        ("기여 2 — 방법론 구현", C_BLUE,
         "Block bias 1 줄 교체 (Pivot A) + data-side KM20 stratified sampling native 구현 (+228 줄, Pivot C)\nPython counterfactual 검증 → native 재현의 2 단계 validation pipeline"),
        ("기여 3 — 정량 효과 (복수 유효 anchor)", C_GREEN,
         "Phase 4 SYSTEM→BERNOULLI: 4 selectivity 구간 p<0.001\nPhase 6 Step 4 KM20 native: s=0.500 p=4.01e-05 + 5-seed CI [+1.11, +2.18]\n3 데이터셋 외적 재현 (DEEP 1M/8M CONSISTENT, SIFT 2배 CI 분리)\nTwo-Level 분해: Level 2 단독 +19.60%p (s=0.010)\nHHI/CV 정량: CV 68% 격차 = KM20 2배 격차의 직접 원인"),
        ("기여 4 — 경계 조건", C_GRAY_MID,
         "s ≥ 0.050 안정 하한 확정 — Phase 7 boundary negative finding 의 학술적 가치"),
    ]

    y = Inches(1.3)
    for title, color, body in contribs:
        # 왼쪽 색띠
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.5), y, Inches(0.15), Inches(1.3))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()

        # 본문
        tb = slide.shapes.add_textbox(Inches(0.75), y, Inches(12.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run_font(r, size=15, bold=True, color=color)
        for line in body.split("\n"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(1)
            r2 = p2.add_run()
            r2.text = line
            set_run_font(r2, size=12, color=C_BLACK)

        y += Inches(1.45)

    set_notes(slide,
        "결론입니다. 본 연구의 기여는 네 가지입니다. Exqutor 의 사각지대 발견, data-side KM20 stratified 의 "
        "native 구현, 5-seed CI + 3 데이터셋 + Two-Level 분해의 복수 유효 anchor, 그리고 s ≥ 0.050 boundary "
        "조건의 정량 확정입니다. 더 쏠린 SIFT 에서 KM20 효과가 DEEP 의 2 배로 확대됨을 HHI/CV 정량 지표로 "
        "인과적으로 설명했습니다. 최종발표에서 wiki 768d 검증과 RQ3 Recovery Rate 7-way 비교로 완성하겠습니다. "
        "감사합니다.")

    # -------- S15 부록 --------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Slide 15 (부록) · Appendix",
               "추가 Design Constraint iii~v + RQ3 7-way 설계")

    # 상단 — Design Constraint iii~v
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.5), Inches(1.7), [
        [("추가 Design Constraint (iii~v)", {"size": 14, "bold": True, "color": C_ORANGE})],
        [("(iii) q-error inf 발산 — single-table hook 우회 시 이중 sampling 으로 첫 sample 매칭 0 건", {"size": 12})],
        [("(iv) sample_size NaN 발산 — 50 query 시퀀스 중 Adaptive update 가 NaN 으로 발산", {"size": 12})],
        [("(v) Adaptive update path SIGSEGV — ", {"size": 12}),
         ("update_sample_size = off", {"size": 11, "font": FONT_CODE, "color": C_ORANGE}),
         (" GUC 로 회피", {"size": 12})],
        [("→ 셋 모두 Adaptive Sampling loop 가 multi-table only design intent 를 가정한 결과", {"size": 12, "color": C_BLUE})],
    ])

    # 하단 — RQ3 7-way 표
    add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12.5), Inches(0.4),
                [[("RQ3 Recovery Rate 프레임워크 — 7 가지 방법",
                   {"size": 14, "bold": True, "color": C_GREEN})]])

    headers_rq3 = ["방법", "패러다임", "기대 Recovery", "사전 학습"]
    rows_rq3 = [
        ["C. Random Projection", "Offline", "10~40%", "사영 1회"],
        ["A. LSH", "Offline", "30~60%", "해시 1회"],
        ["E. Hilbert Curve", "Offline", "20~60%", "index 1회"],
        ["F. Mini-batch K-means", "Offline", "75~95%", "1~5% 데이터"],
        ["G. Distance-Shell", "Online", "25~50%", "불필요"],
        ["B. KDE-pilot", "Online", "50~80%", "불필요"],
        ["H. Importance Sampling", "Weight-based", "30~70%", "불필요"],
    ]
    add_table(slide, Inches(0.5), Inches(3.6), Inches(12.5), Inches(3.2),
              headers_rq3, rows_rq3, font_size=11)

    # 하단 설명
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.4), [
        [("Recovery Rate = (방법X − RANDOM20) / (KM20 − RANDOM20).  1.0 = KM20 동등, 0.0 = RANDOM20 동등.",
          {"size": 11, "color": C_BLACK})],
    ])

    set_notes(slide,
        "부록 슬라이드입니다. Design constraint 5 가지 중 본문에 담지 못한 3 가지는 모두 Adaptive Sampling "
        "loop 의 수치 안정성 관련 기술 발견이며, 본 연구는 update_sample_size = off 로 sample_size=385 fixed "
        "한 채 sampling method 효과만 분리 측정했습니다. RQ3 의 7 가지 방법은 세 패러다임으로 분류되며, "
        "Recovery Rate 지표로 사전 학습 없이 KM20 의 공간 인식 효과를 얼마나 회수하는지 정량 비교합니다.")

    # -------- 저장 --------
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"✓ {OUT} ({prs.slide_width}, {prs.slide_height})")
    print(f"  슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
