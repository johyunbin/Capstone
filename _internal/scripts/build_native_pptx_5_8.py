#!/usr/bin/env python3
"""Native PPTX builder for 5/8 회의 자료 — fully editable text frames, shapes, tables.

Reproduces the Academic v3 visual style (white BG, navy accent, JetBrains Mono eyebrow,
numbered ink badge, header pills, table heatmap) as native PPTX elements.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---------------------------------------------------------------------------
# Design tokens — exact match to academic-deck/index.html :root + Slides.jsx C{}
# ---------------------------------------------------------------------------
# Brand colors (Slides.jsx C{} object, lines 4-12)
NAVY = RGBColor(0x1B, 0x3D, 0xAD)         # --brand-navy = #1B3DAD
NAVY_DEEP = RGBColor(0x14, 0x30, 0x7F)    # --brand-navy-deep = #14307F
BLUE = RGBColor(0x4A, 0x7B, 0xD8)         # --brand-blue = #4A7BD8
BLUE_SOFT = RGBColor(0xE4, 0xEC, 0xF8)    # --brand-blue-soft = #E4ECF8
RED = RGBColor(0xE0, 0x3A, 0x3A)          # --brand-red = #E03A3A
RED_SOFT = RGBColor(0xFB, 0xE3, 0xE3)     # --brand-red-soft = #FBE3E3
GREEN = RGBColor(0x2A, 0x9D, 0x6E)        # --green = #2A9D6E
GOLD = RGBColor(0xD9, 0xA5, 0x3B)         # --gold = #D9A53B

INK = RGBColor(0x0B, 0x0F, 0x1C)          # --ink = #0B0F1C
INK_SOFT = RGBColor(0x1A, 0x22, 0x38)     # --ink-soft = #1A2238

# Gray scale
GRAY_50 = RGBColor(0xFA, 0xFB, 0xFD)
GRAY_100 = RGBColor(0xF2, 0xF4, 0xF8)
GRAY_200 = RGBColor(0xDE, 0xE2, 0xEC)     # --gray-200 = LINE
GRAY_300 = RGBColor(0xC9, 0xCD, 0xD8)
GRAY_400 = RGBColor(0xA4, 0xAB, 0xBC)     # --line-strong
GRAY_500 = RGBColor(0x6E, 0x78, 0x91)     # mono labels / muted body
GRAY_600 = RGBColor(0x4B, 0x54, 0x70)     # body text muted
GRAY_700 = RGBColor(0x2E, 0x36, 0x50)

# Aliases for backward compatibility with existing slide bodies
LINE = GRAY_200
MUTED = GRAY_600
SOFT = GRAY_500
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)      # cards are white in HTML; .card.tint uses GRAY_50
BG_CARD_TINT = GRAY_50
BG_NAVY_TINT = BLUE_SOFT                  # used to mimic .card { background: blueSoft }
BG_HEADER = RGBColor(0xFA, 0xFB, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Heatmap colors — keep custom palette (used for our cell heatmap)
HM_VBAD = RGBColor(0x1F, 0x5B, 0x3A)   # dark green = strongest improve
HM_BAD = RGBColor(0x2F, 0x7D, 0x4A)    # green = improve
HM_MID = RGBColor(0xF7, 0xD2, 0x7A)    # mid yellow
HM_LIGHT = RGBColor(0xFF, 0xE8, 0xB3)  # light yellow
HM_NEUT = RGBColor(0xF3, 0xF5, 0xF8)
HM_WARM = RGBColor(0xF4, 0xA9, 0x9A)   # light red
HM_HOT = RGBColor(0xD9, 0x62, 0x56)    # red
HM_VBAD_TXT = WHITE
HM_BAD_TXT = WHITE
HM_MID_TXT = RGBColor(0x5E, 0x45, 0x00)
HM_LIGHT_TXT = RGBColor(0x6E, 0x4D, 0x00)
HM_WARM_TXT = RGBColor(0x5C, 0x1F, 0x0E)

# Pill backgrounds — match Slides.jsx S1 .pill.navy/.pill.red exactly
# .pill.navy { color: navy; border-color: rgba(27,61,173,0.4); background: rgba(27,61,173,0.05); }
# rgba(27,61,173,0.05) on white ≈ #F5F6FC; rgba(.,0.4) ≈ #A8B5DD
PILL_NAVY_BG = RGBColor(0xF5, 0xF6, 0xFC)
PILL_NAVY_BORDER = RGBColor(0xA8, 0xB5, 0xDD)
# .pill.red { background: rgba(224,58,58,0.06); border-color: rgba(224,58,58,0.4); }
PILL_RED_BG = RGBColor(0xFD, 0xF2, 0xF2)
PILL_RED_BORDER = RGBColor(0xEC, 0xA8, 0xA8)

# Tag colors (used for inline navy/grn/red tag chips in slides)
TAG_GRN_BG = RGBColor(0xE8, 0xF6, 0xEC)
TAG_GRN_BORDER = RGBColor(0xB8, 0xD8, 0xC1)
TAG_RED_BG = RGBColor(0xFF, 0xEF, 0xEF)
TAG_RED_BORDER = PILL_RED_BORDER
TAG_LINE_BG = WHITE

# Fonts — installed 5/7: Pretendard / Inter / JetBrains Mono (sub-agent task #1)
# HTML reference: "Pretendard Variable" / "Inter" / "JetBrains Mono"
KOR_FONT = "Pretendard"          # matches HTML font-sans (한글 + 영문 sans)
KOR_FONT_BOLD = "Pretendard"
MONO_FONT = "JetBrains Mono"     # mono (eyebrow / labels / page numbers)
NUM_FONT = "Inter"               # big numerals (font-num)

# 16:9 = 13.333 × 7.5 inches  (= 1280 × 720 px @ 96 dpi)
# HTML viewport: 1280 × 720, a-pad: 56 72 80 80 (top right bottom left)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.833)        # 80px = 0.833" (left/bottom padding)
MARGIN_RIGHT = Inches(0.75)     # 72px = 0.75" (right padding)
HEADER_Y = Inches(0.583)        # 56px = 0.583" top
FOOTER_Y_TOP = Inches(7.1)      # bottom rule area

OUT_PATH = Path("/Users/hyunbin/Capstone/submission/_drafts/속도는벡터_5월8일회의_v1.pptx")

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_run(run, text, font=KOR_FONT, size=14, bold=False, color=INK, italic=False,
            spc_em=None):
    """Set a run with optional letter-spacing.
    spc_em: float (em). e.g. -0.06 for tight headings (matches HTML letter-spacing).
            HTML em is relative to font size; PPTX `spc` attribute is in 1/100 pt.
    """
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    # set East Asian font
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.find(qn('a:ea'))
    if eastAsia is None:
        eastAsia = etree.SubElement(rPr, qn('a:ea'))
    eastAsia.set('typeface', font)
    # letter-spacing — em → 1/100 pt
    if spc_em is not None:
        spc_pt = float(spc_em) * float(size)
        spc_centi = int(round(spc_pt * 100))
        rPr.set('spc', str(spc_centi))


def add_text_box(slide, x, y, w, h, *, fill=None, line=None, line_w=None, line_dash=None):
    """Add an empty text box / rectangle. Returns the shape; caller fills tf manually."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = line_w
        if line_dash is not None:
            shape.line.dash_style = line_dash
    shape.shadow.inherit = False
    # Reset shape properties
    shape.text_frame.margin_left = Inches(0.1)
    shape.text_frame.margin_right = Inches(0.1)
    shape.text_frame.margin_top = Inches(0.05)
    shape.text_frame.margin_bottom = Inches(0.05)
    shape.text_frame.word_wrap = True
    return shape


def add_simple_text(slide, x, y, w, h, text, *, font=KOR_FONT, size=14, bold=False,
                    color=INK, align=None, anchor=None, spc_em=None):
    """Quick text box with a single run of text."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if anchor is not None:
        tf.vertical_anchor = anchor
    run = p.add_run()
    set_run(run, text, font=font, size=size, bold=bold, color=color, spc_em=spc_em)
    return box


def add_runs_text(slide, x, y, w, h, runs, *, align=None, anchor=None, line_spacing=None):
    """Multi-run text box.
    runs: list of dicts: {text, font?, size?, bold?, color?, newline?}
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for i, r in enumerate(runs):
        if r.get("newline"):
            p = tf.add_paragraph()
            if align is not None:
                p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            if r.get("text"):
                run = p.add_run()
                set_run(run, r["text"], font=r.get("font", KOR_FONT),
                        size=r.get("size", 13), bold=r.get("bold", False),
                        color=r.get("color", INK),
                        spc_em=r.get("spc_em"))
        else:
            run = p.add_run()
            set_run(run, r["text"], font=r.get("font", KOR_FONT),
                    size=r.get("size", 13), bold=r.get("bold", False),
                    color=r.get("color", INK),
                    spc_em=r.get("spc_em"))
    return box


def add_navy_stripe(slide, x=None, y=None, w=None, h=None):
    """The .a-stripe element — small navy vertical bar near the top-left.
    HTML: position: absolute; top: 56px; left: 60px; width: 5px; height: 60px;
    """
    if x is None:
        x = Inches(0.625)   # 60px
    if y is None:
        y = Inches(0.583)   # 56px
    if w is None:
        w = Inches(0.052)   # 5px
    if h is None:
        h = Inches(0.625)   # 60px
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_header(slide, page=None, total=14):
    """Top header: navy stripe + eyebrow + two pill badges (Slides.jsx S1 cover style)."""
    # The .a-stripe near top-left
    add_navy_stripe(slide)

    # Left: eyebrow (CAPSTONE 2026 …)
    add_simple_text(slide, MARGIN_X, Inches(0.6),
                    Inches(7.0), Inches(0.3),
                    "캡스톤 2026 · 속도는벡터",
                    font=MONO_FONT, size=9, bold=False, color=GRAY_500)

    # Right: pill badges (cover-style)
    # Slides.jsx .pill: padding: 4px 9px; font-size: 10px; mono; letter-spacing 0.1em
    pill_y = Inches(0.55)
    pill_h = Inches(0.27)

    # BDAI LAB pill — rightmost
    pill2_w = Inches(0.95)
    pill2_x = SLIDE_W - MARGIN_RIGHT - pill2_w
    p2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pill2_x, pill_y, pill2_w, pill_h)
    p2.fill.solid(); p2.fill.fore_color.rgb = PILL_RED_BG
    p2.line.color.rgb = PILL_RED_BORDER; p2.line.width = Pt(0.75)
    p2.shadow.inherit = False
    tf2 = p2.text_frame
    tf2.margin_left = Inches(0.05); tf2.margin_right = Inches(0.05)
    tf2.margin_top = Inches(0.01); tf2.margin_bottom = Inches(0.01)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2p = tf2.paragraphs[0]; p2p.alignment = PP_ALIGN.CENTER
    set_run(p2p.add_run(), "BDAI LAB", font=MONO_FONT, size=9, bold=False, color=RED)

    # YONSEI · CSE pill — to the left of BDAI
    pill1_w = Inches(1.15)
    pill1_x = pill2_x - pill1_w - Inches(0.08)   # gap: 8px
    p1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pill1_x, pill_y, pill1_w, pill_h)
    p1.fill.solid(); p1.fill.fore_color.rgb = PILL_NAVY_BG
    p1.line.color.rgb = PILL_NAVY_BORDER; p1.line.width = Pt(0.75)
    p1.shadow.inherit = False
    tf1 = p1.text_frame
    tf1.margin_left = Inches(0.05); tf1.margin_right = Inches(0.05)
    tf1.margin_top = Inches(0.01); tf1.margin_bottom = Inches(0.01)
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1p = tf1.paragraphs[0]; p1p.alignment = PP_ALIGN.CENTER
    set_run(p1p.add_run(), "YONSEI · CSE", font=MONO_FONT, size=9, bold=False, color=NAVY)


def add_footer(slide, page=None, total=14):
    """The .a-footer element — bottom-edge mono caption, uppercase 0.14em.
    HTML: position: absolute; left: 80px; right: 80px; bottom: 22px;
          font-size: 10px; color: var(--gray-500); letter-spacing: 0.14em; uppercase
    """
    # The .a-footer is at bottom: 22px → y ≈ 7.5 - 0.229 - 0.18 ≈ 7.18"
    foot_y = Inches(7.21)

    # Left text
    add_simple_text(slide, MARGIN_X, foot_y,
                    Inches(6.0), Inches(0.22),
                    "속도는벡터 · CAPSTONE 2026",
                    font=MONO_FONT, size=8, color=GRAY_500)

    # Right text — page number / date
    if page is not None:
        page_str = f"중간 측정 결과 · 2026.05.08    {page:02d} / {total:02d}"
        right_w = Inches(5.0)
        add_simple_text(slide, SLIDE_W - MARGIN_RIGHT - right_w, foot_y,
                        right_w, Inches(0.22),
                        page_str,
                        font=MONO_FONT, size=8, color=GRAY_500,
                        align=PP_ALIGN.RIGHT)


def add_numbered_badge(slide, x, y, num_text, *, size_in=None):
    """The .a-num element — 44x44 ink-filled square with white mono numerals.
    HTML: width: 44px; height: 44px; background: ink; font-num 18px / 700.
    """
    if size_in is None:
        size_in = Inches(0.46)   # 44px
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size_in, size_in)
    badge.fill.solid(); badge.fill.fore_color.rgb = INK
    badge.line.fill.background()
    badge.shadow.inherit = False
    btf = badge.text_frame
    btf.margin_left = Inches(0); btf.margin_right = Inches(0)
    btf.margin_top = Inches(0); btf.margin_bottom = Inches(0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    set_run(bp.add_run(), num_text, font=NUM_FONT, size=14, bold=True, color=WHITE)
    return badge


def add_lead_heading(slide, y, badge_text, heading, subhead):
    """The .a-head element — .a-num badge + .a-title (30px/800) + .a-eyebrow (11px mono).
    The HTML has gap: 18px between badge and title block, badge has margin-top: 2px.
    """
    # Black square badge — 44×44px
    if badge_text:
        badge_x = MARGIN_X
        badge_y = y + Inches(0.02)   # margin-top: 2px
        add_numbered_badge(slide, badge_x, badge_y, badge_text)
        title_x = badge_x + Inches(0.46) + Inches(0.19)   # 44 + 18 px gap
    else:
        title_x = MARGIN_X

    title_w = SLIDE_W - title_x - MARGIN_RIGHT - Inches(2.5)   # leave room for meta on right

    # Heading text — .a-title { font-size: 30px; weight: 800; line-height: 1.12 }
    # 30px ≈ 22.5pt; we want crisp and not too tall in 16:9 — spc_em -0.03 (tracked tight)
    add_simple_text(slide, title_x, y - Inches(0.04),
                    title_w, Inches(0.55),
                    heading,
                    font=KOR_FONT, size=22, bold=True, color=INK,
                    anchor=MSO_ANCHOR.TOP, spc_em=-0.03)

    # Subheading mono — .a-eyebrow { font-size: 11px; mono; 0.18em letter-spacing }
    add_simple_text(slide, title_x, y + Inches(0.46),
                    title_w, Inches(0.3),
                    subhead,
                    font=MONO_FONT, size=8, color=GRAY_500, spc_em=0.18)

    # Add a thin rule line below — .a-rule
    rule_y = y + Inches(0.85)
    ln = slide.shapes.add_connector(1, MARGIN_X, rule_y,
                                    SLIDE_W - MARGIN_RIGHT, rule_y)
    ln.line.color.rgb = LINE
    ln.line.width = Pt(0.6)


def add_card(slide, x, y, w, h, *, fill=BG_CARD, border=LINE, dashed=False):
    """The .card element.
    HTML: border: 1px solid var(--line); background: white; border-radius: 3px; padding: 18px 20px;
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.025   # 3px corner ≈ very small radius on a typical card
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.shadow.inherit = False
    tf = shape.text_frame
    # padding: 18px 20px → 0.21" top/bottom, 0.21" left/right
    tf.margin_left = Inches(0.21); tf.margin_right = Inches(0.21)
    tf.margin_top = Inches(0.19); tf.margin_bottom = Inches(0.19)
    tf.word_wrap = True
    return shape


def add_card_navy(slide, x, y, w, h):
    """Card with .navy-top — 3px navy top border instead of full navy bar.
    HTML: border-top: 3px solid var(--brand-navy); rest of border same as .card.
    """
    # Card body — same as regular card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.025
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(0.21); tf.margin_right = Inches(0.21)
    tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.19)   # 3px extra top for the navy band
    tf.word_wrap = True

    # Top navy band — 3px tall along the top
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.04))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    return shape


def add_card_red_top(slide, x, y, w, h):
    """Card with .red-top — 3px red top border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.025
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(0.21); tf.margin_right = Inches(0.21)
    tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.19)
    tf.word_wrap = True

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.04))
    band.fill.solid(); band.fill.fore_color.rgb = RED
    band.line.fill.background()
    band.shadow.inherit = False
    return shape


def card_set_content(shape, eyebrow=None, heading=None, body_runs=None,
                     dash_items=None, body_size=10, head_size=14,
                     eyebrow_color=None):
    """Fill a card shape with eyebrow (.label-mono) / heading / body / dash list.
    HTML reference:
      .label-mono { mono 11px / gray-500 / 0.18em / uppercase }
      heading: typically 16px/700, color: ink
      body: 13~14px / gray-600 / 1.55 line-height
      dash square = .sq-bullet (8x8px navy box, marginRight 10px)
    """
    if eyebrow_color is None:
        eyebrow_color = GRAY_500   # HTML .label-mono default
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    first = True
    if eyebrow:
        if first:
            first = False
        else:
            p = tf.add_paragraph()
        # eyebrow uses 0.18em letter-spacing (HTML .label-mono)
        set_run(p.add_run(), eyebrow, font=MONO_FONT, size=8, bold=False,
                color=eyebrow_color, spc_em=0.18)
    if heading:
        if first:
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(5)
        # card heading: tracked tight (-0.03em) for typographic consistency
        set_run(p.add_run(), heading, font=KOR_FONT, size=head_size, bold=True,
                color=INK, spc_em=-0.03)
    if body_runs:
        if first:
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(5)
        for r in body_runs:
            run = p.add_run()
            set_run(run, r["text"], font=r.get("font", KOR_FONT),
                    size=r.get("size", body_size),
                    bold=r.get("bold", False),
                    color=r.get("color", GRAY_600))
    if dash_items:
        for item in dash_items:
            p = tf.add_paragraph()
            p.space_before = Pt(3)
            # Use a small bullet square (HTML .sq-bullet) — render as square unicode + space
            run_d = p.add_run()
            set_run(run_d, "▪ ", font=KOR_FONT, size=body_size, bold=True, color=NAVY)
            if isinstance(item, str):
                run_t = p.add_run()
                set_run(run_t, item, font=KOR_FONT, size=body_size, color=GRAY_600)
            else:
                # list of runs
                for r in item:
                    run_t = p.add_run()
                    set_run(run_t, r["text"], font=r.get("font", KOR_FONT),
                            size=r.get("size", body_size),
                            bold=r.get("bold", False),
                            color=r.get("color", GRAY_600))


def add_table_styled(slide, x, y, w, h, headers, rows,
                     col_widths=None, header_fill=BG_HEADER,
                     cell_fills=None, cell_text_colors=None,
                     header_size=8.5, cell_size=10, mono_cols=None,
                     row_h=Inches(0.32)):
    """Build a styled table.
    headers: list of strings
    rows: list of lists of strings
    cell_fills: 2D list (rows × cols) of RGBColor or None
    cell_text_colors: same shape, RGBColor or None
    mono_cols: set of column indices to render in mono font
    """
    cols = len(headers)
    nrows = len(rows) + 1  # +1 for header
    table_shape = slide.shapes.add_table(nrows, cols, x, y, w, h)
    table = table_shape.table
    if col_widths is not None:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    mono_cols = mono_cols or set()

    # Header row
    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if i in mono_cols:
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        set_run(p.add_run(), hdr,
                font=MONO_FONT, size=header_size, bold=True, color=SOFT)

    # Body rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            fill = (cell_fills[r][c] if cell_fills else None)
            txt_color = (cell_text_colors[r][c] if cell_text_colors else None) or INK_SOFT
            if fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if c in mono_cols:
                p.alignment = PP_ALIGN.RIGHT
                set_run(p.add_run(), val, font=MONO_FONT, size=cell_size,
                        bold=True if fill in (HM_VBAD, HM_BAD, HM_HOT, HM_WARM, HM_LIGHT, HM_MID) else False,
                        color=txt_color)
            else:
                p.alignment = PP_ALIGN.LEFT
                # Mark first column as bold-ish for label
                bold = (c == 0)
                set_run(p.add_run(), val, font=KOR_FONT, size=cell_size,
                        bold=bold, color=txt_color)
    # Set row heights
    for r in range(nrows):
        try:
            table.rows[r].height = row_h
        except Exception:
            pass
    return table


def add_implication_bar(slide, text, *, y=None, runs=None):
    """The .a-impl element — full-width navy bar with white text, arrow prefix.
    HTML: position: absolute; left: 0; right: 0; bottom: 0;
          background: navy (#1B3DAD); color: #fff; padding: 16px 80px;
          font-size: 15px; line-height: 1.45; gap: 14px;
          ::before { content: "→"; color: blueSoft (#E4ECF8); font-mono; }
    runs: optional list of {text, bold?, color?} — overrides text param for richer formatting.
    """
    if y is None:
        y = Inches(6.79)   # leaves 0.71" at bottom — matches HTML 56px+impl height
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), y,
                                  SLIDE_W, Inches(0.71))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    # padding: 16px 80px
    tf.margin_left = MARGIN_X; tf.margin_right = MARGIN_RIGHT
    tf.margin_top = Inches(0.17); tf.margin_bottom = Inches(0.17)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Arrow prefix → in mono blueSoft, then a small gap, then body text
    arrow = p.add_run()
    set_run(arrow, "→  ", font=MONO_FONT, size=11, bold=True, color=BLUE_SOFT)

    if runs:
        for r in runs:
            run = p.add_run()
            set_run(run, r["text"],
                    font=r.get("font", KOR_FONT),
                    size=r.get("size", 12),
                    bold=r.get("bold", False),
                    color=r.get("color", WHITE))
    else:
        run = p.add_run()
        set_run(run, text, font=KOR_FONT, size=12, bold=False, color=WHITE)


def add_stat_box(slide, x, y, w, h, label, number, sub):
    """A .stat-w widget: small mono label / big navy number / small body caption.
    HTML: .stat-w .label { mono 10px / gray-500 / 0.18em / uppercase / mb 4 }
          .stat-w .v     { font-num 700 / -0.03em letter-spacing / navy }
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.025
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.13); tf.margin_bottom = Inches(0.13)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # label-mono: 0.18em letter-spacing
    set_run(p.add_run(), label, font=MONO_FONT, size=8, bold=False, color=GRAY_500, spc_em=0.18)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    # .stat .v: -0.03em letter-spacing (HTML .num-mid)
    set_run(p2.add_run(), number, font=NUM_FONT, size=22, bold=True, color=NAVY, spc_em=-0.03)
    p3 = tf.add_paragraph(); p3.space_before = Pt(3)
    set_run(p3.add_run(), sub, font=KOR_FONT, size=9, color=GRAY_600)


CHART_DIR = Path("/Users/hyunbin/Capstone/experiments/figures/native_pptx_charts")


def add_chart_image(slide, slide_num, x, y, w=None, h=None,
                    *, eyebrow=None, caption=None):
    """Embed a pre-rendered matplotlib chart for slide S{slide_num}.

    The PNG sits inside a thin gray-200 border that mirrors the .card chrome.
    """
    img_path = CHART_DIR / f"S{slide_num}.png"
    if not img_path.exists():
        # Render-time fallback — surface clearly instead of silently skipping
        raise FileNotFoundError(f"Missing chart image: {img_path}")

    pic = slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)
    # Card-style border around the picture
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, y, pic.width, pic.height)
    border.fill.background()
    border.line.color.rgb = LINE
    border.line.width = Pt(0.5)
    border.shadow.inherit = False

    if eyebrow:
        add_simple_text(slide, x, y - Inches(0.32), pic.width, Inches(0.28),
                        eyebrow, font=MONO_FONT, size=8, bold=True,
                        color=NAVY, spc_em=0.18)
    if caption:
        add_simple_text(slide, x, y + pic.height + Inches(0.05),
                        pic.width, Inches(0.28),
                        caption, font=MONO_FONT, size=7.5, color=GRAY_500)
    return pic


def add_tag(slide, x, y, w, h, text, kind="navy"):
    """Inline-style tag (small colored rectangle)."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if kind == "navy":
        box.fill.solid(); box.fill.fore_color.rgb = NAVY
        text_color = WHITE
        box.line.fill.background()
    elif kind == "grn":
        box.fill.solid(); box.fill.fore_color.rgb = TAG_GRN_BG
        box.line.color.rgb = TAG_GRN_BORDER; box.line.width = Pt(0.5)
        text_color = GREEN
    elif kind == "red":
        box.fill.solid(); box.fill.fore_color.rgb = TAG_RED_BG
        box.line.color.rgb = TAG_RED_BORDER; box.line.width = Pt(0.5)
        text_color = RED
    else:
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LINE; box.line.width = Pt(0.5)
        text_color = INK_SOFT
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, font=MONO_FONT, size=8, bold=True, color=text_color)


# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # =====================================================
    # SLIDE 1 — COVER (Slides.jsx S1: bigger spacing, eyebrow + pills, big 60px title)
    # =====================================================
    s = prs.slides.add_slide(blank)
    # No standard header on cover — but we still want navy stripe at top-left
    add_navy_stripe(s)

    # Top-left eyebrow (mono, no stripe element) — spc_em 0.18 for tracked uppercase
    add_simple_text(s, MARGIN_X, Inches(0.6),
                    Inches(7.0), Inches(0.3),
                    "캡스톤 2026 · 속도는벡터",
                    font=MONO_FONT, size=9, bold=False, color=GRAY_500, spc_em=0.18)

    # Top-right pills — same as add_header
    pill_y = Inches(0.55)
    pill_h = Inches(0.27)
    pill2_w = Inches(0.95)
    pill2_x = SLIDE_W - MARGIN_RIGHT - pill2_w
    p2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pill2_x, pill_y, pill2_w, pill_h)
    p2.fill.solid(); p2.fill.fore_color.rgb = PILL_RED_BG
    p2.line.color.rgb = PILL_RED_BORDER; p2.line.width = Pt(0.75)
    p2.shadow.inherit = False
    tf2 = p2.text_frame
    tf2.margin_left = Inches(0.05); tf2.margin_right = Inches(0.05)
    tf2.margin_top = Inches(0.01); tf2.margin_bottom = Inches(0.01)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2p = tf2.paragraphs[0]; p2p.alignment = PP_ALIGN.CENTER
    set_run(p2p.add_run(), "BDAI LAB", font=MONO_FONT, size=9, color=RED)

    pill1_w = Inches(1.15)
    pill1_x = pill2_x - pill1_w - Inches(0.08)
    p1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, pill1_x, pill_y, pill1_w, pill_h)
    p1.fill.solid(); p1.fill.fore_color.rgb = PILL_NAVY_BG
    p1.line.color.rgb = PILL_NAVY_BORDER; p1.line.width = Pt(0.75)
    p1.shadow.inherit = False
    tf1 = p1.text_frame
    tf1.margin_left = Inches(0.05); tf1.margin_right = Inches(0.05)
    tf1.margin_top = Inches(0.01); tf1.margin_bottom = Inches(0.01)
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1p = tf1.paragraphs[0]; p1p.alignment = PP_ALIGN.CENTER
    set_run(p1p.add_run(), "YONSEI · CSE", font=MONO_FONT, size=9, color=NAVY)

    # Big title — Slides.jsx S1 has fontSize: 60, weight: 800, lineHeight: 1.04
    # Title is broken into 3 lines:
    #   Skew-Aware
    #   Stratified Sampling for
    #   <Vector-Augmented> Analytical Query
    title_y = Inches(1.85)
    line_h = Inches(0.85)   # close to 60px * 1.04 ≈ 62px
    # title-l: -0.03em letter-spacing, weight 800
    add_simple_text(s, MARGIN_X, title_y,
                    Inches(11.5), line_h,
                    "Skew-Aware",
                    font=KOR_FONT, size=44, bold=True, color=INK, spc_em=-0.03)
    add_simple_text(s, MARGIN_X, title_y + line_h,
                    Inches(11.5), line_h,
                    "Stratified Sampling for",
                    font=KOR_FONT, size=44, bold=True, color=INK, spc_em=-0.03)
    # Last line with two color runs — same letter-spacing
    add_runs_text(s, MARGIN_X, title_y + 2 * line_h,
                  Inches(12), line_h,
                  [
                      {"text": "Vector-Augmented", "font": KOR_FONT, "size": 44, "bold": True, "color": NAVY, "spc_em": -0.03},
                      {"text": " Analytical Query", "font": KOR_FONT, "size": 44, "bold": True, "color": INK, "spc_em": -0.03},
                  ])

    # Subtitle — gray-600, smaller body text
    add_runs_text(s, MARGIN_X, title_y + 3 * line_h + Inches(0.15),
                  Inches(11.5), Inches(0.45),
                  [
                      {"text": "중간 측정 결과 종합 — ", "size": 13, "color": GRAY_600},
                      {"text": "15 cell partial measurement", "size": 13, "bold": True, "color": NAVY},
                      {"text": " · 5/8 비대면 회의 자료", "size": 13, "color": GRAY_600},
                  ])

    # 4-col footer — TEAM / ADVISOR / REFERENCE / DATE
    # HTML: paddingTop 22, borderTop 1px gray-200; gridTemplateColumns: repeat(4, 1fr); gap 24
    foot_y = Inches(5.85)
    # top border line
    ln = s.shapes.add_connector(1, MARGIN_X, foot_y - Inches(0.18),
                                  SLIDE_W - MARGIN_RIGHT, foot_y - Inches(0.18))
    ln.line.color.rgb = LINE; ln.line.width = Pt(0.6)

    foot_total_w = SLIDE_W - MARGIN_X - MARGIN_RIGHT
    foot_w = (foot_total_w - Inches(0.75)) / 4   # 24px gap × 3 ≈ 0.75
    foot_gap = Inches(0.25)

    # v1_kind: "kor" (Pretendard) / "mono" (JetBrains Mono) / "num" (Inter, for numerals)
    cols = [
        ("TEAM", "속도는벡터", "박세은 · 강재현 · 조현빈 · 이동욱", "kor"),
        ("ADVISOR", "박광현 교수", "BDAI Research Lab", "kor"),
        ("REFERENCE", "arXiv:2512.09695v2", "Exqutor (BDAI)", "mono"),
        ("DATE", "2026.05.08", "회의 자료 (중간 측정)", "num"),
    ]
    for i, (lbl, v1, v2, v1_kind) in enumerate(cols):
        cx = MARGIN_X + (foot_w + foot_gap) * i
        # HTML: label-mono { 11px mono / gray-500 / 0.18em / uppercase } / mb 6
        # eyebrow with letter-spacing for emphasis (cover line)
        add_simple_text(s, cx, foot_y, foot_w, Inches(0.26),
                        lbl, font=MONO_FONT, size=8, color=GRAY_500, spc_em=0.18)
        # Primary value — fontSize 13 (or 11 for mono/num)
        if v1_kind == "kor":
            font_for_v1 = KOR_FONT; v1_size = 13
        elif v1_kind == "num":
            font_for_v1 = NUM_FONT; v1_size = 13     # Inter for date numerals
        else:
            font_for_v1 = MONO_FONT; v1_size = 11
        add_simple_text(s, cx, foot_y + Inches(0.26), foot_w, Inches(0.3),
                        v1, font=font_for_v1, size=v1_size, bold=True, color=INK)
        # Secondary value — fontSize 11, color gray-600, mono
        add_simple_text(s, cx, foot_y + Inches(0.56), foot_w, Inches(0.26),
                        v2, font=MONO_FONT, size=8, color=GRAY_600)

    add_footer(s, page=1)

    # =====================================================
    # SLIDE 2 — TOC
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.95), "",
                     "오늘의 구성",
                     "TABLE OF CONTENTS · 10 SECTIONS · 14 SLIDES · 회의 자료")

    toc_items = [
        ("01", "Problem", "Exqutor 위치"),
        ("02", "Prior Work", "ECQO + Adaptive Sampling"),
        ("03", "Approach", "분포 인지 stratification"),
        ("04", "RQ1 · Diagnostic", "selectivity gradient"),
        ("05", "RQ2 · Aware", "KM20 oracle + Anti-Neyman"),
        ("06", "RQ3 · Agnostic", "25 method tier elim"),
        ("07", "측정 매트릭스", "10 cell heatmap"),
        ("08", "Sweet Spot", "SIFT skew → −32%"),
        ("09", "Multi-relation", "multi-vector + natural join"),
        ("10", "Future Work", "sf100 + 자문 메일"),
    ]
    # HTML: gridTemplateColumns repeat(5, 1fr); gap 14px; minHeight 138px
    # Title: 16/700/ink; Sub: 11/mono/gray-500.
    grid_x0 = MARGIN_X
    grid_y0 = Inches(2.5)
    grid_w_total = SLIDE_W - MARGIN_X - MARGIN_RIGHT
    cell_gap = Inches(0.15)   # 14px ≈ 0.146"
    cell_w = (grid_w_total - cell_gap * 4) / 5
    cell_h = Inches(1.7)      # >138px to fit two-line items comfortably
    row_gap = Inches(0.2)

    for idx, (num, title, sub) in enumerate(toc_items):
        col = idx % 5
        row = idx // 5
        x = grid_x0 + (cell_w + cell_gap) * col
        y = grid_y0 + (cell_h + row_gap) * row
        # Card — white BG, gray-200 border, 3px radius, 18px 20px padding
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, cell_h)
        card.adjustments[0] = 0.025
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False

        # b-num — 24×24px black ink square with white mono number
        # 24px → 0.25" but Slides.jsx uses 11px font; render b-num at 0.25" with size 9
        bn_x = x + Inches(0.21)
        bn_y = y + Inches(0.21)
        bn = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                bn_x, bn_y,
                                Inches(0.27), Inches(0.27))
        bn.fill.solid(); bn.fill.fore_color.rgb = INK
        bn.line.fill.background()
        bn.shadow.inherit = False
        btf = bn.text_frame
        btf.margin_left = Inches(0); btf.margin_right = Inches(0)
        btf.margin_top = Inches(0); btf.margin_bottom = Inches(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        set_run(bp.add_run(), num, font=NUM_FONT, size=9, bold=True, color=WHITE)

        # Title (16px/700/ink) — anchored bottom-left of card
        add_simple_text(s, x + Inches(0.21),
                        y + cell_h - Inches(0.7),
                        cell_w - Inches(0.4), Inches(0.32),
                        title, font=KOR_FONT, size=12, bold=True, color=INK)
        # Sub (11px mono/gray-500)
        add_simple_text(s, x + Inches(0.21),
                        y + cell_h - Inches(0.4),
                        cell_w - Inches(0.4), Inches(0.26),
                        sub, font=MONO_FONT, size=8, color=GRAY_500)

    add_footer(s, page=2)

    # =====================================================
    # SLIDE 3 — Problem (Exqutor 위치)
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "01",
                     "본 연구의 위치 — Exqutor 가 미해결한 영역",
                     "PROBLEM · ADAPTIVE SAMPLING 의 전 단계 분배 전략 정량화")

    body_y = Inches(1.95)
    col_w = Inches(5.7)
    gap = Inches(0.25)

    # LEFT card with table
    left = add_card(s, MARGIN_X, body_y, col_w, Inches(3.5), fill=WHITE, border=LINE)
    # eyebrow
    add_simple_text(s, MARGIN_X + Inches(0.18), body_y + Inches(0.14),
                    col_w - Inches(0.36), Inches(0.3),
                    "EXQUTOR 매트릭스",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    # table inside left card
    add_table_styled(s, MARGIN_X + Inches(0.18), body_y + Inches(0.55),
                     col_w - Inches(0.36), Inches(1.5),
                     headers=["", "인덱스 ✓", "인덱스 ✗"],
                     rows=[
                         ["multi-table", "ECQO · range query", "Adaptive Sampling 영역"],
                         ["single-table", "ECQO 부분", "Adaptive Sampling · momentum"],
                     ],
                     col_widths=[Inches(1.6), Inches(1.85), Inches(1.85)],
                     row_h=Inches(0.42))
    add_simple_text(s, MARGIN_X + Inches(0.18), body_y + Inches(2.6),
                    col_w - Inches(0.36), Inches(0.7),
                    "Exqutor 본 논문 (BDAI-Research) 의 핵심 contribution 영역 — multi-table + indexed 영역에 집중.",
                    font=KOR_FONT, size=10, color=MUTED)

    # RIGHT navy card
    right_x = MARGIN_X + col_w + gap
    right_card = add_card_navy(s, right_x, body_y, col_w, Inches(2.0))
    card_set_content(
        right_card,
        eyebrow="본 연구의 공략점",
        heading="단일 + 비인덱스 영역의\n분배 전략 (KM20 vs BERN) 가치 정량",
        head_size=14,
        dash_items=[
            [
                {"text": "Exqutor Adaptive Sampling 의 ", "size": 10, "color": INK_SOFT},
                {"text": "전 단계", "size": 10, "bold": True, "color": INK_SOFT},
                {"text": " sample 분배 전략", "size": 10, "color": INK_SOFT},
            ],
            [
                {"text": "proportional vs Neyman vs ", "size": 10, "color": INK_SOFT},
                {"text": "distribution-aware", "size": 10, "bold": True, "color": NAVY},
            ],
            [
                {"text": "분포 모를 때 ", "size": 10, "color": INK_SOFT},
                {"text": "production-ready 4강", "size": 10, "bold": True, "color": NAVY},
                {"text": " 도출", "size": 10, "color": INK_SOFT},
            ],
            "중간 측정 결과 = 단일 정확성의 필요조건 입증",
        ]
    )

    # Two stat boxes below right card
    stat_y = body_y + Inches(2.15)
    stat_w = (col_w - Inches(0.15)) / 2
    add_stat_box(s, right_x, stat_y, stat_w, Inches(1.35),
                 "DATASET", "5", "DEEP·SIFT·SSN++·WIKI·YFCC")
    add_stat_box(s, right_x + stat_w + Inches(0.15), stat_y, stat_w, Inches(1.35),
                 "CELL", "13", "10 단일 100% + 3 multi 부록")

    add_footer(s, page=3)

    # =====================================================
    # SLIDE 4 — Prior Work
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "02",
                     "선행 연구 — Exqutor: ECQO + Adaptive Sampling",
                     "PRIOR WORK · arXiv:2512.09695v2 · BDAI Research Lab")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    # LEFT card — ECQO
    left = add_card(s, MARGIN_X, body_y, col_w, Inches(2.4), fill=BG_CARD, border=LINE)
    card_set_content(left,
                     eyebrow="ECQO — INDEXED RANGE QUERY",
                     heading="HNSW range query 1~2 ms",
                     head_size=14,
                     dash_items=[
                         "인덱스 + multi-table 영역 정량 처리",
                         "pgvector 33.3% / VBASE 50% / DuckDB 100% 의 고정 비율 한계 극복",
                         "본 논문 contribution 의 정량 핵심",
                     ])

    # RIGHT card — Adaptive Sampling
    right_x = MARGIN_X + col_w + Inches(0.25)
    right = add_card(s, right_x, body_y, col_w, Inches(2.4), fill=BG_CARD, border=LINE)
    card_set_content(right,
                     eyebrow="ADAPTIVE SAMPLING — UNINDEXED",
                     heading="momentum 기반 동적 sample size 1.2~3.2× speedup",
                     head_size=14,
                     dash_items=[
                         "인덱스 없는 단일 테이블 영역 처리",
                         [{"text": "sample 분배 전략은 ", "size": 10, "color": INK_SOFT},
                          {"text": "정량 분석 대상이 아님", "size": 10, "bold": True, "color": RED}],
                         [{"text": "본 연구가 ", "size": 10, "color": INK_SOFT},
                          {"text": "이 단계의 가치", "size": 10, "bold": True, "color": INK_SOFT},
                          {"text": " 를 정량화", "size": 10, "color": INK_SOFT}],
                     ])

    # Bottom dashed insight box
    insight_y = body_y + Inches(2.6)
    box = add_card(s, MARGIN_X, insight_y, SLIDE_W - 2*MARGIN_X, Inches(1.85),
                   fill=WHITE, border=LINE, dashed=True)
    card_set_content(box,
                     eyebrow="본 연구의 보완 위치",
                     body_runs=[
                         {"text": "Exqutor 의 ", "size": 11, "color": INK_SOFT},
                         {"text": "Adaptive Sampling 모듈", "size": 11, "bold": True, "color": INK},
                         {"text": " 이 momentum 으로 sample size ", "size": 11, "color": INK_SOFT},
                         {"text": "크기", "size": 11, "bold": True, "color": INK},
                         {"text": " 를 조정하는 단계의 ", "size": 11, "color": INK_SOFT},
                         {"text": "전 단계", "size": 11, "bold": True, "color": INK},
                         {"text": ", 즉 sample ", "size": 11, "color": INK_SOFT},
                         {"text": "분배 전략 (proportional vs Neyman vs distribution-aware)", "size": 11, "bold": True, "color": INK},
                         {"text": " 의 가치를 정량 입증한다. 중간 측정 결과 = 5 dataset × 2 scale = 10 단일 cell + 3 multi cell 의 ", "size": 11, "color": INK_SOFT},
                         {"text": "15 cell partial", "size": 11, "bold": True, "color": NAVY},
                         {"text": " 측정으로, sf100 (80M) 은 5/8 회의 후 자문 합의 결과 반영하여 별도 측정.", "size": 11, "color": INK_SOFT},
                     ])

    add_footer(s, page=4)

    # =====================================================
    # SLIDE 5 — Approach
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "03",
                     "Approach — 분포 인지 stratification + 4강 method",
                     "APPROACH · KM20 ORACLE + DISTRIBUTION-AGNOSTIC 4 WINNERS")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    # LEFT card — RQ2
    left = add_card(s, MARGIN_X, body_y, col_w, Inches(3.0), fill=WHITE, border=LINE)
    card_set_content(left,
                     eyebrow="RQ2 · DISTRIBUTION-AWARE",
                     heading="KM20 oracle stratification",
                     head_size=14,
                     body_runs=[
                         {"text": "full K-means K=20 학습 → cluster ID 로 stratum 분배. ", "size": 10, "color": INK_SOFT},
                         {"text": "sample size 100/385/1000/3000 모두 BERN 우월", "size": 10, "bold": True, "color": NAVY},
                         {"text": ".", "size": 10, "color": INK_SOFT},
                     ],
                     dash_items=[
                         "Proportional (default) vs Neyman vs Anti-Neyman ablation",
                         "K-aware: K∈{10,20,50,100,200} sweep · K_optimal per dataset",
                         "σ_i 신호 약함 honest 입증 (Wilcoxon p>0.5, d<0.1)",
                     ])

    # RIGHT card — RQ3
    right_x = MARGIN_X + col_w + Inches(0.25)
    right = add_card(s, right_x, body_y, col_w, Inches(3.0), fill=WHITE, border=LINE)
    card_set_content(right,
                     eyebrow="RQ3 · DISTRIBUTION-AGNOSTIC",
                     heading="30 method → 4강 (Tier 1 = 17, Tier 2 = 2)",
                     head_size=14,
                     body_runs=[
                         {"text": "30 method × 10 cell × 5 sel = 1500 measurement, 단일 100% finalize. ", "size": 10, "color": INK_SOFT},
                         {"text": "Tier 1 = 17 (-8.04~-6.78), Wave 0 = 3, Pruned = 7", "size": 10, "bold": True, "color": NAVY},
                         {"text": ".", "size": 10, "color": INK_SOFT},
                     ],
                     dash_items=[
                         [{"text": "Hilbert", "size": 10, "bold": True, "color": INK},
                          {"text": " — learning-free + 결정론", "size": 10, "color": INK_SOFT}],
                         [{"text": "MiniBatch_partial", "size": 10, "bold": True, "color": INK},
                          {"text": " — OLTP partial_fit, ARI 1.000", "size": 10, "color": INK_SOFT}],
                         [{"text": "Hybrid", "size": 10, "bold": True, "color": INK},
                          {"text": " — KMeans + Hilbert", "size": 10, "color": INK_SOFT}],
                         [{"text": "HDBSCAN", "size": 10, "bold": True, "color": INK},
                          {"text": " — density-based skew 가치", "size": 10, "color": INK_SOFT}],
                     ])

    # Three stats below
    stat_y = body_y + Inches(3.2)
    stat_w = (SLIDE_W - 2*MARGIN_X - Inches(0.4)) / 3
    add_stat_box(s, MARGIN_X, stat_y, stat_w, Inches(1.4),
                 "METHOD", "33", "17 base + 11 NEW + 5 Wave1")
    add_stat_box(s, MARGIN_X + stat_w + Inches(0.2), stat_y, stat_w, Inches(1.4),
                 "TIER", "33→4", "33 → 28 → 12 → 4")
    add_stat_box(s, MARGIN_X + 2*(stat_w + Inches(0.2)), stat_y, stat_w, Inches(1.4),
                 "FINAL WINNER", "4", "Hilbert·Hybrid·MiniBatch_partial·HDBSCAN")

    add_footer(s, page=5)

    # =====================================================
    # SLIDE 6 — RQ1 Diagnostic
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "04",
                     "RQ1 · Diagnostic — Selectivity Gradient 단조성",
                     "PER-SEED SPEARMAN ρ + PAIRED BOOTSTRAP 95% CI")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    # LEFT — claim card + table
    claim = add_card_navy(s, MARGIN_X, body_y, col_w, Inches(1.55))
    card_set_content(claim,
                     eyebrow="CLAIM",
                     heading="BERN 부정확성은\nselectivity 가 작을수록 단조 증가",
                     head_size=14,
                     body_runs=[
                         {"text": "per-seed Spearman ρ + 95% CI (paired bootstrap, B=2000). 5 dataset × sf1/sf10 = 10 cell 부호 일관성 입증.",
                          "size": 10, "color": INK_SOFT},
                     ])

    # Table for ρ values
    rho_y = body_y + Inches(1.7)
    rho_rows = [
        ["DEEP", "−0.584", "−0.596", "✓"],
        ["SIFT", "−0.366", "−0.471", "✓"],
        ["SSN++", "−0.599", "−0.609", "✓"],
        ["WIKI", "−0.440", "−0.576", "✓"],
        ["YFCC", "−0.527", "−0.589", "✓"],
    ]
    add_table_styled(s, MARGIN_X, rho_y, col_w, Inches(2.1),
                     headers=["Dataset", "sf1 ρ", "sf10 ρ", "일관"],
                     rows=rho_rows,
                     col_widths=[Inches(1.7), Inches(1.4), Inches(1.4), Inches(1.2)],
                     mono_cols={1, 2}, row_h=Inches(0.32))
    add_simple_text(s, MARGIN_X, rho_y + Inches(2.15), col_w, Inches(0.3),
                    "10 cell 100% 부호 일관 (ρ < 0) · ρ ∈ [−0.609, −0.366] · 단일 100% 측정 완료",
                    font=MONO_FONT, size=8, color=SOFT)

    # Forest plot under the table (extends into the bottom band before .a-impl)
    chart6_y = rho_y + Inches(2.5)
    add_chart_image(s, 6, MARGIN_X, chart6_y,
                    w=col_w, h=Inches(1.55),
                    eyebrow="FOREST PLOT — Spearman ρ + 95% CI")

    # RIGHT — three small cards
    right_x = MARGIN_X + col_w + Inches(0.25)
    right_w = col_w
    rc1 = add_card(s, right_x, body_y, right_w, Inches(1.3), fill=WHITE, border=LINE)
    card_set_content(rc1,
                     eyebrow="METRIC DEFINITION",
                     body_runs=[
                         {"text": "Selectivity Gradient", "size": 11, "bold": True, "color": INK},
                         {"text": " g(s) = relative error 의 selectivity 도함수 부호. 단조 감소 (ρ < 0) 시 BERN 은 좁은 sel 에서 점점 더 부정확.",
                          "size": 11, "color": INK_SOFT},
                     ])

    rc2 = add_card(s, right_x, body_y + Inches(1.45), right_w, Inches(1.5), fill=WHITE, border=LINE)
    card_set_content(rc2,
                     eyebrow="중간 측정 결과 (5/8 AM)",
                     dash_items=[
                         '5/5 부호 일관 → "BERN 부정확성 단조성" 일반화 입증',
                         "4/5 부호 일관 → honest 별도 보고 + dataset 별 caveat",
                         "numpy estimator (Phase 7) 빠른 반복 환경",
                     ])

    rc3 = add_card(s, right_x, body_y + Inches(3.1), right_w, Inches(1.3), fill=WHITE, border=LINE)
    card_set_content(rc3,
                     eyebrow="2-LEVEL DECOMPOSITION",
                     body_runs=[
                         {"text": "상위: cell-level error (block vs row). 하위: stratum-level variance contribution. 단조성은 ",
                          "size": 11, "color": INK_SOFT},
                         {"text": "상위 수준", "size": 11, "bold": True, "color": INK},
                         {"text": " 의 BERN 한계 진단.", "size": 11, "color": INK_SOFT},
                     ])

    add_footer(s, page=6)

    # =====================================================
    # SLIDE 7 — RQ2 Aware
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "05",
                     "RQ2 · Aware — KM20 oracle + Anti-Neyman + K-aware",
                     "SAMPLE-SIZE ROBUSTNESS · σ_i SIGNAL HONESTY")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    # LEFT — sample-size robustness
    rs_card = add_card(s, MARGIN_X, body_y, col_w, Inches(2.2), fill=WHITE, border=LINE)
    add_simple_text(s, MARGIN_X + Inches(0.18), body_y + Inches(0.12),
                    col_w - Inches(0.36), Inches(0.3),
                    "KM20 ORACLE — SAMPLE-SIZE ROBUSTNESS",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_table_styled(s,
                     MARGIN_X + Inches(0.2), body_y + Inches(0.55),
                     col_w - Inches(0.4), Inches(1.4),
                     headers=["Sample", "DEEP sf10", "SIFT sf10"],
                     rows=[
                         ["n=100", "−1.42%", "−9.22%"],
                         ["n=385", "−1.76%", "−10.47%"],
                         ["n=1000", "−1.81%", "−10.20%"],
                         ["n=3000", "−1.79%", "−10.18%"],
                     ],
                     col_widths=[Inches(1.6), Inches(1.7), Inches(1.7)],
                     mono_cols={1, 2}, row_h=Inches(0.27))

    # Anti-Neyman ablation
    an_y = body_y + Inches(2.35)
    an_card = add_card(s, MARGIN_X, an_y, col_w, Inches(2.2), fill=WHITE, border=LINE)
    add_simple_text(s, MARGIN_X + Inches(0.18), an_y + Inches(0.12),
                    col_w - Inches(0.36), Inches(0.3),
                    "ANTI-NEYMAN ABLATION (s=0.01)",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_table_styled(s,
                     MARGIN_X + Inches(0.2), an_y + Inches(0.55),
                     col_w - Inches(0.4), Inches(1.45),
                     headers=["Dataset", "Anti-Neyman effect"],
                     rows=[
                         ["DEEP", "+10.85%"],
                         ["SIFT", "+8.95%"],
                         ["SSN++", "+16.45%"],
                         ["WIKI", "+19.47%"],
                         ["YFCC", "+9.69%"],
                     ],
                     col_widths=[Inches(2.5), Inches(2.5)],
                     mono_cols={1}, row_h=Inches(0.27))

    # RIGHT — K-optimal
    right_x = MARGIN_X + col_w + Inches(0.25)
    right_w = col_w
    k_card = add_card_navy(s, right_x, body_y, right_w, Inches(2.55))
    add_simple_text(s, right_x + Inches(0.22), body_y + Inches(0.12),
                    right_w - Inches(0.4), Inches(0.3),
                    "K-AWARE — K_OPTIMAL PER DATASET",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_simple_text(s, right_x + Inches(0.22), body_y + Inches(0.4),
                    right_w - Inches(0.4), Inches(0.55),
                    "K∈{10, 20, 50, 100, 200} sweep · 저차원 → K=20, 고차원 (256~768d) → K=50/100. K=20 default robustness 정량.",
                    font=KOR_FONT, size=10, color=INK_SOFT)
    add_table_styled(s,
                     right_x + Inches(0.22), body_y + Inches(1.0),
                     right_w - Inches(0.42), Inches(1.5),
                     headers=["Dataset", "dim", "K_opt"],
                     rows=[
                         ["DEEP", "96", "K=20 default"],
                         ["SIFT", "128", "K=20 default"],
                         ["YFCC", "192", "K=20 default"],
                         ["SSN++", "256", "K=50 (보강)"],
                         ["WIKI", "768", "K=50 (보강)"],
                     ],
                     col_widths=[Inches(1.7), Inches(1.0), Inches(2.6)],
                     mono_cols={1, 2}, row_h=Inches(0.27))

    sigma_card = add_card(s, right_x, body_y + Inches(2.7), right_w, Inches(1.85),
                           fill=BG_CARD, border=LINE)
    card_set_content(sigma_card,
                     eyebrow="σ_i SIGNAL HONESTY",
                     body_runs=[
                         {"text": "paired Wilcoxon ", "size": 11, "color": INK_SOFT},
                         {"text": "p > 0.5", "size": 11, "bold": True, "color": INK},
                         {"text": ", Cohen's d ", "size": 11, "color": INK_SOFT},
                         {"text": "< 0.1", "size": 11, "bold": True, "color": INK},
                         {"text": ". σ_i 신호의 ", "size": 11, "color": INK_SOFT},
                         {"text": "비결정적 가치", "size": 11, "bold": True, "color": NAVY},
                         {"text": " 명시. Anti-Neyman 은 negative control narrative 으로 사용.", "size": 11, "color": INK_SOFT},
                     ])

    add_footer(s, page=7)

    # =====================================================
    # SLIDE 8 — RQ3 Tier elimination (5/8 11:40 가지치기 + 4강 결정 confirm)
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "06",
                     "RQ3 · Agnostic — 33 → 16 살아남기 → 4강 → production-friendly 1",
                     "Tier 1 spread 1.24%p · METHOD CHOICE 보다 분포 인지가 결정적")

    # Funnel diagram visualises the elimination flow that the table below details
    chart8_y = Inches(1.85)
    chart8_w = SLIDE_W - 2 * MARGIN_X
    add_chart_image(s, 8, MARGIN_X, chart8_y,
                    w=chart8_w, h=Inches(1.45),
                    eyebrow="TIER ELIMINATION FUNNEL · 33 → 16 살아남기 → 4강 → production-friendly 1")

    tbl_y = Inches(3.45)
    add_table_styled(s, MARGIN_X, tbl_y, SLIDE_W - 2*MARGIN_X, Inches(2.0),
                     headers=["Tier", "기준", "잔존 method", "N"],
                     rows=[
                         ["Wave 0", "variance explosion outlier (paired Δ% >100%)",
                          "dbscan · lsh · random_proj (즉시 제외)", "3"],
                         ["Tier 1", "avg_Δ% ≥ -5.6 · neg ≥ 7/9 · CI excludes 0 ≥ 7/9",
                          "hdbscan·pca_kmeans·coresets·zorder·kmeans_pp·minibatch·minibatch_partial·gmm·hilbert·pca1d·agglomerative·hybrid·hierarchical_kmeans·sparse_rp·kdtree", "15"],
                         ["Tier 2", "boundary (Tier 1 미달, neg 6/9)",
                          "birch (avg -5.51)", "1"],
                         ["Tier 3", "Wave 1 partial 5 cell SURVIVE",
                          "reservoir (WIKI sf1 -8.30% partial)", "1"],
                         ["Pruned", "magnitude<3% OR sign 불일관 OR CI 약",
                          "pq · halton · hammersley · spectral · sobol · optics", "6"],
                         ["★ 4강", "production criteria (1위×cost×interpretability)",
                          "★1 hdbscan · ★2 hilbert · ★3 minibatch_partial · ★4 hybrid", "4"],
                     ],
                     col_widths=[Inches(1.0), Inches(3.6), Inches(5.4), Inches(0.7)],
                     mono_cols={3}, cell_size=9, row_h=Inches(0.4))

    # Four feature cards (sit beneath the smaller tier table thanks to the new funnel)
    card_y = Inches(5.6)
    card_w = (SLIDE_W - 2*MARGIN_X - Inches(0.45)) / 4
    card_h = Inches(1.15)
    cards = [
        ("★1 HDBSCAN", "Strongest narrative", "avg -8.04 · SIFT 1위 -34.17% · oracle 영역 (4313s, prod X)"),
        ("★2 MINIBATCH_PARTIAL", "OLTP narrative 유일", "avg -7.63 · CI 9/10 강력 · online partial_fit · stream OK"),
        ("★3 HILBERT", "Production sweet spot", "avg -7.54 · 수 초 fit · SIFT -33.53% · CI 9/10"),
        ("★4 HYBRID (MB+Hilbert)", "Mechanism ablation", "avg -7.13 · clustering vs ordering driver 분리"),
    ]
    for i, (eb, head, body) in enumerate(cards):
        cx = MARGIN_X + (card_w + Inches(0.15)) * i
        c = add_card(s, cx, card_y, card_w, card_h, fill=BG_CARD, border=LINE)
        card_set_content(c, eyebrow=eb, heading=head, body_runs=[
            {"text": body, "size": 9, "color": INK_SOFT},
        ], head_size=12)

    add_footer(s, page=8)

    # =====================================================
    # SLIDE 9 — 측정 매트릭스 Heatmap (Slide 10 in spec)
    #          5/8 11:40 가지치기 결과 confirm 후 4강 결정 narrative 강화
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "07",
                     "W4 Matrix — ★1 hdbscan · ★2 mb_partial · ★3 hilbert · ★4 hybrid × 10 cell",
                     "PAIRED Δ% vs BERN @ sel=0.10 · 단일 100% finalize 5/8 14:13 · 4강 production criteria")

    # Heatmap table — MAIN: 10 cell (5 dataset × sf1/sf10). YFCC_DL 폐기 결정 (5/8) 반영.
    body_y = Inches(1.95)
    # Define cells: (cell, hilbert, hybrid, mb_p, hdbscan)
    heatmap_data = [
        ("DEEP_sf1", "−1.07%", "−1.71%", "−1.99%", "−2.48%"),
        ("DEEP_sf10", "−1.98%", "−2.73%", "−2.87%", "−2.51%"),
        ("SIFT_sf1", "−33.53%", "−30.46%", "−33.13%", "−34.17%"),
        ("SIFT_sf10", "−12.02%", "−11.48%", "−11.63%", "−11.79%"),
        ("SSN_sf1", "+1.69%", "+0.64%", "+1.02%", "+0.84%"),
        ("SSN_sf10", "+1.38%", "+0.56%", "+1.35%", "+0.67%"),
        ("WIKI_sf1", "−10.92%", "−8.99%", "−11.30%", "−11.29%"),
        ("WIKI_sf10", "−5.70%", "−5.43%", "−3.77%", "−5.54%"),
        ("YFCC_sf1", "−8.07%", "−6.98%", "−8.37%", "−8.40%"),
        ("YFCC_sf10", "−5.21%", "−4.78%", "−5.62%", "−5.77%"),
    ]

    def color_for_value(v):
        # Parse percent
        s_ = v.replace("−", "-").replace("+", "").rstrip("%")
        try:
            n = float(s_)
        except ValueError:
            return None, INK_SOFT
        # negative = improve (green); positive = hurt (red)
        if n <= -25:
            return HM_VBAD, HM_VBAD_TXT
        if n <= -8:
            return HM_BAD, HM_BAD_TXT
        if n <= -3:
            return HM_MID, HM_MID_TXT
        if n < 0:
            return HM_LIGHT, HM_LIGHT_TXT
        if n == 0:
            return HM_NEUT, INK_SOFT
        if n < 2:
            return HM_WARM, HM_WARM_TXT
        if n < 5:
            return HM_WARM, HM_WARM_TXT
        return HM_HOT, WHITE

    rows = []
    fills = []
    text_colors = []
    for cell, h, hy, mb, hb in heatmap_data:
        rows.append([cell, h, hy, mb, hb])
        f_row = [None]
        c_row = [INK]
        for v in (h, hy, mb, hb):
            f, t = color_for_value(v)
            f_row.append(f)
            c_row.append(t)
        fills.append(f_row)
        text_colors.append(c_row)

    add_table_styled(s, MARGIN_X, body_y, Inches(7.5), Inches(3.65),
                     headers=["Cell", "Hilbert", "Hybrid", "MiniBatch_p", "HDBSCAN"],
                     rows=rows,
                     col_widths=[Inches(1.7), Inches(1.45), Inches(1.45), Inches(1.45), Inches(1.45)],
                     mono_cols={1, 2, 3, 4},
                     cell_fills=fills, cell_text_colors=text_colors,
                     cell_size=9, row_h=Inches(0.30))

    add_simple_text(s, MARGIN_X, body_y + Inches(3.78),
                    Inches(7.5), Inches(0.3),
                    "MAIN 10 cell (5 dataset × sf1/sf10) · 단일 100% 측정 완료 5/8 14:13",
                    font=MONO_FONT, size=8, color=SOFT)

    # Legend / cross-scale right column
    right_x = Inches(8.55)
    right_w = Inches(4.0)

    leg_card = add_card(s, right_x, body_y, right_w, Inches(1.4), fill=WHITE, border=LINE)
    add_simple_text(s, right_x + Inches(0.18), body_y + Inches(0.12),
                    right_w - Inches(0.36), Inches(0.3),
                    "HEATMAP LEGEND",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    # Manually add legend swatches inside the card area (under the eyebrow)
    leg_labels = [
        ("≤ −25%", HM_VBAD, WHITE),
        ("−8 ~ −25%", HM_BAD, WHITE),
        ("−3 ~ −8%", HM_MID, HM_MID_TXT),
        ("0 ~ −3%", HM_LIGHT, HM_LIGHT_TXT),
        ("0 ~ +5%", HM_WARM, HM_WARM_TXT),
        ("> +5%", HM_HOT, WHITE),
    ]
    sw_y = body_y + Inches(0.55)
    for i, (lbl, fill, text) in enumerate(leg_labels):
        col = i % 3
        row = i // 3
        sx = right_x + Inches(0.18) + Inches(1.25) * col
        sy = sw_y + Inches(0.4) * row
        sw = slide_swatch(s, sx, sy, Inches(1.2), Inches(0.32), lbl, fill, text)

    cs_card = add_card_navy(s, right_x, body_y + Inches(1.55), right_w, Inches(2.2))
    card_set_content(cs_card,
                     eyebrow="4강 결정 + CROSS-SCALE 일관성",
                     dash_items=[
                         "★1 hdbscan -8.04 · SIFT 1위 -34.17% (oracle)",
                         "★2 mb_partial -7.63 · CI 9/10 · OLTP 유일",
                         "★3 hilbert -7.54 · 수 초 fit · CI 9/10",
                         "★4 hybrid -7.13 · MB+Hilbert ablation 분리",
                         "Tier 1 spread 1.21%p · 분포 인지 boundary 결정적",
                         "5/5 dataset sf1 ↔ sf10 부호 동일 · 단일 100% finalize",
                     ])

    add_footer(s, page=9)

    # =====================================================
    # SLIDE 10 — Sweet Spot
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "08",
                     "Distribution Sweet Spot — Skew → 큰 효과, Balanced → 미세 hurt",
                     "CLUSTER RATIO ↑ → EFFECT MAGNITUDE ↑ · HONEST REPORT")

    body_y = Inches(1.95)

    # LEFT — distribution table
    rows = [
        ["SIFT_sf1", "9.41", "−32.08%", "SKEW"],
        ["SIFT_sf10", "~9", "−10.72%", "SKEW"],
        ["WIKI_sf1", "~3.5", "−9.61%", "SKEW"],
        ["WIKI_sf10", "~3.5", "−4.48%", "SKEW"],
        ["YFCC_sf1", "~2.4", "−6.88%", "MID"],
        ["DEEP_sf1", "~3", "−0.43%", "MILD"],
        ["SSN_sf1", "1.29", "+2.34%", "BALANCED"],
        ["SSN_sf10", "~1.3", "+2.06%", "BALANCED"],
    ]
    cell_fills = []
    text_colors = []
    for r in rows:
        eff = r[2]
        is_pos = eff.startswith("+")
        f_row = [None, None,
                 HM_LIGHT if (not is_pos) else HM_WARM if eff.endswith("%") else None,
                 None]
        c_row = [INK, INK_SOFT,
                 HM_LIGHT_TXT if (not is_pos) else HM_WARM_TXT,
                 INK]
        cell_fills.append(f_row)
        text_colors.append(c_row)

    add_table_styled(s, MARGIN_X, body_y, Inches(7.0), Inches(3.4),
                     headers=["Cell", "Cluster ratio", "Effect (Hilbert)", "분류"],
                     rows=rows,
                     col_widths=[Inches(1.85), Inches(1.7), Inches(2.1), Inches(1.35)],
                     mono_cols={1, 2}, row_h=Inches(0.4),
                     cell_fills=cell_fills, cell_text_colors=text_colors)

    # RIGHT — narrative card (top half) + sweet-spot scatter (bottom half)
    right_x = Inches(8.05)
    right_w = Inches(4.5)
    nar = add_card_navy(s, right_x, body_y, right_w, Inches(2.05))
    card_set_content(nar,
                     eyebrow="SWEET SPOT NARRATIVE",
                     dash_items=[
                         [{"text": "SIFT skew (ratio 9.41)", "size": 10, "bold": True, "color": INK},
                          {"text": " → −28~−32% LARGE", "size": 10, "color": INK_SOFT}],
                         [{"text": "WIKI/YFCC mid-skew", "size": 10, "bold": True, "color": INK},
                          {"text": " → −5~−10% consistent", "size": 10, "color": INK_SOFT}],
                         [{"text": "SSN++ balanced (ratio 1.29)", "size": 10, "bold": True, "color": INK},
                          {"text": " → +1~+2% mild hurt — ", "size": 10, "color": INK_SOFT},
                          {"text": "honest 별도 보고", "size": 10, "bold": True, "color": NAVY}],
                         "Sweet spot boundary: ratio > 1.3",
                     ])

    chart10_y = body_y + Inches(2.2)
    add_chart_image(s, 10, right_x, chart10_y,
                    w=right_w, h=Inches(1.55),
                    eyebrow="SCATTER · CLUSTER RATIO vs Δ%")

    # Bottom implication bar
    add_implication_bar(s, "SSN++ ceiling = balanced + low intrinsic dim + KM20-to-BERN headroom 0.5% (vs SIFT 34.5%)")

    add_footer(s, page=10)

    # =====================================================
    # SLIDE 11 — Multi-vector + Multi-table
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "09",
                     "Multi-vector + Multi-table natural join — Exqutor 직접 매칭",
                     "3 MULTI CELL · DIRECT EXQUTOR MULTI-TABLE 영역 매칭")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    # LEFT — Multi-vector × 4강 method (sel=0.10 paired Δ% vs BERN)
    mv_card = add_card(s, MARGIN_X, body_y, col_w, Inches(2.2), fill=WHITE, border=LINE)
    add_simple_text(s, MARGIN_X + Inches(0.18), body_y + Inches(0.12),
                    col_w - Inches(0.36), Inches(0.3),
                    "MULTI-VECTOR × 4강 (sel=0.10 paired Δ%)",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_table_styled(s,
                     MARGIN_X + Inches(0.2), body_y + Inches(0.55),
                     col_w - Inches(0.4), Inches(1.5),
                     headers=["Cell (96+dim)", "hdbscan", "hilbert", "hybrid", "mb_partial"],
                     rows=[
                         ["partsupp_deep_sift_10", "−1.02%", "−0.48%", "+0.31%", "−1.30%"],
                         ["partsupp_deep_wiki_10", "+1.15%", "+0.06%", "+0.08%", "+0.99%"],
                     ],
                     col_widths=[Inches(2.0), Inches(0.85), Inches(0.85), Inches(0.85), Inches(0.95)],
                     mono_cols={1, 2, 3, 4}, row_h=Inches(0.32), cell_size=9)
    add_simple_text(s, MARGIN_X + Inches(0.2), body_y + Inches(1.85),
                    col_w - Inches(0.4), Inches(0.3),
                    "STAGE 1+2 완료 (5/8) · n=2500/cell · sign 3/8 neg · |Δ| < 1.5% marginal",
                    font=KOR_FONT, size=10, color=MUTED)

    # LEFT bottom — Multi-table (compact) + grouped bar chart underneath
    mt_y = body_y + Inches(2.4)
    mt_card = add_card(s, MARGIN_X, mt_y, col_w, Inches(1.4), fill=WHITE, border=LINE)
    add_simple_text(s, MARGIN_X + Inches(0.18), mt_y + Inches(0.12),
                    col_w - Inches(0.36), Inches(0.3),
                    "MULTI-TABLE NATURAL JOIN × 4강 — STAGE 3 측정 진행 중",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_table_styled(s,
                     MARGIN_X + Inches(0.2), mt_y + Inches(0.55),
                     col_w - Inches(0.4), Inches(0.6),
                     headers=["Cell", "type", "STAGE 3 4강 Δ%"],
                     rows=[
                         ["partsupp_deep_10 ⨝ part_wiki_10", "TPC-H", "측정 중 (회의 후)"],
                     ],
                     col_widths=[Inches(2.7), Inches(1.0), Inches(1.4)],
                     mono_cols={1, 2}, row_h=Inches(0.3), cell_size=9)
    add_simple_text(s, MARGIN_X + Inches(0.2), mt_y + Inches(1.0),
                    col_w - Inches(0.4), Inches(0.4),
                    "ps_partkey = p_partkey 자연 외래키 · DEEP 96 ↔ WIKI 768 dim 격차",
                    font=KOR_FONT, size=10, color=MUTED)

    # Bottom strip — grouped bar chart for the 3 multi cells × 4 winners
    chart11_y = mt_y + Inches(1.55)
    add_chart_image(s, 11, MARGIN_X, chart11_y,
                    w=SLIDE_W - 2 * MARGIN_X, h=Inches(0.95),
                    eyebrow="MULTI-CELL · 4강 method paired Δ%")

    # RIGHT navy card — Magnitude shrinkage (single → multi)
    right_x = MARGIN_X + col_w + Inches(0.25)
    right_w = col_w
    ex_card = add_card_navy(s, right_x, body_y, right_w, Inches(2.5))
    card_set_content(ex_card,
                     eyebrow="MAGNITUDE SHRINKAGE — SINGLE → MULTI",
                     body_runs=[
                         {"text": "단일 sweet spot 4강 평균 |Δ%| ", "size": 11, "color": INK_SOFT},
                         {"text": "17.13%", "size": 11, "bold": True, "color": INK},
                         {"text": " (SIFT/WIKI/YFCC sf1) → multi-vector 4강 평균 |Δ%| ", "size": 11, "color": INK_SOFT},
                         {"text": "0.67%", "size": 11, "bold": True, "color": INK},
                         {"text": " (sel=0.10). magnitude ", "size": 11, "color": INK_SOFT},
                         {"text": "25.4× 약화", "size": 11, "bold": True, "color": NAVY},
                         {"text": ". 단일 정확성 = multi 정확성의 ", "size": 11, "color": INK_SOFT},
                         {"text": "필요조건만", "size": 11, "bold": True, "color": NAVY},
                         {"text": " (충분조건 X).", "size": 11, "color": INK_SOFT},
                     ],
                     dash_items=[
                         "sel=0.01: 8/8 hurt (sample budget narrow)",
                         "sel=0.10: 3/8 neg · 5/8 pos (boundary, |Δ| < 1.5%)",
                         "sel=0.50: 7/8 neg (수렴, magnitude marginal)",
                         "method ranking 보존 X — multi 환경 노이즈 수준",
                     ])

    # Multi-relation takeaway card — future work narrative
    mr_card = add_card(s, right_x, body_y + Inches(2.7), right_w, Inches(1.4),
                        fill=BG_CARD, border=LINE)
    card_set_content(mr_card,
                     eyebrow="MULTI-RELATION FUTURE WORK",
                     body_runs=[
                         {"text": "단일 sweet spot의 강력 improve가 multi-vector에서 marginal로 약화 → multi-relation 일반화는 ", "size": 10, "color": INK_SOFT},
                         {"text": "joint-aware clustering", "size": 10, "bold": True, "color": INK},
                         {"text": " 또는 ", "size": 10, "color": INK_SOFT},
                         {"text": "multi-vector decomposition", "size": 10, "bold": True, "color": INK},
                         {"text": " 별도 설계 필요. STAGE 3 multi-table join × 4강 측정은 회의 후 보강.",
                          "size": 10, "color": INK_SOFT},
                     ])

    add_footer(s, page=11)

    # =====================================================
    # SLIDE 12 — Cross-Scale (1M ↔ 8M)  [renumbered after YFCC Verify removal]
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "10",
                     "Cross-Scale — sf1 (800K) ↔ sf10 (8M) 일관성",
                     "SIGN CONSISTENCY · MAGNITUDE ATTENUATION · NECESSARY COND")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    cs_rows = [
        ["DEEP", "−1.07%", "−1.98%", "✓"],
        ["SIFT", "−33.53%", "−12.02%", "✓"],
        ["SSN++", "+1.69%", "+1.38%", "✓ (hurt)"],
        ["WIKI", "−10.92%", "−5.70%", "✓"],
        ["YFCC", "−8.07%", "−5.21%", "✓"],
    ]
    add_table_styled(s, MARGIN_X, body_y, col_w, Inches(2.5),
                     headers=["Dataset", "sf1 (Hilbert)", "sf10 (Hilbert)", "부호 일관"],
                     rows=cs_rows,
                     col_widths=[Inches(1.4), Inches(1.6), Inches(1.6), Inches(1.1)],
                     mono_cols={1, 2}, row_h=Inches(0.36))
    add_simple_text(s, MARGIN_X, body_y + Inches(2.65), col_w, Inches(0.3),
                    "단일 10 cell 100% 측정 완료 (5/5 dataset 부호 일관) · 5/8 14:13",
                    font=MONO_FONT, size=8, color=SOFT)

    right_x = MARGIN_X + col_w + Inches(0.25)
    right_w = col_w

    cs_card = add_card_navy(s, right_x, body_y, right_w, Inches(2.4))
    card_set_content(cs_card,
                     eyebrow="CROSS-SCALE TAKEAWAY",
                     dash_items=[
                         [{"text": "SIFT: sf1 −33.5% → sf10 −12.0% · effect ", "size": 11, "color": INK_SOFT},
                          {"text": "1/3 attenuation", "size": 11, "bold": True, "color": INK}],
                         "YFCC: sf1 −8.07% → sf10 −5.21% · 부호 일관 (단일 100%)",
                         "SSN++ hurt: sf1↔sf10 ceiling effect honest reporting",
                         "sf100 (80M) 추가 시 cross-scale 3-point validation 완성",
                     ])

    he_card = add_card(s, right_x, body_y + Inches(2.55), right_w, Inches(1.7),
                       fill=BG_CARD, border=LINE)
    card_set_content(he_card,
                     eyebrow="EFFECT HONESTY",
                     body_runs=[
                         {"text": "본 연구는 단일 정확성을 multi 정확성의 ", "size": 11, "color": INK_SOFT},
                         {"text": "필요조건만", "size": 11, "bold": True, "color": INK},
                         {"text": " 입증. 충분조건 일반화 + sf100 cross-scale 3-point 검증은 future work. 5/8 회의 후 자문 합의 결과 반영하여 sf100 launch.",
                          "size": 11, "color": INK_SOFT},
                     ])

    add_footer(s, page=12)

    # =====================================================
    # SLIDE 13 — Honest Limitations  [renumbered]
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "11",
                     "Honest Limitations — 중간 측정 결과 기준",
                     "8 LIMITATION · MAPPED TO FUTURE WORK")

    body_y = Inches(1.95)
    lim_rows = [
        ["1", "단일 → multi-table generalization", "partsupp_deep_sift / deep_10⨝part_wiki_10 부분 입증, multi-relation 일반화는 future"],
        ["2", "NPY-only mode 의 RQ2 dependency", "적재본 부재 시 RQ2 skip · NPY 파이프라인 정비 future"],
        ["3", "YFCC PCA basis 의존성 caveat", "YFCC distribution shape 는 PCA fit sample 의존 — sf100 launch 시 sample sweep 검증"],
        ["4", "σ_i 신호 약함", "Anti-Neyman vs Proportional Wilcoxon p>0.5, d<0.1 honest"],
        ["5", "IS NaN sel=0.01 발산 (80~95%)", "분할 X + weight only invalid → negative control narrative"],
        ["6", "K-sweep upper bound K=200", "WIKI/SSN++ 고차원 K>200 영역 차후"],
        ["7", "sf100 (80M) deferred", "5/8 자문 합의 후 5/22 까지 launch"],
        ["8", "SSN++ balanced ceiling 가설", "cluster ratio 1.29 미만 구조에서는 분배 효과 ceiling — 추가 검증"],
    ]
    # Status bar chart sits to the right of the limitations table — same 8 rows,
    # visualised by resolution status (PARTIAL / FUTURE / HONEST)
    table_w = Inches(7.2)
    add_table_styled(s, MARGIN_X, body_y, table_w, Inches(4.5),
                     headers=["#", "Limitation", "대응 / future work"],
                     rows=lim_rows,
                     col_widths=[Inches(0.5), Inches(2.5), Inches(4.2)],
                     mono_cols={0}, row_h=Inches(0.42), cell_size=9)

    chart14_x = MARGIN_X + table_w + Inches(0.2)
    chart14_w = SLIDE_W - chart14_x - MARGIN_RIGHT
    add_chart_image(s, 14, chart14_x, body_y,
                    w=chart14_w, h=Inches(4.5),
                    eyebrow="LIMITATION STATUS BAR")

    add_footer(s, page=13)

    # =====================================================
    # SLIDE 14 — Future Work + 자문 메일  [renumbered]
    # =====================================================
    s = prs.slides.add_slide(blank)
    add_header(s)
    add_lead_heading(s, Inches(0.85), "12",
                     "Future Work — sf100 plan + 자문 메일 + 5/27 narrative",
                     "L1: SF100 · L2: 자문 · L3: 5/27 발표 · L4: 차후 학기")

    body_y = Inches(1.95)
    col_w = Inches(5.7)

    # LEFT card — sf100 plan
    sf_card = add_card(s, MARGIN_X, body_y, col_w, Inches(2.4), fill=WHITE, border=LINE)
    add_simple_text(s, MARGIN_X + Inches(0.18), body_y + Inches(0.12),
                    col_w - Inches(0.36), Inches(0.3),
                    "L1 · SF100 (80M) PLAN",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_table_styled(s,
                     MARGIN_X + Inches(0.2), body_y + Inches(0.55),
                     col_w - Inches(0.4), Inches(1.7),
                     headers=["Step", "ETA"],
                     rows=[
                         ["5/8 회의 자문 메일 초안 합의", "5/8 19:00"],
                         ["채림 + 지도교수 메일 발송", "~5/12"],
                         ["자문 회신 ETA", "~5/15"],
                         ["sf100 launch (실험 파이프라인 확장)", "5/15~5/21"],
                         ["측정 5 dataset × ~2-4h", "~10-20h overnight"],
                     ],
                     col_widths=[Inches(3.9), Inches(1.2)],
                     mono_cols={1}, row_h=Inches(0.3), cell_size=9)

    # LEFT bottom — L4 차후 학기
    fut_y = body_y + Inches(2.6)
    fut_card = add_card(s, MARGIN_X, fut_y, col_w, Inches(2.0), fill=WHITE, border=LINE)
    card_set_content(fut_card,
                     eyebrow="L4 · 차후 학기",
                     dash_items=[
                         "Exqutor multi-table 영역 일반화",
                         "vector.c C-level integration (Phase 6 SQL D)",
                         "K>200 K-sweep · WIKI 768d 고차원",
                         "Distribution shift online detection",
                         "DuckDB native fixed-rate baseline 직접 통합",
                     ])

    # RIGHT navy — 자문 메일
    right_x = MARGIN_X + col_w + Inches(0.25)
    adv_card = add_card_navy(s, right_x, body_y, col_w, Inches(2.4))
    add_simple_text(s, right_x + Inches(0.22), body_y + Inches(0.12),
                    col_w - Inches(0.4), Inches(0.3),
                    "L2 · 자문 메일 핵심 질문",
                    font=MONO_FONT, size=8, bold=True, color=NAVY)
    add_table_styled(s,
                     right_x + Inches(0.22), body_y + Inches(0.55),
                     col_w - Inches(0.42), Inches(1.7),
                     headers=["대상", "핵심 질문"],
                     rows=[
                         ["채림 석사", "(1) 5 dataset 분포 진단 metric (skew score · cluster ratio) 신뢰도\n(2) multi-table natural join 영역의 sampling rate dependency"],
                         ["지도교수", "(1) sf100 (80M) 5 dataset 측정의 시간 자원 합의\n(2) 중간 측정 narrative (4강 + 25 method tier) 발표 정합성"],
                     ],
                     col_widths=[Inches(1.3), Inches(3.8)],
                     row_h=Inches(0.85), cell_size=9)

    # RIGHT bottom — 5/27 narrative
    nar_card = add_card(s, right_x, fut_y, col_w, Inches(2.0), fill=BG_CARD, border=LINE)
    card_set_content(nar_card,
                     eyebrow="L3 · 5/27 발표 NARRATIVE",
                     body_runs=[
                         {"text": "5/27 = 중간 측정 (10 cell main) + sf100 (5 cell) = ", "size": 11, "color": INK_SOFT},
                         {"text": "15 cell narrative", "size": 11, "bold": True, "color": NAVY},
                         {"text": ". 본 deck 의 ToC 흐름 (Problem → Prior → Approach → RQ1/2/3 → Multi → Honest → Future) 를 5/27 12 슬라이드로 압축 + sf100 추가 1 슬라이드.",
                          "size": 11, "color": INK_SOFT},
                     ])

    add_footer(s, page=14)

    # =====================================================
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Wrote: {OUT_PATH}  ({OUT_PATH.stat().st_size:,} bytes)")
    print(f"Slides: {len(prs.slides)}")


def slide_card_blank(slide, x, y, w, h):
    """Blank rounded white card for TOC."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.04
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE; shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def slide_swatch(slide, x, y, w, h, label, fill, text_color):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), label, font=MONO_FONT, size=8, bold=True, color=text_color)
    return box


if __name__ == "__main__":
    main()
