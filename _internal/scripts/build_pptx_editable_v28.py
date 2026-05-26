#!/usr/bin/env python3
"""
PPTX (a) editable — deck_v28 16 슬라이드 native python-pptx 작성.
영상 슬라이드 2개에 Claude Design URL hyperlink + ▶ 큰 버튼.

Capstone Design System:
- font: Apple SD Gothic Neo
- 색상: navy #1B2A4E (anchor), cyan #4FB3D9 (highlight), red #D85A5A (catastrophic), 흰 배경, 회색 #6B7280
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# 색상
NAVY = RGBColor(0x1B, 0x2A, 0x4E)
CYAN = RGBColor(0x4F, 0xB3, 0xD9)
RED = RGBColor(0xD8, 0x5A, 0x5A)
INK = RGBColor(0x1F, 0x2A, 0x3F)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GREY = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xD9, 0x84, 0x4F)
LIME = RGBColor(0x6F, 0xA8, 0x4F)
VIOLET = RGBColor(0x6B, 0x4F, 0xA8)

FONT_KR = "Apple SD Gothic Neo"

ANIM1_URL = "https://claude.ai/design/p/455bb91a-3c1e-48a4-b8ee-8932f4d8881d?file=index.html&present=1"
ANIM2_URL = "https://claude.ai/design/p/585e3fa7-daae-4cfc-8d8f-a0cf966299ea?file=video.html&present=1"
DECK_URL = "https://claude.ai/design/p/019e1a41-701c-7134-9ce1-1247262c1563?file=deck_v28.html"


def add_text(slide, left, top, width, height, text, *, font=FONT_KR, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_box(slide, left, top, width, height, fill=None, line_color=None, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    return sh


def add_hyperlink(slide, left, top, width, height, text, url, *, font=FONT_KR, size=14,
                  fill=NAVY, text_color=WHITE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    r.hyperlink.address = url
    return sh


def add_brand_chip(slide):
    """좌상단 brand chip"""
    add_text(slide, Inches(0.4), Inches(0.25), Inches(2.5), Inches(0.3),
             "■ 속도는벡터 · Capstone 2026-1", size=10, color=GREY)


def add_footer(slide, footer_text="DEEP sf=10 · sel=0.001 · Q3 · 15-trial trim mean"):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(13), Inches(0.3),
             footer_text, size=9, color=GREY, align=PP_ALIGN.LEFT)


# =============================================================================
# 슬라이드별 빌드 함수
# =============================================================================
def s1_cover(slide):
    """Slide 1 — 표지"""
    add_brand_chip(slide)
    # title
    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
             "Vector-augmented Analytical Query에서", size=24, color=GREY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.4),
             "카디널리티 추정 정확도가 실행 계획을 결정한다", size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
             "Exqutor §V-B Adaptive Sampling의 표본 선택 단계 — 분포 인지 결합 효과 검증",
             size=14, color=INK, align=PP_ALIGN.CENTER, italic=True)
    # team
    add_text(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
             "팀 속도는벡터 — 박세은 · 강재현 · 조현빈 · 이동욱", size=13, color=GREY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.95), Inches(11.3), Inches(0.4),
             "지도교수 박광현 · 지도연구원 임채림 · 멘토 박성원 (Samsung AI Center)",
             size=11, color=GREY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.55), Inches(11.3), Inches(0.3),
             "연세대학교 컴퓨터과학과 · 2026-1학기 캡스톤 디자인 최종발표",
             size=10, color=GREY, align=PP_ALIGN.CENTER)


def s2_motivation(slide):
    """Slide 2 — 동기"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "① 동기", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "VAQ 시대, planner의 카디 추정이 실행 계획을 결정한다", size=28, bold=True, color=NAVY)
    add_text(slide, Inches(0.6), Inches(2.3), Inches(12.5), Inches(0.5),
             "벡터 술어가 결합된 분석 쿼리 = pgvector·DuckDB·VBASE에서 카디 추정 고정 비율(33%·50%·100%) → 잘못된 plan",
             size=14, color=INK)
    # 3 카드
    cards = [("pgvector", "33.3% 고정", "v_pred 매칭 1/3"),
             ("VBASE", "50% 고정", "v_pred 매칭 1/2"),
             ("DuckDB", "100% 고정", "v_pred 매칭 전부")]
    for i, (n, r, d) in enumerate(cards):
        x = Inches(0.6 + i * 4.3)
        add_box(slide, x, Inches(3.5), Inches(4.0), Inches(2.5), fill=WHITE, line_color=NAVY, line_w=1.5)
        add_text(slide, x, Inches(3.7), Inches(4.0), Inches(0.5), n, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(4.5), Inches(4.0), Inches(0.5), r, size=24, bold=True, color=RED, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(5.2), Inches(4.0), Inches(0.5), d, size=12, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(slide, "Exqutor (arXiv:2512.09695v2) §IV — 기존 엔진 분석")


def s3_problem(slide):
    """Slide 3 — 문제 정의"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "② 문제", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "Adaptive Sampling의 표본 선택 단계가 약점", size=28, bold=True, color=NAVY)
    # Exqutor 박스
    add_box(slide, Inches(0.6), Inches(2.7), Inches(12.7), Inches(1.5), fill=LIGHT_GREY, line_color=GREY, line_w=0.5)
    add_text(slide, Inches(0.9), Inches(2.85), Inches(12.0), Inches(0.4),
             "Exqutor §V-B Adaptive Sampling (논문)", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(0.9), Inches(3.3), Inches(12.0), Inches(0.8),
             "표본 N=385 · Bernoulli random sampling · 모멘텀 식 1~6으로 동적 조정",
             size=14, color=INK)
    # 약점 박스
    add_box(slide, Inches(0.6), Inches(4.5), Inches(12.7), Inches(2.0), fill=None, line_color=RED, line_w=1.5)
    add_text(slide, Inches(0.9), Inches(4.65), Inches(12.0), Inches(0.4),
             "★ 약점 — 표본 선택 (sample selection) 자체는 분포 무관 random Bernoulli",
             size=13, bold=True, color=RED)
    add_text(slide, Inches(0.9), Inches(5.1), Inches(12.0), Inches(1.4),
             "• 선택적 쿼리 (sel=0.001 등)에서 표본 안 매칭 0건 발생 가능 → 추정 0 → catastrophic plan\n"
             "• 데이터 분포가 skew일수록 random sampling은 표본 효율 저하\n"
             "• 분포 인지 stratification으로 표본 선택만 강화하면 회복 가능한가? — 본 연구 질문",
             size=13, color=INK)
    add_footer(slide, "Exqutor §V-B (식 1-6) · 본 연구 = 표본 선택 단계 개입의 전 변인 검증")


def s4_research_question(slide):
    """Slide 4 — Research Question"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "③ Research Question", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "분포 인지 결합으로 catastrophic plan을 회복할 수 있는가?",
             size=28, bold=True, color=NAVY)
    # 3-way 구조
    add_text(slide, Inches(0.6), Inches(2.7), Inches(12.5), Inches(0.5),
             "3-way matched 측정 (한 측정 = 세 mode 동시 산출)", size=14, bold=True, color=INK)
    rows = [
        ("B1 (대조군)", "Exqutor Bernoulli 그대로", "논문 원본 재현", NAVY),
        ("CaseA (음성 대조)", "Bernoulli → method 표본 완전 대체", "표본 출처만 바꿈", GREY),
        ("CaseB (본 연구)", "est_final = (B1 + method) / 2", "분포 인지 산술 평균 결합", CYAN),
    ]
    for i, (n, d, e, c) in enumerate(rows):
        y = Inches(3.4 + i * 1.0)
        add_box(slide, Inches(0.6), y, Inches(12.7), Inches(0.85), fill=WHITE, line_color=c, line_w=1.5)
        add_text(slide, Inches(0.9), y + Inches(0.1), Inches(3.0), Inches(0.65),
                 n, size=15, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(4.0), y + Inches(0.1), Inches(5.0), Inches(0.65),
                 d, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(9.0), y + Inches(0.1), Inches(4.0), Inches(0.65),
                 e, size=12, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "13 method × 5 dataset × 5 조작변인 = 1,508 측정 (3-way matched)")


def s5_animation1(slide):
    """Slide 5 — Animation 1 영상 슬라이드 (왜 정확한 카디 추정이 중요한가)"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "④ 영상 ① — 왜 정확한 카디 추정이 중요한가", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "33행 vs 7,603행 — plan 분기 · 자원 7배 차이", size=24, bold=True, color=NAVY)
    # 큰 영상 placeholder 박스
    add_box(slide, Inches(0.6), Inches(2.4), Inches(12.7), Inches(4.0), fill=NAVY, line_color=None)
    # ▶ 큰 버튼 영역
    add_text(slide, Inches(0.6), Inches(2.9), Inches(12.7), Inches(1.0),
             "▶", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.1), Inches(12.7), Inches(0.5),
             "카드 클릭 → Claude Design Animation 1 fullscreen 재생 (32초)",
             size=16, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.7), Inches(12.7), Inches(0.5),
             "32초 timeline motion · split-screen 두 분기 (부정확 7,242ms vs 정확 1,000ms)",
             size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.2), Inches(12.7), Inches(0.4),
             "키보드 컨트롤: space 재생/정지 · ←→ 0.1s 이동 · shift+←→ 1s · 0 처음",
             size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    # 큰 hyperlink 버튼
    add_hyperlink(slide, Inches(5.0), Inches(5.7), Inches(3.9), Inches(0.55),
                  "▶ 영상 재생 (Claude Design)", ANIM1_URL, size=14, fill=CYAN, text_color=NAVY)
    add_footer(slide, "raw · phase2 · DEEP sf=10 · Q3 · sel=0.001 · 15-trial trim mean")


def s6_method_design(slide):
    """Slide 6 — 방법론"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑤ 방법", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "13 method × 7 paradigm — 분포 인지 stratification", size=24, bold=True, color=NAVY)
    paradigms = [
        ("Space-filling curve", "hilbert_real · skilling_hilbert · zorder_morton", RED),
        ("Dimensionality reduction", "sparse_rp · pca1d · rsvd · ica_fastica", NAVY),
        ("Stratified sampling", "cum_sqrtf · lavallee_hidiroglou", AMBER),
        ("Quantization / grid bucketing", "mhist2 · rabitq_strat", LIME),
        ("Weighted reservoir sampling", "chao_weighted ★ winner", CYAN),
        ("Hash bucketing", "hyperloglog", VIOLET),
        ("Clustering", "(본 13에서 제외 · 정합성 audit)", GREY),
    ]
    for i, (n, m, c) in enumerate(paradigms):
        y = Inches(2.6 + i * 0.65)
        add_text(slide, Inches(0.8), y, Inches(0.3), Inches(0.5),
                 "●", size=18, color=c, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.3), y, Inches(4.5), Inches(0.5),
                 n, size=15, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(6.0), y, Inches(7.3), Inches(0.5),
                 m, size=13, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "method audit 5/19 · _internal/METHOD_REGISTRY.md")


def s7_experiment(slide):
    """Slide 7 — 실험 설계"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑥ 실험 설계", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "1,508 측정 (3-way matched) + 156 cell latency", size=24, bold=True, color=NAVY)
    # 표
    headers = ["축", "값", "비고"]
    rows = [
        ("Dataset", "DEEP · SIFT · SSN · YFCC · WIKI / +DEEP+WIKI · DEEP+YFCC concat", "5 base + 2 concat"),
        ("Scale Factor", "sf=1 / sf=10 / sf=100", "전 변인"),
        ("Selectivity", "sel=0.001 / 0.01 / 0.1", "전 변인 (선택적·중간·일반)"),
        ("Query", "Q3 / Q9 / Q10 / Q12 (TPC-H 벡터 술어)", "qid 0·1·2 각 × 3"),
        ("Method", "13 method × 7 paradigm", "registry 5/19"),
        ("Trials / Latency", "15-trial trim mean · n_warmup=1 · statement_timeout 600s", "Exqutor patched PG 55435"),
    ]
    # header
    for i, h in enumerate(headers):
        x = [Inches(0.6), Inches(2.6), Inches(10.8)][i]
        w = [Inches(2.0), Inches(8.2), Inches(2.5)][i]
        add_box(slide, x, Inches(2.6), w, Inches(0.5), fill=NAVY)
        add_text(slide, x, Inches(2.65), w, Inches(0.4), h,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # body
    for ri, row in enumerate(rows):
        y = Inches(3.1 + ri * 0.55)
        for i, v in enumerate(row):
            x = [Inches(0.6), Inches(2.6), Inches(10.8)][i]
            w = [Inches(2.0), Inches(8.2), Inches(2.5)][i]
            add_box(slide, x, y, w, Inches(0.55), fill=WHITE, line_color=LIGHT_GREY, line_w=0.5)
            add_text(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.45),
                     v, size=11, color=INK if i < 2 else GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "phase2 · phase3 · phase4_extension = 30 latency cells")


def s8_rq3_combined(slide):
    """Slide 8 — RQ3 결과 핵심"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑦ 결과", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "결합 vs 대조군 — 89.1% case에서 Q-error 개선", size=24, bold=True, color=NAVY)
    # 핵심 수치 카드 3개
    cards = [
        ("89.1%", "결합 vs B1 better", "1,344 / 1,508 case", CYAN),
        ("−4.38%", "Q-error 중앙값 개선", "B1 대비 paired Δ", CYAN),
        ("p < 0.001", "Wilcoxon signed-rank", "통계 유의", NAVY),
    ]
    for i, (big, mid, sub, c) in enumerate(cards):
        x = Inches(0.6 + i * 4.3)
        add_box(slide, x, Inches(2.6), Inches(4.0), Inches(2.8), fill=WHITE, line_color=c, line_w=2)
        add_text(slide, x, Inches(2.9), Inches(4.0), Inches(1.2),
                 big, size=44, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(4.15), Inches(4.0), Inches(0.5),
                 mid, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(4.7), Inches(4.0), Inches(0.5),
                 sub, size=11, color=GREY, align=PP_ALIGN.CENTER)
    # negative control
    add_box(slide, Inches(0.6), Inches(5.8), Inches(12.7), Inches(1.0),
            fill=LIGHT_GREY, line_color=GREY, line_w=0.5)
    add_text(slide, Inches(0.9), Inches(5.95), Inches(12.0), Inches(0.4),
             "Negative control (CaseA · 완전 대체)", size=12, bold=True, color=GREY)
    add_text(slide, Inches(0.9), Inches(6.4), Inches(12.0), Inches(0.4),
             "Better 35.2% (530 / 1,508) — 단순 표본 출처 교체만으로는 효과 없음 → 결합 평균이 핵심",
             size=12, color=INK)
    add_footer(slide, "v13 정본 · 3-way matched 1,508 측정 · raw _internal/cache/rq3/v13_summary.md")


def s9_paradigm_winner(slide):
    """Slide 9 — 11.5 신규 = 7 paradigm 대표 + winner"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑧ 7 paradigm 대표 method", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "Winner = chao_weighted (Weighted reservoir, Chao 1982)", size=22, bold=True, color=NAVY)
    headers = ["Paradigm", "대표 method", "추정", "Q-err", "lat (ms)", "Δ%", "회복"]
    widths = [Inches(2.8), Inches(2.7), Inches(1.0), Inches(0.95), Inches(1.2), Inches(1.05), Inches(0.85)]
    rows = [
        ("Space-filling curve", "hilbert_real ⚠️", "8,546", "1.1241", "1,063.75", "+4.49%", "✗"),
        ("Dimensionality red.", "sparse_rp", "5,105", "1.4892", "1,032.09", "+1.38%", "✓"),
        ("Stratified sampling", "cum_sqrtf", "6,226", "1.2211", "990.01", "−2.76%", "✓"),
        ("Quantization", "mhist2", "3,457", "2.1994", "944.53", "−7.22%", "✓"),
        ("Weighted reservoir ★", "chao_weighted ★", "7,158", "1.0622", "927.05", "−8.94%", "✓"),
        ("Hash bucketing", "hyperloglog", "3,222", "2.3596", "1,021.75", "+0.36%", "✓"),
        ("Clustering", "—", "—", "—", "—", "—", "—"),
    ]
    # header
    x = Inches(0.6)
    for i, h in enumerate(headers):
        add_box(slide, x, Inches(2.5), widths[i], Inches(0.5), fill=NAVY)
        add_text(slide, x, Inches(2.55), widths[i], Inches(0.4), h,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + widths[i]
    # body
    for ri, row in enumerate(rows):
        y = Inches(3.0 + ri * 0.5)
        x = Inches(0.6)
        is_winner = "chao_weighted" in row[1]
        bg = CYAN if is_winner else (LIGHT_GREY if row[1] == "—" else WHITE)
        for i, v in enumerate(row):
            add_box(slide, x, y, widths[i], Inches(0.5), fill=bg, line_color=LIGHT_GREY, line_w=0.5)
            txt_color = NAVY if is_winner else (GREY if row[1] == "—" else INK)
            add_text(slide, x + Inches(0.05), y + Inches(0.05), widths[i] - Inches(0.1), Inches(0.4), v,
                     size=10, bold=is_winner, color=txt_color,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x = x + widths[i]
    add_text(slide, Inches(0.6), Inches(6.6), Inches(12.7), Inches(0.4),
             "★ chao_weighted = Q-error 1.0622 · latency −8.94% · oracle plan 회복 ✓ 동시 달성 (DEEP sf=10 sel=0.001 Q3 qid=0)",
             size=11, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    add_footer(slide, "raw · DEEP sf=10 sel=0.001 Q3 qid=0 · ⚠️ hilbert_real = PCA 2D lex sort alias honest")


def s10_scatter(slide):
    """Slide 10 — 11.6 신규 = 13 method scatter 텍스트 ver"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑨ 13 method 전체 비교", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "Q-error 개선 ≠ latency 개선 — 12/13 oracle plan 회복", size=22, bold=True, color=NAVY)
    headers = ["#", "Method", "추정", "Q-err", "lat (ms)", "Δ%", "회복"]
    widths = [Inches(0.5), Inches(2.8), Inches(1.4), Inches(1.4), Inches(1.6), Inches(1.4), Inches(1.0)]
    rows = [
        ("1", "chao_weighted ★", "7,158", "1.0622", "927.05", "−8.94%", "✓"),
        ("2", "hilbert_real ⚠️", "8,546", "1.1241", "1,063.75", "+4.49%", "✗"),
        ("3", "cum_sqrtf", "6,226", "1.2211", "990.01", "−2.76%", "✓"),
        ("4", "sparse_rp", "5,105", "1.4892", "1,032.09", "+1.38%", "✓"),
        ("5", "lavallee_hidiroglou", "5,069", "1.5000", "984.67", "−3.28%", "✓"),
        ("6", "zorder_morton", "4,211", "1.8055", "1,020.81", "+0.27%", "✓"),
        ("7", "mhist2", "3,457", "2.1994", "944.53", "−7.22%", "✓"),
        ("8", "rsvd", "3,367", "2.2580", "1,008.68", "−0.92%", "✓"),
        ("9", "hyperloglog", "3,222", "2.3596", "1,021.75", "+0.36%", "✓"),
        ("10", "pca1d", "3,158", "2.4076", "950.14", "−6.67%", "✓"),
        ("11", "skilling_hilbert", "3,053", "2.4906", "979.70", "−3.77%", "✓"),
        ("12", "rabitq_strat", "1,330", "5.7175", "1,006.40", "−1.15%", "✓"),
        ("13", "ica_fastica", "1,053", "7.2228", "932.47", "−8.41%", "✓"),
    ]
    # header
    x = Inches(0.6)
    for i, h in enumerate(headers):
        add_box(slide, x, Inches(2.4), widths[i], Inches(0.4), fill=NAVY)
        add_text(slide, x, Inches(2.45), widths[i], Inches(0.3), h,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + widths[i]
    # body
    for ri, row in enumerate(rows):
        y = Inches(2.8 + ri * 0.33)
        x = Inches(0.6)
        is_winner = ri == 0
        bg = CYAN if is_winner else WHITE
        for i, v in enumerate(row):
            add_box(slide, x, y, widths[i], Inches(0.33), fill=bg, line_color=LIGHT_GREY, line_w=0.3)
            color = NAVY if is_winner else INK
            add_text(slide, x + Inches(0.05), y, widths[i] - Inches(0.1), Inches(0.33), v,
                     size=9, bold=is_winner, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x = x + widths[i]
    add_text(slide, Inches(0.6), Inches(7.1), Inches(12.7), Inches(0.3),
             "★ chao_weighted = Q-error 최저 + latency 최저 동시 winner · ⚠️ hilbert_real만 oracle plan 미회복",
             size=10, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    add_footer(slide, "raw · phase2 · DEEP sf=10 · Q3 · sel=0.001 · qid=0 · 15-trial trim mean")


def s11_latency_corrected(slide):
    """Slide 11 — latency 배수 정정 후"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑩ Latency 결과 — Adaptive Sampling 효과", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "기본 5,677 ms → Adaptive 약 980 ms (5.81× 더 빠름)", size=24, bold=True, color=NAVY)
    # 3 막대
    bars = [
        ("pgvector 기본", "5,677 ms", "(기준 1.0×)", 12.0, GREY),
        ("Exqutor B1 (논문 원본)", "977.6 ms", "5.81× 더 빠름 ↓", 2.07, NAVY),
        ("결합 CaseB (본 연구)", "983.5 ms", "5.77× 더 빠름 ↓", 2.08, CYAN),
    ]
    for i, (label, val, multiplier, bar_inch, color) in enumerate(bars):
        y = Inches(2.6 + i * 1.0)
        add_text(slide, Inches(0.6), y, Inches(4.0), Inches(0.4),
                 label, size=13, bold=True, color=INK)
        add_text(slide, Inches(0.6), y + Inches(0.35), Inches(4.0), Inches(0.4),
                 multiplier, size=11, color=color, italic=True)
        add_box(slide, Inches(4.7), y + Inches(0.05), Inches(bar_inch), Inches(0.55), fill=color)
        # 값 표시
        text_x = Inches(4.7 + bar_inch + 0.15)
        add_text(slide, text_x, y + Inches(0.1), Inches(2.0), Inches(0.45),
                 val, size=14, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    # 메시지
    add_box(slide, Inches(0.6), Inches(6.0), Inches(12.7), Inches(0.85), fill=LIGHT_GREY, line_color=CYAN, line_w=1)
    add_text(slide, Inches(0.9), Inches(6.15), Inches(12.0), Inches(0.55),
             "baseline 과 결합 latency 사실상 동등 (격차 +0.60%, 5.9 ms)",
             size=14, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "★ 강재현 지적 정정: 5.77→5.81 / 5.70→5.77 / ↑→↓ · raw 30 cell latency")


def s12_animation2(slide):
    """Slide 12 — Animation 2 영상 슬라이드 (본 연구 결합 동작)"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑪ 영상 ② — 본 연구 결합 엔진 동작", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "Exqutor B1 vs CaseB — 분포 인지 결합으로 plan 회복", size=22, bold=True, color=NAVY)
    # 큰 영상 placeholder 박스
    add_box(slide, Inches(0.6), Inches(2.4), Inches(12.7), Inches(4.0), fill=CYAN, line_color=None)
    add_text(slide, Inches(0.6), Inches(2.9), Inches(12.7), Inches(1.0),
             "▶", size=72, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.1), Inches(12.7), Inches(0.5),
             "카드 클릭 → Claude Design Animation 2 fullscreen 재생 (45초)",
             size=16, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.7), Inches(12.7), Inches(0.5),
             "45초 timeline · DEEP 점 cloud + B1(매칭 0건) vs CaseB(양수 매칭) + plan 회복 도넛 94.9% vs 58.3%",
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.2), Inches(12.7), Inches(0.4),
             "키보드: space 재생/정지 · ←→ 0.1s · shift+←→ 1s · 0 처음",
             size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_hyperlink(slide, Inches(5.0), Inches(5.7), Inches(3.9), Inches(0.55),
                  "▶ 영상 재생 (Claude Design)", ANIM2_URL, size=14, fill=NAVY, text_color=CYAN)
    add_footer(slide, "raw · phase2~4 · 156 cell · 13 method · DEEP/SIFT/SSN/YFCC sf=10")


def s13_honest_limitation(slide):
    """Slide 13 — Honest limitations"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑫ Honest limitations", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "정직한 한계 — 효과 크기·해석·일반화", size=24, bold=True, color=NAVY)
    items = [
        ("중앙값 효과는 작음", "Q-error 중앙값 개선 −4.38% — 통계적 유의 vs 실용적 효과의 차이",
         "응답 시간 개선은 cell-level noise 안에 묶임 (12/13 method oracle plan 회복 시)"),
        ("paradigm 일관성", "P1 Cluster 본 13에서 제외 — 정합성·정합성 audit 통과 못 함",
         "minibatch_partial · gmm · faiss_ivf 별도 분석"),
        ("hilbert_real ★3", "Faloutsos 1989 ❌ → PCA 2D lex sort alias",
         "honest 표시 + zorder_morton honest 후보로 병기"),
        ("concat sf=100 부분 미측정", "DEEP+WIKI / DEEP+YFCC concat sf=100 일부 cell 미측정",
         "측정 시간·리소스 제약 — 보고서 §10에 명시"),
    ]
    for i, (t, d, r) in enumerate(items):
        y = Inches(2.4 + i * 1.05)
        add_box(slide, Inches(0.6), y, Inches(12.7), Inches(0.95),
                fill=WHITE, line_color=LIGHT_GREY, line_w=0.5)
        add_text(slide, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(slide, Inches(0.9), y + Inches(0.5), Inches(3.5), Inches(0.4),
                 d, size=11, color=GREY)
        add_text(slide, Inches(4.7), y + Inches(0.25), Inches(8.5), Inches(0.55),
                 r, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "보고서 §10 · _internal/method_audit/ · 학술 정직성 우선")


def s14_future_work(slide):
    """Slide 14 — Future Work ('두 갈래' 제거)"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑬ Future Work", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "Future Work", size=28, bold=True, color=NAVY)
    # 2 카드
    groups = [
        ("Group A — 메커니즘 후속",
         ["plan-level 효과 분해 — Q-error → plan 회복 → latency 인과 사슬",
          "본 연구 cell 28% (latency 차이 발생 cell)에서 plan tree 비교",
          "PG planner heuristic 임계값 추정 — 카디 추정 어느 정도면 다른 plan 선택?",
          "AAA conference 후보 — VLDB · SIGMOD"],
         CYAN),
        ("Group B — 산업 적용 framework",
         ["pgvector · DuckDB · Milvus · VBASE 4 엔진 통합 측정 (6/11 보고서)",
          "method registry → production-ready library 패키지",
          "Exqutor 패치 → 분포 인지 결합 옵션 — 1줄 PG config",
          "산업체 평가 — Naver · Kakao · Samsung 한국 vector DB 도입 case"],
         NAVY),
    ]
    for i, (title, items, c) in enumerate(groups):
        x = Inches(0.6 + i * 6.45)
        add_box(slide, x, Inches(2.5), Inches(6.25), Inches(4.0), fill=WHITE, line_color=c, line_w=2)
        add_text(slide, x + Inches(0.3), Inches(2.7), Inches(5.95), Inches(0.5),
                 title, size=16, bold=True, color=c)
        for j, item in enumerate(items):
            add_text(slide, x + Inches(0.4), Inches(3.4 + j * 0.65), Inches(5.85), Inches(0.6),
                     "• " + item, size=11, color=INK)
    add_footer(slide, "★ 강재현 지적 정정: '두 갈래' 텍스트 제거 — Group A·B 시각 자체가 두 갈래 표현")


def s15_team(slide):
    """Slide 15 — 팀 / 기여"""
    add_brand_chip(slide)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12.5), Inches(0.5),
             "⑭ 팀 · 기여", size=12, bold=True, color=CYAN)
    add_text(slide, Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.0),
             "팀 속도는벡터 — Capstone 2026-1", size=24, bold=True, color=NAVY)
    team = [
        ("박세은", "팀장", "최종발표 · 슬라이드 디자인 · narrative", NAVY),
        ("강재현", "팀원", "데이터 검증 · figure · 결과 시각화", CYAN),
        ("조현빈", "팀원", "측정 자동화 · v13 정본 · method audit", AMBER),
        ("이동욱", "팀원", "엔진 적용 검증 · 분석", LIME),
    ]
    for i, (name, role, contrib, c) in enumerate(team):
        x = Inches(0.6 + i * 3.2)
        add_box(slide, x, Inches(2.6), Inches(3.0), Inches(2.5), fill=WHITE, line_color=c, line_w=1.5)
        add_text(slide, x, Inches(2.8), Inches(3.0), Inches(0.5), name,
                 size=20, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(3.4), Inches(3.0), Inches(0.4), role,
                 size=12, color=GREY, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(4.0), Inches(3.0), Inches(1.0), contrib,
                 size=10, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.5), Inches(12.5), Inches(0.5),
             "지도교수 박광현 · 지도연구원 임채림 · 멘토 박성원 (Samsung AI Center)",
             size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(6.0), Inches(12.5), Inches(0.4),
             "연세대학교 컴퓨터과학과 · BDAI Lab", size=12, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(slide, "github.com/johyunbin/Capstone · notion 팀 페이지")


def s16_thanks(slide):
    """Slide 16 — 감사합니다"""
    add_brand_chip(slide)
    # 큰 감사합니다
    add_text(slide, Inches(0.6), Inches(2.8), Inches(12.7), Inches(2.0),
             "감사합니다", size=88, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.8), Inches(12.7), Inches(0.5),
             "Q&A 환영합니다", size=20, color=GREY, align=PP_ALIGN.CENTER)
    # Claude Design 링크
    add_hyperlink(slide, Inches(4.8), Inches(5.8), Inches(4.4), Inches(0.5),
                  "📎 Claude Design deck v28 (원본)", DECK_URL, size=12, fill=LIGHT_GREY, text_color=NAVY)
    add_footer(slide, "DEEP/SIFT/SSN/YFCC sf=10 · 1,508 측정 · 13 method · 7 paradigm · v13 정본")


# =============================================================================
# Main
# =============================================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    builders = [s1_cover, s2_motivation, s3_problem, s4_research_question,
                s5_animation1, s6_method_design, s7_experiment, s8_rq3_combined,
                s9_paradigm_winner, s10_scatter, s11_latency_corrected, s12_animation2,
                s13_honest_limitation, s14_future_work, s15_team, s16_thanks]

    blank_layout = prs.slide_layouts[6]
    for i, build in enumerate(builders, start=1):
        slide = prs.slides.add_slide(blank_layout)
        # 배경 흰색
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        build(slide)

    out = "/Users/hyunbin/Capstone/submission/_drafts/속도는벡터_기말발표_editable_v4_영상hyperlink_20260526_154845.pptx"
    prs.save(out)
    print(f"✓ {out}")
    print(f"슬라이드: {len(prs.slides)}")


if __name__ == "__main__":
    main()
