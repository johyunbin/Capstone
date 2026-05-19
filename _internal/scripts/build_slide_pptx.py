#!/usr/bin/env python3
"""build_slide_pptx.py — 슬라이드 PNG 묶음 → 16:9 이미지-슬라이드 pptx.

사용법: python3 build_slide_pptx.py <png_dir> <output.pptx>
각 PNG(이름순)를 한 슬라이드에 full-bleed로 배치 → PowerPoint/Keynote
슬라이드쇼 녹화 기능으로 내레이션을 입혀 MP4로 내보낼 수 있는 deck.
"""

import glob
import os
import sys

from pptx import Presentation
from pptx.util import Inches


def build(png_dir, out_pptx):
    pngs = sorted(glob.glob(os.path.join(png_dir, "*.png")))
    if not pngs:
        print(f"✗ {png_dir} 에 PNG 없음")
        sys.exit(1)
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]       # 빈 레이아웃
    for p in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            p, Inches(0), Inches(0),
            width=prs.slide_width, height=prs.slide_height,
        )
    prs.save(out_pptx)
    print(f"✓ {out_pptx}  ({len(pngs)}장, 16:9)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("사용법: python3 build_slide_pptx.py <png_dir> <output.pptx>")
        sys.exit(1)
    build(sys.argv[1], sys.argv[2])
