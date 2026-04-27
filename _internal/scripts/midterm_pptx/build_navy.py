"""
Version A — Navy Corporate (Samsung Research 정통)
어두운 표지 + 흰 본문 + 파란 헤더 바 + 화살표 시사점 박스.
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from common import (
    new_presentation, blank_slide, set_bg, add_text, add_paragraph,
    add_rect, add_round_rect, add_oval, slide_number, add_figure, SLIDE_W,
)
import content as C
from theme import NAVY as TH

TOTAL = 17


# ─────────────────────────────────────────
# Helper — 본문 슬라이드 공통 헤더
# ─────────────────────────────────────────
def content_header(slide, title, subtitle=None):
    set_bg(slide, TH.bg_light)
    # 상단 얇은 컬러 바
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(5), TH.primary)
    # 좌측 vertical bar (제목·부제 영역의 왼쪽 align 라인) — 위/아래 정렬 문제 해결
    add_rect(slide, Inches(0.4), Inches(0.30), Pt(3), Inches(0.85), TH.primary)
    # 제목 (vertical anchor middle 로 박스 가운데 정렬)
    add_text(slide, Inches(0.6), Inches(0.20), Inches(11.5), Inches(0.55),
             title, size=22, bold=True, color=TH.primary, font=TH.font_main,
             anchor=MSO_ANCHOR.MIDDLE)
    # 부제
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(12.0), Inches(0.35),
                 subtitle, size=11, color=TH.text_dim, font=TH.font_main,
                 anchor=MSO_ANCHOR.MIDDLE)
    # 헤더 하단 분리선 — y=1.25 → 1.20 (제목/부제와 0.07 spacing 확보)
    add_rect(slide, Inches(0.6), Inches(1.20), Inches(12.1), Pt(1), TH.border)
    # 분리선 ↔ 본문 spacing 0.20 inch — 본문 함수는 y=1.45 부터 시작


# ─────────────────────────────────────────
# S1. Cover
# ─────────────────────────────────────────
def s1_cover(prs):
    s = blank_slide(prs)
    set_bg(s, TH.bg_dark)

    # 좌상단 작은 accent line
    add_rect(s, Inches(0.8), Inches(1.6), Inches(1.2), Pt(3), TH.accent)

    # 카테고리 라벨
    add_text(s, Inches(0.8), Inches(1.75), Inches(8), Inches(0.35),
             "CAPSTONE 2026-1 · MIDTERM PRESENTATION",
             size=11, bold=True, color=TH.accent, font=TH.font_main)

    # 메인 타이틀
    add_text(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.0),
             C.S1_COVER["title"], size=32, bold=True, color=TH.text_white,
             font=TH.font_main, line_spacing=1.3)

    # 부제
    add_text(s, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.6),
             C.S1_COVER["subtitle"], size=15, color=TH.accent,
             font=TH.font_main, line_spacing=1.4)

    # 분리선
    add_rect(s, Inches(0.8), Inches(5.0), Inches(5.5), Pt(1),
             TH.text_dim)

    # 팀 이름 + 소속
    add_text(s, Inches(0.8), Inches(5.2), Inches(8), Inches(0.5),
             f"{C.S1_COVER['team']}",
             size=22, bold=True, color=TH.text_white, font=TH.font_main)
    add_text(s, Inches(0.8), Inches(5.85), Inches(10), Inches(0.35),
             C.S1_COVER["course"], size=12, color=TH.text_dim, font=TH.font_main)
    add_text(s, Inches(0.8), Inches(6.20), Inches(10), Inches(0.35),
             C.S1_COVER["advisor"], size=11, color=TH.text_dim, font=TH.font_main)

    # 우측 발표 정보
    add_text(s, Inches(10.0), Inches(6.20), Inches(3), Inches(0.35),
             C.S1_COVER["date"],
             size=11, color=TH.text_dim, align=PP_ALIGN.RIGHT, font=TH.font_main)
    add_text(s, Inches(8.0), Inches(6.55), Inches(5), Inches(0.35),
             " · ".join(C.S1_COVER["members"]),
             size=10, color=TH.text_dim, align=PP_ALIGN.RIGHT, font=TH.font_main)


# ─────────────────────────────────────────
# S2. Table of Contents
# ─────────────────────────────────────────
def s2_toc(prs):
    s = blank_slide(prs)
    content_header(s, C.S2_TOC["title"])
    slide_number(s, 2, TOTAL, TH.text_dim)

    for i, (num, title, desc) in enumerate(C.S2_TOC["sections"]):
        y = Inches(1.65 + i * 1.05)
        # 번호 배지
        add_round_rect(s, Inches(0.8), y, Inches(0.85), Inches(0.7), TH.bg_dark)
        add_text(s, Inches(0.8), y + Inches(0.13), Inches(0.85), Inches(0.5),
                 num, size=18, bold=True, color=TH.text_white,
                 align=PP_ALIGN.CENTER, font=TH.font_main)
        # 제목
        add_text(s, Inches(1.85), y + Inches(0.05), Inches(3.5), Inches(0.5),
                 title, size=18, bold=True, color=TH.text, font=TH.font_main)
        # 설명
        add_text(s, Inches(5.4), y + Inches(0.13), Inches(7.5), Inches(0.5),
                 desc, size=12, color=TH.text_dim, font=TH.font_main)


# ─────────────────────────────────────────
# S3. Problem
# ─────────────────────────────────────────
def s3_problem(prs):
    s = blank_slide(prs)
    content_header(s, C.S3_PROBLEM["title"])
    slide_number(s, 3, TOTAL, TH.text_dim)

    # 헤드라인
    add_text(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.5),
             C.S3_PROBLEM["headline"], size=17, bold=True,
             color=TH.text, font=TH.font_main)

    # 좌측: 고정 비율 카드 3개
    add_text(s, Inches(0.6), Inches(2.15), Inches(5.5), Inches(0.35),
             "기존 시스템 — 고정 비율 추정", size=12, bold=True,
             color=TH.primary, font=TH.font_main)
    for i, (sysname, ratio, note) in enumerate(C.S3_PROBLEM["fixed_ratio"]):
        y = Inches(2.55 + i * 0.95)
        add_round_rect(s, Inches(0.6), y, Inches(5.5), Inches(0.85),
                       TH.card, border=TH.border)
        # 시스템 이름
        add_text(s, Inches(0.85), y + Inches(0.10), Inches(2.0), Inches(0.4),
                 sysname, size=15, bold=True, color=TH.primary, font=TH.font_main)
        # 비율 (큰 숫자)
        add_text(s, Inches(2.7), y + Inches(0.05), Inches(1.4), Inches(0.55),
                 ratio, size=24, bold=True, color=TH.text, font=TH.font_main)
        # 설명
        add_text(s, Inches(0.85), y + Inches(0.50), Inches(5.0), Inches(0.3),
                 note, size=10, color=TH.text_dim, font=TH.font_main)

    # 우측: 실제 분포의 폭 + 영향
    add_text(s, Inches(6.5), Inches(2.15), Inches(6.3), Inches(0.35),
             "실제 — 0.001% ~ 90% 의 광범위한 분포", size=12, bold=True,
             color=TH.warn, font=TH.font_main)

    add_round_rect(s, Inches(6.5), Inches(2.55), Inches(6.3), Inches(2.7),
                   TH.card, border=TH.border)
    tf = add_text(s, Inches(6.7), Inches(2.7), Inches(6.0), Inches(0.4),
                  "■ " + C.S3_PROBLEM["reality"][0],
                  size=12, color=TH.text, font=TH.font_main, line_spacing=1.4)
    for line in C.S3_PROBLEM["reality"][1:]:
        add_paragraph(tf, "■ " + line, size=12, color=TH.text,
                      font=TH.font_main, space_after=8)

    # 하단: 핵심 메시지 박스
    add_round_rect(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.05),
                   TH.primary)
    add_text(s, Inches(0.85), Inches(5.7), Inches(11.5), Inches(0.85),
             C.S3_PROBLEM["tagline"], size=14, bold=True, color=TH.text_white,
             align=PP_ALIGN.CENTER, font=TH.font_main, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.4)


# ─────────────────────────────────────────
# S4. Background — Exqutor 두 해법
# ─────────────────────────────────────────
def s4_background(prs):
    s = blank_slide(prs)
    content_header(s, C.S4_BACKGROUND["title"], C.S4_BACKGROUND["subtitle"])
    slide_number(s, 4, TOTAL, TH.text_dim)

    for i, sol in enumerate(C.S4_BACKGROUND["solutions"]):
        x = Inches(0.6 + i * 6.15)
        y = Inches(1.55)
        # 카드 외곽
        add_round_rect(s, x, y, Inches(5.95), Inches(4.3),
                       TH.card, border=TH.border)
        # 상단 컬러 바
        add_rect(s, x, y, Inches(5.95), Inches(0.45), TH.primary)
        # 솔루션 이름
        add_text(s, x + Inches(0.25), y + Inches(0.05), Inches(5.5), Inches(0.4),
                 sol["name"], size=16, bold=True, color=TH.text_white,
                 font=TH.font_main)
        # 헤드라인
        add_text(s, x + Inches(0.25), y + Inches(0.65), Inches(5.5), Inches(0.4),
                 sol["headline"], size=13, bold=True, color=TH.primary,
                 font=TH.font_main)
        # 본문
        add_text(s, x + Inches(0.25), y + Inches(1.15), Inches(5.5), Inches(1.5),
                 sol["detail"], size=12, color=TH.text, font=TH.font_main,
                 line_spacing=1.5)
        # 메트릭 박스
        add_round_rect(s, x + Inches(0.25), y + Inches(2.85), Inches(5.45),
                       Inches(1.20), TH.bg_light, border=TH.accent, border_width=1.5)
        add_text(s, x + Inches(0.25), y + Inches(2.95), Inches(5.45), Inches(0.65),
                 sol["metric"], size=30, bold=True, color=TH.primary,
                 align=PP_ALIGN.CENTER, font=TH.font_main)
        add_text(s, x + Inches(0.25), y + Inches(3.55), Inches(5.45), Inches(0.4),
                 sol["metric_label"], size=11, color=TH.text_dim,
                 align=PP_ALIGN.CENTER, font=TH.font_main)

    # 하단: 종합 성과
    add_round_rect(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.65),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(6.18), Inches(11.5), Inches(0.5),
             "→  " + C.S4_BACKGROUND["achievements"], size=12, bold=True,
             color=TH.accent, anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S5. Limitation — 사각지대
# ─────────────────────────────────────────
def s5_limitation(prs):
    s = blank_slide(prs)
    content_header(s, C.S5_LIMITATION["title"], C.S5_LIMITATION["subtitle"])
    slide_number(s, 5, TOTAL, TH.text_dim)

    # 좌측 (6inch): findings 카드 3 개 (텍스트 좁힘)
    for i, f in enumerate(C.S5_LIMITATION["findings"]):
        y = Inches(1.55 + i * 1.40)
        add_round_rect(s, Inches(0.6), y, Inches(5.6), Inches(1.25),
                       TH.card, border=TH.border)
        add_rect(s, Inches(0.6), y, Inches(0.45), Inches(1.25), TH.primary)
        add_text(s, Inches(0.6), y + Inches(0.4), Inches(0.45), Inches(0.5),
                 f"#{i+1}", size=15, bold=True, color=TH.text_white,
                 align=PP_ALIGN.CENTER, font=TH.font_main)
        add_text(s, Inches(1.15), y + Inches(0.10), Inches(5.0), Inches(0.35),
                 f["tag"], size=12, bold=True, color=TH.primary, font=TH.font_main)
        add_text(s, Inches(1.15), y + Inches(0.42), Inches(5.0), Inches(0.32),
                 f["code"], size=8, color=TH.warn, font=TH.font_mono)
        add_text(s, Inches(1.15), y + Inches(0.72), Inches(5.0), Inches(0.55),
                 f["detail"], size=9, color=TH.text, font=TH.font_main,
                 line_spacing=1.3)

    # 우측 (6.4inch): vector.c snippet figure — 비율 2.54
    add_text(s, Inches(6.4), Inches(1.50), Inches(6.4), Inches(0.35),
             "vector.c — Hook Trigger 사각지대 (line 243)", size=11, bold=True,
             color=TH.primary, font=TH.font_main)
    add_figure(s, C.FIGURES["vector_c"], Inches(6.4), Inches(1.90),
               Inches(6.35), Inches(3.0),
               caption=C.FIGURE_CAPTIONS["vector_c"],
               caption_color=TH.text_dim, font=TH.font_main)

    # 하단 시사점
    add_round_rect(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.65),
                   TH.warn)
    add_text(s, Inches(0.85), Inches(6.18), Inches(11.5), Inches(0.5),
             "★  " + C.S5_LIMITATION["implication"], size=12, bold=True,
             color=TH.text_white, anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S6. Differentiator
# ─────────────────────────────────────────
def s6_differentiator(prs):
    s = blank_slide(prs)
    content_header(s, C.S6_DIFFERENTIATOR["title"], C.S6_DIFFERENTIATOR["subtitle"])
    slide_number(s, 6, TOTAL, TH.text_dim)

    # 두 Pivot 카드
    for i, lane in enumerate(C.S6_DIFFERENTIATOR["lanes"]):
        x = Inches(0.6 + i * 6.15)
        y = Inches(1.55)
        add_round_rect(s, x, y, Inches(5.95), Inches(2.4), TH.card, border=TH.border)
        # 좌측 컬러 라인
        add_rect(s, x, y, Inches(0.18), Inches(2.4), TH.accent)
        # 이름
        add_text(s, x + Inches(0.4), y + Inches(0.08), Inches(2.5), Inches(0.4),
                 lane["name"], size=11, bold=True, color=TH.accent, font=TH.font_main)
        # 제목
        add_text(s, x + Inches(0.4), y + Inches(0.42), Inches(5.4), Inches(0.5),
                 lane["title"], size=17, bold=True, color=TH.text, font=TH.font_main)
        # 설명
        add_text(s, x + Inches(0.4), y + Inches(0.92), Inches(5.4), Inches(0.7),
                 lane["detail"], size=11, color=TH.text, font=TH.font_main,
                 line_spacing=1.35)
        # 역할 라벨 + 메트릭
        add_text(s, x + Inches(0.4), y + Inches(1.6), Inches(5.4), Inches(0.3),
                 lane["role"], size=10, color=TH.text_dim, font=TH.font_main)
        add_text(s, x + Inches(0.4), y + Inches(1.85), Inches(2.5), Inches(0.5),
                 lane["metric"], size=22, bold=True, color=TH.primary,
                 font=TH.font_main)
        add_text(s, x + Inches(2.95), y + Inches(1.95), Inches(2.9), Inches(0.4),
                 lane["metric_label"], size=10, color=TH.text_dim,
                 font=TH.font_main)

    # 직교성 메시지
    add_round_rect(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.55),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(4.13), Inches(11.5), Inches(0.5),
             "✦  " + C.S6_DIFFERENTIATOR["tagline"], size=11, bold=True,
             color=TH.accent, anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)

    # 본 연구 기여 3 줄 (좌측) + Two-Level figure (우측) — 검은 띠와 figure 사이 spacing 0.3 inch 확보
    add_text(s, Inches(0.6), Inches(4.95), Inches(4.6), Inches(0.35),
             "본 연구의 핵심 기여", size=12, bold=True, color=TH.primary,
             font=TH.font_main)
    contrib_y = Inches(5.35)
    for i, c in enumerate(C.S6_DIFFERENTIATOR["contribution"]):
        y = contrib_y + Inches(i * 0.5)
        add_oval(s, Inches(0.7), y + Inches(0.09), Inches(0.13),
                 Inches(0.13), TH.accent)
        add_text(s, Inches(0.95), y, Inches(4.4), Inches(0.45),
                 c, size=9, color=TH.text, font=TH.font_main, line_spacing=1.3)

    # 우측: Two-Level decomposition figure — y 4.6→4.95 (검은 띠와 0.30 inch spacing), height 2.7→2.4
    add_figure(s, C.FIGURES["two_level"], Inches(5.4), Inches(4.95),
               Inches(7.3), Inches(2.4),
               caption=C.FIGURE_CAPTIONS["two_level"],
               caption_color=TH.text_dim, font=TH.font_main)


# ─────────────────────────────────────────
# S7. RQ1
# ─────────────────────────────────────────
def s7_rq1(prs):
    s = blank_slide(prs)
    content_header(s, C.S7_RQ1["title"], C.S7_RQ1["subtitle"])
    slide_number(s, 7, TOTAL, TH.text_dim)

    # 좌측: 실험 셋업 + Findings 컴팩트
    add_text(s, Inches(0.6), Inches(1.45), Inches(5.5), Inches(0.35),
             "실험 셋업", size=11, bold=True, color=TH.primary, font=TH.font_main)
    for i, (k, v) in enumerate(C.S7_RQ1["setup"]):
        y = Inches(1.78 + i * 0.42)
        add_text(s, Inches(0.6), y + Inches(0.05), Inches(1.4), Inches(0.35),
                 k, size=9, color=TH.text_dim, font=TH.font_main)
        add_text(s, Inches(2.1), y + Inches(0.03), Inches(4.0), Inches(0.35),
                 v, size=10, bold=True, color=TH.text, font=TH.font_main)

    # 좌측 하단: Findings
    add_text(s, Inches(0.6), Inches(3.55), Inches(5.5), Inches(0.35),
             "주요 발견 (Phase 1 + Phase 5)", size=11, bold=True,
             color=TH.primary, font=TH.font_main)
    add_round_rect(s, Inches(0.6), Inches(3.9), Inches(5.55), Inches(2.0),
                   TH.card, border=TH.border)
    tf = add_text(s, Inches(0.75), Inches(4.0), Inches(5.25), Inches(0.4),
                  "①  " + C.S7_RQ1["findings"][0],
                  size=10, color=TH.text, font=TH.font_main, line_spacing=1.45)
    for i, line in enumerate(C.S7_RQ1["findings"][1:], start=2):
        prefix = ["②", "③"][i - 2]
        add_paragraph(tf, f"{prefix}  {line}", size=10, color=TH.text,
                      font=TH.font_main, space_after=6)

    # 우측: figure (phase5 heatmap) — 비율 2.47, max 6.35×2.6 + caption
    add_text(s, Inches(6.4), Inches(1.45), Inches(6.4), Inches(0.35),
             "Phase 5 Local Skew × Selectivity Spearman ρ",
             size=11, bold=True, color=TH.primary, font=TH.font_main)
    add_figure(s, C.FIGURES["phase5_heatmap"], Inches(6.4), Inches(1.85),
               Inches(6.35), Inches(2.95),
               caption=C.FIGURE_CAPTIONS["phase5_heatmap"],
               caption_color=TH.text_dim, font=TH.font_main)

    # 해석
    add_round_rect(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.55),
                   TH.warn)
    add_text(s, Inches(0.85), Inches(6.22), Inches(11.5), Inches(0.4),
             "→  " + C.S7_RQ1["interpretation"],
             size=11, bold=True, color=TH.text_white,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main, line_spacing=1.3)


# ─────────────────────────────────────────
# S8. RQ2 BERNOULLI 표
# ─────────────────────────────────────────
def s8_rq2_bernoulli(prs):
    s = blank_slide(prs)
    content_header(s, C.S8_RQ2_BERNOULLI["title"], C.S8_RQ2_BERNOULLI["subtitle"])
    slide_number(s, 8, TOTAL, TH.text_dim)

    rows = C.S8_RQ2_BERNOULLI["table_rows"]
    headers = C.S8_RQ2_BERNOULLI["table_headers"]
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_x, table_y = Inches(0.6), Inches(1.55)
    # 표 (좌측 6.8 inch) + figure (우측 5.4 inch) 분할 layout
    table_w, table_h = Inches(6.8), Inches(0.36 * n_rows + 0.2)
    tbl_shape = s.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, table_h)
    tbl = tbl_shape.table

    col_widths = [0.85, 1.2, 1.2, 0.85, 1.35, 1.35]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    # 헤더
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = TH.primary
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = TH.font_main; run.font.size = Pt(10)
        run.font.bold = True; run.font.color.rgb = TH.text_white

    # 본문
    for r, row in enumerate(rows, start=1):
        is_sig = row[5] in ("★★", "★★★")
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH.card if r % 2 else TH.bg_light
            if c == 5 and is_sig:
                cell.fill.fore_color.rgb = TH.accent
            tf = cell.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            run.font.name = TH.font_main; run.font.size = Pt(10)
            run.font.bold = (c == 5 and is_sig) or (c == 0)
            run.font.color.rgb = TH.text_white if (c == 5 and is_sig) else TH.text

    # 우측: figure (phase4 scatter 2x2 grid) — 새 비율 1.18, max 5.4×4.6
    add_figure(s, C.FIGURES["phase4_scatter"], Inches(7.5), Inches(1.55),
               Inches(5.3), Inches(4.5),
               caption=C.FIGURE_CAPTIONS["phase4_scatter"],
               caption_color=TH.text_dim, font=TH.font_main)

    # 하단: takeaways (왼쪽 표 아래)
    add_text(s, Inches(0.6), Inches(4.50), Inches(6.9), Inches(0.32),
             "Takeaways", size=11, bold=True, color=TH.primary, font=TH.font_main)
    add_round_rect(s, Inches(0.6), Inches(4.85), Inches(6.85), Inches(1.85),
                   TH.card, border=TH.border)
    tf = add_text(s, Inches(0.78), Inches(4.95), Inches(6.65), Inches(0.4),
                  "①  " + C.S8_RQ2_BERNOULLI["takeaways"][0],
                  size=10, color=TH.text, font=TH.font_main, line_spacing=1.4)
    for i, t in enumerate(C.S8_RQ2_BERNOULLI["takeaways"][1:], start=2):
        prefix = "②" if i == 2 else "③"
        add_paragraph(tf, f"{prefix}  {t}", size=10, color=TH.text,
                      font=TH.font_main, space_after=6)


# ─────────────────────────────────────────
# S9. RQ2 KM20
# ─────────────────────────────────────────
def s9_rq2_km20(prs):
    s = blank_slide(prs)
    content_header(s, C.S9_RQ2_KM20["title"], C.S9_RQ2_KM20["subtitle"])
    slide_number(s, 9, TOTAL, TH.text_dim)

    rows = C.S9_RQ2_KM20["table_rows"]
    headers = C.S9_RQ2_KM20["table_headers"]
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.55),
                                   Inches(7.5), Inches(0.4 * n_rows + 0.3))
    tbl = tbl_shape.table
    widths = [1.5, 2.0, 2.5, 1.5]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = TH.primary
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = TH.font_main; run.font.size = Pt(11)
        run.font.bold = True; run.font.color.rgb = TH.text_white

    for r, row in enumerate(rows, start=1):
        is_5seed = "5/5" in row[3]
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH.card if r % 2 else TH.bg_light
            tf = cell.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            run.font.name = TH.font_main; run.font.size = Pt(11)
            run.font.bold = (c <= 1 and is_5seed)
            run.font.color.rgb = TH.success if is_5seed and c == 1 else TH.text

    # 우측: figure (phase6 box) — 비율 1.65, max 4.5×3.0
    add_text(s, Inches(8.3), Inches(1.55), Inches(4.5), Inches(0.35),
             "BERNOULLI vs STRATIFIED (KM20)",
             size=11, bold=True, color=TH.primary, font=TH.font_main)
    add_figure(s, C.FIGURES["phase6_box"], Inches(8.3), Inches(1.95),
               Inches(4.5), Inches(2.85),
               caption=C.FIGURE_CAPTIONS["phase6_box"],
               caption_color=TH.text_dim, font=TH.font_main)

    # takeaways
    add_round_rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.2),
                   TH.card, border=TH.border)
    tf = add_text(s, Inches(0.85), Inches(5.1), Inches(11.5), Inches(0.45),
                  "①  " + C.S9_RQ2_KM20["takeaways"][0],
                  size=10, color=TH.text, font=TH.font_main, line_spacing=1.45)
    for t in C.S9_RQ2_KM20["takeaways"][1:]:
        prefix = "②" if "Bonferroni" in t else "③"
        add_paragraph(tf, f"{prefix}  {t}", size=10, color=TH.text,
                      font=TH.font_main, space_after=4)

    # 하단: vector.c 구현 강조
    add_round_rect(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.45),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(6.33), Inches(11.5), Inches(0.4),
             "vector.c  +228 줄 native  ·  GUC 3 개  ·  stratum_id 사전 부여  "
             "·  Horvitz-Thompson 가중 추정",
             size=10, bold=True, color=TH.accent,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S10. External Validity
# ─────────────────────────────────────────
def s10_external(prs):
    s = blank_slide(prs)
    content_header(s, C.S10_EXTERNAL_VALIDITY["title"],
                   C.S10_EXTERNAL_VALIDITY["subtitle"])
    slide_number(s, 10, TOTAL, TH.text_dim)

    # 좌측: figure (cross_dataset bar) ★ 핵심 시각화 — 비율 1.26
    add_figure(s, C.FIGURES["cross_dataset"], Inches(0.6), Inches(1.55),
               Inches(5.5), Inches(4.4),
               caption=C.FIGURE_CAPTIONS["cross_dataset"],
               caption_color=TH.text_dim, font=TH.font_main)

    # 우측: 3 메트릭 + findings (영역 확대 6.4 inch)
    for i, m in enumerate(C.S10_EXTERNAL_VALIDITY["metrics"]):
        x = Inches(6.4)
        y = Inches(1.55 + i * 1.0)
        bar_color = TH.warn if m["color"] == "orange" else TH.primary
        add_round_rect(s, x, y, Inches(6.35), Inches(0.9),
                       TH.card, border=TH.border)
        add_rect(s, x, y, Inches(0.18), Inches(0.9), bar_color)
        # 라벨
        add_text(s, x + Inches(0.3), y + Inches(0.05), Inches(3.2), Inches(0.32),
                 m["label"], size=11, bold=True, color=TH.text, font=TH.font_main)
        # CI / CV
        add_text(s, x + Inches(0.3), y + Inches(0.42), Inches(3.5), Inches(0.45),
                 f"{m['ci']}  ·  {m['cv']}", size=9, color=TH.text_dim,
                 font=TH.font_main, line_spacing=1.3)
        # 큰 숫자
        add_text(s, x + Inches(4.0), y + Inches(0.10), Inches(2.2), Inches(0.7),
                 m["value"], size=24, bold=True, color=bar_color,
                 align=PP_ALIGN.RIGHT, font=TH.font_main)

    # Findings (하단) — 3 줄 분리로 가독성 보강
    add_text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.28),
             "주요 발견 + 시사점", size=10, bold=True, color=TH.primary,
             font=TH.font_main)
    findings_tf = add_text(s, Inches(0.6), Inches(6.30), Inches(12.1),
                           Inches(0.32),
                           "①  " + C.S10_EXTERNAL_VALIDITY["key_findings"][0],
                           size=9, color=TH.text, font=TH.font_main,
                           line_spacing=1.25)
    add_paragraph(findings_tf, "②  " + C.S10_EXTERNAL_VALIDITY["key_findings"][1],
                  size=9, color=TH.text, font=TH.font_main, space_after=0)
    add_paragraph(findings_tf, "③  " + C.S10_EXTERNAL_VALIDITY["key_findings"][2],
                  size=9, color=TH.text, font=TH.font_main, space_after=0)


# ─────────────────────────────────────────
# S11. 결과 해석 — Two-Level + Cluster Skew (신규)
# ─────────────────────────────────────────
def s11_interpretation(prs):
    s = blank_slide(prs)
    content_header(s, C.S10B_INTERPRETATION["title"],
                   C.S10B_INTERPRETATION["subtitle"])
    slide_number(s, 11, TOTAL, TH.text_dim)

    # 좌측 (6.0 inch): Two-Level decomposition table
    add_text(s, Inches(0.6), Inches(1.5), Inches(6.2), Inches(0.35),
             "▣ " + C.S10B_INTERPRETATION["two_level"]["headline"],
             size=11, bold=True, color=TH.primary, font=TH.font_main)

    rows = C.S10B_INTERPRETATION["two_level"]["rows"]
    headers = C.S10B_INTERPRETATION["two_level"]["headers"]
    n_rows = len(rows) + 1
    tbl_shape = s.shapes.add_table(n_rows, len(headers),
                                   Inches(0.6), Inches(1.9),
                                   Inches(6.2), Inches(0.35 * n_rows + 0.1))
    tbl = tbl_shape.table
    widths = [1.05, 1.0, 1.0, 1.4, 1.75]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = TH.primary
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = TH.font_main; run.font.size = Pt(9)
        run.font.bold = True; run.font.color.rgb = TH.text_white
    for r, row in enumerate(rows, start=1):
        is_star = "★" in row[4]
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH.warn if is_star else (
                TH.card if r % 2 else TH.bg_light)
            tf = cell.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            run.font.name = TH.font_main; run.font.size = Pt(9)
            run.font.bold = is_star or (c == 0)
            run.font.color.rgb = TH.text_white if is_star else TH.text

    # 우측 (6.4 inch): cluster_skew figure
    add_text(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.35),
             "▣ " + C.S10B_INTERPRETATION["skew"]["headline"],
             size=11, bold=True, color=TH.primary, font=TH.font_main)
    add_figure(s, C.FIGURES["cluster_skew"], Inches(7.0), Inches(1.9),
               Inches(5.85), Inches(3.0),
               caption=C.FIGURE_CAPTIONS["cluster_skew"],
               caption_color=TH.text_dim, font=TH.font_main)

    # 우측 하단: skew items
    items_y = Inches(5.1)
    for i, item in enumerate(C.S10B_INTERPRETATION["skew"]["items"]):
        y = items_y + Inches(i * 0.32)
        add_oval(s, Inches(7.0), y + Inches(0.07), Inches(0.12),
                 Inches(0.12), TH.accent)
        add_text(s, Inches(7.25), y, Inches(5.5), Inches(0.32),
                 item, size=10, color=TH.text, font=TH.font_main, line_spacing=1.3)

    # 하단 tagline
    add_round_rect(s, Inches(0.6), Inches(6.20), Inches(12.1), Inches(0.55),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(6.27), Inches(11.5), Inches(0.45),
             "★  " + C.S10B_INTERPRETATION["tagline"],
             size=11, bold=True, color=TH.accent,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S12. RQ3 design (formerly S11)
# ─────────────────────────────────────────
def s11_rq3(prs):
    s = blank_slide(prs)
    content_header(s, C.S11_RQ3_DESIGN["title"], C.S11_RQ3_DESIGN["subtitle"])
    slide_number(s, 12, TOTAL, TH.text_dim)

    # Recovery Rate 식
    add_round_rect(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.65),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(1.53), Inches(11.5), Inches(0.5),
             "▣ " + C.S11_RQ3_DESIGN["framework"],
             size=14, bold=True, color=TH.accent,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)

    # 3 패러다임 카드
    for i, p in enumerate(C.S11_RQ3_DESIGN["paradigms"]):
        x = Inches(0.6 + i * 4.1)
        y = Inches(2.3)
        add_round_rect(s, x, y, Inches(3.95), Inches(3.0),
                       TH.card, border=TH.border)
        # 번호 + 패러다임명
        add_text(s, x + Inches(0.25), y + Inches(0.1), Inches(0.5), Inches(0.4),
                 f"P{i+1}", size=14, bold=True, color=TH.accent,
                 font=TH.font_main)
        add_text(s, x + Inches(0.85), y + Inches(0.12), Inches(3.0), Inches(0.4),
                 p["name"], size=14, bold=True, color=TH.text, font=TH.font_main)
        add_rect(s, x + Inches(0.25), y + Inches(0.65), Inches(3.5), Pt(1),
                 TH.border)

        # 방법 리스트
        for j, m in enumerate(p["methods"]):
            mid = y + Inches(0.85 + j * 0.42)
            add_oval(s, x + Inches(0.3), mid + Inches(0.07), Inches(0.12),
                     Inches(0.12), TH.accent)
            add_text(s, x + Inches(0.55), mid, Inches(3.4), Inches(0.4),
                     m, size=11, color=TH.text, font=TH.font_main)
        # rationale
        add_text(s, x + Inches(0.25), y + Inches(2.35), Inches(3.5), Inches(0.55),
                 p["rationale"], size=9, color=TH.text_dim,
                 font=TH.font_main, line_spacing=1.4)

    # 평가 / 기대치
    add_text(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.35),
             "평가 조건", size=11, bold=True, color=TH.primary, font=TH.font_main)
    add_text(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.35),
             "→  " + C.S11_RQ3_DESIGN["evaluation"],
             size=11, color=TH.text, font=TH.font_main)

    add_round_rect(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.45),
                   TH.success)
    add_text(s, Inches(0.85), Inches(6.33), Inches(11.5), Inches(0.4),
             "★  " + C.S11_RQ3_DESIGN["expected"],
             size=11, bold=True, color=TH.text_white,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S12. Progress
# ─────────────────────────────────────────
def s12_progress(prs):
    s = blank_slide(prs)
    content_header(s, C.S12_PROGRESS["title"])
    slide_number(s, 13, TOTAL, TH.text_dim)

    weeks = C.S12_PROGRESS["weeks"]
    for i, w in enumerate(weeks):
        x = Inches(0.6 + i * 3.05)
        y = Inches(1.55)
        is_current = w["status"] == "▶"
        bar_color = TH.warn if is_current else TH.success

        add_round_rect(s, x, y, Inches(2.95), Inches(4.4),
                       TH.card, border=TH.border)
        add_rect(s, x, y, Inches(2.95), Inches(0.55), bar_color)
        add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(1.5), Inches(0.4),
                 w["id"] + "  " + w["status"], size=14, bold=True,
                 color=TH.text_white, font=TH.font_main)
        add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(2.6), Inches(0.35),
                 w["period"], size=10, color=TH.text_dim, font=TH.font_main)
        add_text(s, x + Inches(0.2), y + Inches(1.05), Inches(2.6), Inches(0.4),
                 w["title"], size=13, bold=True, color=TH.primary,
                 font=TH.font_main)
        # 항목들
        for j, item in enumerate(w["items"]):
            top = y + Inches(1.55 + j * 0.55)
            add_text(s, x + Inches(0.2), top, Inches(2.7), Inches(0.55),
                     "·  " + item, size=9, color=TH.text, font=TH.font_main,
                     line_spacing=1.4)

    # summary
    add_round_rect(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.55),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(6.25), Inches(11.5), Inches(0.5),
             "✓  " + C.S12_PROGRESS["summary"],
             size=12, bold=True, color=TH.accent,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S13. Plan
# ─────────────────────────────────────────
def s13_plan(prs):
    s = blank_slide(prs)
    content_header(s, C.S13_PLAN["title"])
    slide_number(s, 14, TOTAL, TH.text_dim)

    # Gantt 행
    for i, (wid, period, title, detail) in enumerate(C.S13_PLAN["weeks"]):
        y = Inches(1.55 + i * 0.65)
        add_rect(s, Inches(0.6), y + Inches(0.55), Inches(12.1), Pt(1),
                 TH.border)
        # week id
        add_round_rect(s, Inches(0.6), y, Inches(0.85), Inches(0.5),
                       TH.primary)
        add_text(s, Inches(0.6), y + Inches(0.07), Inches(0.85), Inches(0.4),
                 wid, size=12, bold=True, color=TH.text_white,
                 align=PP_ALIGN.CENTER, font=TH.font_main)
        # period
        add_text(s, Inches(1.65), y + Inches(0.05), Inches(1.7), Inches(0.4),
                 period, size=10, color=TH.text_dim, font=TH.font_main)
        # title
        add_text(s, Inches(3.45), y + Inches(0.05), Inches(3.5), Inches(0.4),
                 title, size=13, bold=True, color=TH.text, font=TH.font_main)
        # detail
        add_text(s, Inches(7.05), y + Inches(0.08), Inches(5.7), Inches(0.4),
                 detail, size=10, color=TH.text_dim, font=TH.font_main)

    # 산출물
    add_text(s, Inches(0.6), Inches(6.3), Inches(2.5), Inches(0.4),
             "산출물", size=11, bold=True, color=TH.primary, font=TH.font_main)
    add_text(s, Inches(2.0), Inches(6.3), Inches(10.5), Inches(0.4),
             "  ·  ".join(C.S13_PLAN["deliverables"]),
             size=11, color=TH.text, font=TH.font_main)


# ─────────────────────────────────────────
# S14. Roles
# ─────────────────────────────────────────
def s14_roles(prs):
    """4-step 카드 형태 (Portfolio About Me 변형) — 각자 핵심 산출물 강조."""
    s = blank_slide(prs)
    content_header(s, C.S14_ROLES["title"])
    slide_number(s, 15, TOTAL, TH.text_dim)

    members = C.S14_ROLES["members"]
    for i, m in enumerate(members):
        x = Inches(0.6 + i * 3.05)
        y = Inches(1.6)
        # 카드
        add_round_rect(s, x, y, Inches(2.95), Inches(4.4),
                       TH.card, border=TH.border)
        # 상단 컬러 바
        add_rect(s, x, y, Inches(2.95), Inches(0.5), TH.primary)
        # 번호 배지 (좌상단 흰색 원)
        add_oval(s, x + Inches(0.15), y + Inches(0.13), Inches(0.27),
                 Inches(0.27), TH.text_white)
        add_text(s, x + Inches(0.15), y + Inches(0.13), Inches(0.27),
                 Inches(0.27), str(i+1), size=11, bold=True, color=TH.primary,
                 align=PP_ALIGN.CENTER, font=TH.font_main)
        # 역할 (상단 바 안)
        add_text(s, x + Inches(0.55), y + Inches(0.10), Inches(2.3),
                 Inches(0.32), m["role"], size=12, bold=True,
                 color=TH.text_white, font=TH.font_main)
        # 이름 (큰 글씨)
        add_text(s, x + Inches(0.2), y + Inches(0.65), Inches(2.6),
                 Inches(0.55), m["name"], size=22, bold=True,
                 color=TH.primary, font=TH.font_main)
        # 분리선
        add_rect(s, x + Inches(0.2), y + Inches(1.30), Inches(1.2), Pt(2),
                 TH.accent)
        # 담당
        add_text(s, x + Inches(0.2), y + Inches(1.45), Inches(2.6),
                 Inches(1.55), m["duty"], size=10, color=TH.text,
                 font=TH.font_main, line_spacing=1.4)
        # Highlight 라벨
        add_text(s, x + Inches(0.2), y + Inches(3.15), Inches(2.6),
                 Inches(0.3), "HIGHLIGHT", size=8, bold=True,
                 color=TH.accent, font=TH.font_main)
        # Highlight 박스
        add_round_rect(s, x + Inches(0.2), y + Inches(3.45), Inches(2.55),
                       Inches(0.85), TH.bg_light, border=TH.accent,
                       border_width=0.75)
        add_text(s, x + Inches(0.3), y + Inches(3.50), Inches(2.4),
                 Inches(0.75), m["highlight"], size=10, bold=True,
                 color=TH.text, font=TH.font_main, line_spacing=1.35)

    # 공유 원칙
    add_round_rect(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.55),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(6.20), Inches(11.5), Inches(0.45),
             "▣  " + C.S14_ROLES["shared"],
             size=10, bold=True, color=TH.accent,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S16 (신규) — Conclusion (4 핵심 발견)
# ─────────────────────────────────────────
def s15_conclusion(prs):
    s = blank_slide(prs)
    content_header(s, C.S15_CONCLUSION["title"], C.S15_CONCLUSION["subtitle"])
    slide_number(s, 16, TOTAL, TH.text_dim)

    findings = C.S15_CONCLUSION["findings"]
    for i, f in enumerate(findings):
        x = Inches(0.6 + i * 3.05)
        y = Inches(1.55)
        add_round_rect(s, x, y, Inches(2.95), Inches(4.5),
                       TH.card, border=TH.border)
        # 상단 컬러 바
        add_rect(s, x, y, Inches(2.95), Inches(0.55), TH.primary)
        add_text(s, x + Inches(0.2), y + Inches(0.10), Inches(2.6),
                 Inches(0.4), f["tag"], size=14, bold=True,
                 color=TH.text_white, font=TH.font_main)
        # 헤드라인
        add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(2.6),
                 Inches(0.85), f["headline"], size=12, bold=True,
                 color=TH.text, font=TH.font_main, line_spacing=1.3)
        # 큰 메트릭 (강조 박스)
        add_round_rect(s, x + Inches(0.2), y + Inches(1.65),
                       Inches(2.55), Inches(1.05), TH.bg_light,
                       border=TH.accent, border_width=1.5)
        add_text(s, x + Inches(0.2), y + Inches(1.70), Inches(2.55),
                 Inches(0.55), f["metric"], size=22, bold=True,
                 color=TH.primary, align=PP_ALIGN.CENTER, font=TH.font_main)
        add_text(s, x + Inches(0.2), y + Inches(2.25), Inches(2.55),
                 Inches(0.45), f["metric_label"], size=9, color=TH.text_dim,
                 align=PP_ALIGN.CENTER, font=TH.font_main, line_spacing=1.3)
        # 디테일
        add_text(s, x + Inches(0.2), y + Inches(2.85), Inches(2.6),
                 Inches(1.55), f["detail"], size=9, color=TH.text,
                 font=TH.font_main, line_spacing=1.4)

    # 하단 tagline + References footer
    add_round_rect(s, Inches(0.6), Inches(6.20), Inches(12.1), Inches(0.55),
                   TH.bg_dark)
    add_text(s, Inches(0.85), Inches(6.25), Inches(11.5), Inches(0.45),
             "★  " + C.S15_CONCLUSION["tagline"],
             size=11, bold=True, color=TH.accent,
             anchor=MSO_ANCHOR.MIDDLE, font=TH.font_main)


# ─────────────────────────────────────────
# S17. Thanks
# ─────────────────────────────────────────
def s15_thanks(prs):
    s = blank_slide(prs)
    set_bg(s, TH.bg_dark)

    add_rect(s, Inches(5.5), Inches(2.6), Inches(2.4), Pt(3), TH.accent)
    add_text(s, Inches(0), Inches(2.85), SLIDE_W, Inches(1.2),
             C.S16_THANKS["title"], size=58, bold=True, color=TH.text_white,
             align=PP_ALIGN.CENTER, font=TH.font_main)
    add_text(s, Inches(0), Inches(4.1), SLIDE_W, Inches(0.6),
             C.S16_THANKS["subtitle"], size=18, color=TH.accent,
             align=PP_ALIGN.CENTER, font=TH.font_main)
    add_text(s, Inches(0), Inches(5.5), SLIDE_W, Inches(0.4),
             C.S16_THANKS["team"], size=15, bold=True, color=TH.text_white,
             align=PP_ALIGN.CENTER, font=TH.font_main)
    add_text(s, Inches(0), Inches(5.95), SLIDE_W, Inches(0.4),
             " · ".join(C.S16_THANKS["members"]),
             size=11, color=TH.text_dim,
             align=PP_ALIGN.CENTER, font=TH.font_main)


# ─────────────────────────────────────────
# Main
# ─────────────────────────────────────────
def build():
    prs = new_presentation()
    s1_cover(prs)
    s2_toc(prs)
    s3_problem(prs)
    s4_background(prs)
    s5_limitation(prs)
    s6_differentiator(prs)
    s7_rq1(prs)
    s8_rq2_bernoulli(prs)
    s9_rq2_km20(prs)
    s10_external(prs)
    s11_interpretation(prs)
    s11_rq3(prs)
    s12_progress(prs)
    s13_plan(prs)
    s14_roles(prs)
    s15_conclusion(prs)
    s15_thanks(prs)
    return prs


if __name__ == "__main__":
    import sys
    from pathlib import Path
    out = Path(sys.argv[1] if len(sys.argv) > 1 else
               "/Users/hyunbin/Capstone/submission/_drafts/속도는벡터_중간발표_navy.pptx")
    prs = build()
    prs.save(str(out))
    print(f"saved: {out}")
