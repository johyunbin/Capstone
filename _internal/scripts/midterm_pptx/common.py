"""
공통 유틸 — pptx 객체 생성, 한국어 폰트 강제, 형상 헬퍼.
모든 디자인 빌드가 이 모듈을 import.
"""
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as etree  # type: ignore

# 16:9 와이드
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def force_kr_font(run, font_name: str):
    """한국어 (CJK) 폰트를 latin/east-asia 모두 강제."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs", "a:latin"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font_name)


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color: Optional[RGBColor] = None,
             align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo",
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """텍스트 박스 추가 — 단일 paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    force_kr_font(run, font)
    return tf


def add_paragraph(tf, text, *,
                  size=12, bold=False, color: Optional[RGBColor] = None,
                  align=PP_ALIGN.LEFT, font="Apple SD Gothic Neo",
                  space_before=0, space_after=4, level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    force_kr_font(run, font)
    return p


def add_rect(slide, x, y, w, h, color: Optional[RGBColor] = None,
             border: Optional[RGBColor] = None, border_width=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if color is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = color
    else:
        s.fill.background()
    if border is not None:
        s.line.color.rgb = border
        s.line.width = Pt(border_width)
    else:
        s.line.fill.background()
    return s


def add_round_rect(slide, x, y, w, h, color: Optional[RGBColor] = None,
                   border: Optional[RGBColor] = None, border_width=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    s.adjustments[0] = 0.08  # 모서리 둥글기
    if color is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = color
    else:
        s.fill.background()
    if border is not None:
        s.line.color.rgb = border
        s.line.width = Pt(border_width)
    else:
        s.line.fill.background()
    return s


def add_oval(slide, x, y, w, h, color: Optional[RGBColor] = None,
             border: Optional[RGBColor] = None, border_width=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.shadow.inherit = False
    if color is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = color
    else:
        s.fill.background()
    if border is not None:
        s.line.color.rgb = border
        s.line.width = Pt(border_width)
    else:
        s.line.fill.background()
    return s


def add_line(slide, x1, y1, x2, y2, color: RGBColor, width=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def slide_number(slide, n, total, color: RGBColor, font="Apple SD Gothic Neo"):
    add_text(slide, Inches(12.3), Inches(7.05), Inches(1.0), Inches(0.35),
             f"{n} / {total}", size=9, color=color, align=PP_ALIGN.RIGHT, font=font)


def page_footer(slide, text, color: RGBColor, font="Apple SD Gothic Neo"):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(10.5), Inches(0.35),
             text, size=9, color=color, font=font)


def add_figure(slide, path, x, y, max_w, max_h, caption=None,
               caption_color=None, font="Apple SD Gothic Neo",
               align="center"):
    """원본 aspect ratio 를 유지하며 (x,y,max_w,max_h) bound 안에 fit.
    align: 'center' | 'top' | 'left'
    """
    from pathlib import Path
    from PIL import Image

    p = Path(path)
    if not p.exists():
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, max_w, max_h)
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        s.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tf = s.text_frame
        tf.text = f"[Figure missing: {p.name}]"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        return s

    img = Image.open(str(p))
    img_ratio = img.width / img.height
    box_ratio = max_w / max_h

    # 캡션 자리 확보 (있으면 이미지 영역 30pt 줄임)
    cap_h = Inches(0.30) if caption else Inches(0)
    avail_h = max_h - cap_h

    if img_ratio > (max_w / avail_h):
        # 이미지가 더 wide → 너비 max_w 로 맞춤
        actual_w = max_w
        actual_h = max_w / img_ratio
    else:
        # 이미지가 더 tall → 높이 avail_h 로 맞춤
        actual_h = avail_h
        actual_w = avail_h * img_ratio

    # 박스 안 정렬
    if align == "center":
        ox = (max_w - actual_w) / 2
        oy = (avail_h - actual_h) / 2
    elif align == "top":
        ox = (max_w - actual_w) / 2
        oy = 0
    else:  # left
        ox = 0
        oy = (avail_h - actual_h) / 2

    pic = slide.shapes.add_picture(str(p), x + ox, y + oy, actual_w, actual_h)

    if caption:
        cap_y = y + avail_h + Inches(0.05)
        cc = caption_color if caption_color is not None else RGBColor(0x5B, 0x71, 0x8D)
        add_text(slide, x, cap_y, max_w, Inches(0.25),
                 caption, size=9, color=cc, align=PP_ALIGN.CENTER, font=font)
    return pic
